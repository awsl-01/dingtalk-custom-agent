"""
调试模板结构
"""
import os
import sys
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')
sys.stderr.reconfigure(encoding='utf-8')


def analyze_template():
    """分析模板结构"""
    template_path = "templates/清新欧美风.pptx"

    if not os.path.exists(template_path):
        print(f"模板文件不存在: {template_path}")
        return

    prs = Presentation(template_path)
    slides = list(prs.slides)

    print("=== 模板结构分析 ===\n")

    # 分析第2页（目录页）
    if len(slides) >= 2:
        toc_slide = slides[1]
        print(f"第2页（目录页）形状分析：")
        print(f"形状数量: {len(toc_slide.shapes)}\n")

        # 按位置分组
        shapes_by_section = {}
        for i, shape in enumerate(toc_slide.shapes):
            if shape.has_text_frame:
                # 获取位置
                top_inches = shape.top / 914400
                left_inches = shape.left / 914400

                # 获取文本
                text = ""
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        text = para.text.strip()
                        break

                if text:
                    print(f"形状{i+1}:")
                    print(f"  位置: ({left_inches:.3f}, {top_inches:.3f})")
                    print(f"  尺寸: ({shape.width / 914400:.3f} x {shape.height / 914400:.3f})")
                    print(f"  文本: {text[:50]}")
                    print()


if __name__ == "__main__":
    analyze_template()
