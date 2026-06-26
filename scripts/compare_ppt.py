"""
对比模板文件和生成文件的差异
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt

sys.stdout.reconfigure(encoding='utf-8')
sys.stderr.reconfigure(encoding='utf-8')


def analyze_ppt(file_path, name):
    """分析PPT文件"""
    print(f"\n=== {name} ===")
    print(f"文件: {file_path}")

    if not os.path.exists(file_path):
        print(f"文件不存在!")
        return

    prs = Presentation(file_path)

    print(f"幻灯片数量: {len(prs.slides)}")
    print(f"幻灯片宽度: {prs.slide_width / 914400:.3f} 英寸")
    print(f"幻灯片高度: {prs.slide_height / 914400:.3f} 英寸")

    for i, slide in enumerate(prs.slides):
        print(f"\n--- 第{i+1}页 ---")

        # 背景
        try:
            bg = slide.background
            fill = bg.fill
            if fill.type is not None:
                if fill.type == 1:  # SOLID_FILL
                    print(f"  背景颜色: {fill.fore_color.rgb}")
        except Exception as e:
            print(f"  背景: 无法读取")

        # 形状
        print(f"  形状数量: {len(slide.shapes)}")
        for j, shape in enumerate(slide.shapes):
            print(f"    形状{j+1}: {shape.shape_type}")
            print(f"      位置: ({shape.left / 914400:.3f}, {shape.top / 914400:.3f})")
            print(f"      尺寸: ({shape.width / 914400:.3f} x {shape.height / 914400:.3f})")

            if shape.has_text_frame:
                for k, para in enumerate(shape.text_frame.paragraphs):
                    if para.text.strip():
                        text = para.text[:50]
                        if len(para.text) > 50:
                            text += "..."
                        print(f"      文本{k+1}: {text}")

                        if para.runs:
                            run = para.runs[0]
                            if run.font.size:
                                print(f"        字号: {run.font.size.pt}pt")


def main():
    """主函数"""
    # 分析模板文件
    template_path = "templates/default_template.pptx"
    analyze_ppt(template_path, "模板文件")

    # 分析生成的文件（使用最新的）
    generated_path = "test_output/从百草园到三味书屋课文讲解.pptx"
    if os.path.exists(generated_path):
        analyze_ppt(generated_path, "生成的文件")
    else:
        # 如果最新文件不存在，尝试其他文件
        import glob
        pptx_files = glob.glob("test_output/*.pptx")
        if pptx_files:
            generated_path = max(pptx_files, key=os.path.getctime)
            analyze_ppt(generated_path, "生成的文件")
        else:
            print(f"\n生成的文件不存在")

    # 对比
    print("\n=== 对比分析 ===")
    print("请检查以下差异：")
    print("1. 背景颜色是否一致")
    print("2. 形状位置和尺寸是否一致")
    print("3. 字体和字号是否一致")
    print("4. 装饰元素是否保留")


if __name__ == "__main__":
    main()
