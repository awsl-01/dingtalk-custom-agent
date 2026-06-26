"""
调试文本修改问题
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt

sys.stdout.reconfigure(encoding='utf-8')
sys.stderr.reconfigure(encoding='utf-8')


def test_modify():
    """测试文本修改"""
    template_path = "templates/default_template.pptx"

    if not os.path.exists(template_path):
        print(f"模板文件不存在: {template_path}")
        return

    # 加载模板
    prs = Presentation(template_path)

    # 获取目录页（第2页）
    slides = list(prs.slides)
    if len(slides) < 2:
        print("模板幻灯片数量不足")
        return

    toc_slide = slides[1]

    print("=== 目录页形状分析 ===")
    for i, shape in enumerate(toc_slide.shapes):
        print(f"\n形状{i+1}:")
        print(f"  类型: {shape.shape_type}")
        print(f"  位置: ({shape.left / 914400:.3f}, {shape.top / 914400:.3f})")
        print(f"  尺寸: ({shape.width / 914400:.3f} x {shape.height / 914400:.3f})")

        if shape.has_text_frame:
            print(f"  有文本框: True")
            for j, para in enumerate(shape.text_frame.paragraphs):
                if para.text.strip():
                    print(f"    段落{j+1}: '{para.text}'")
                    if para.runs:
                        print(f"      Runs数量: {len(para.runs)}")
                        for k, run in enumerate(para.runs):
                            print(f"      Run{k+1}: '{run.text}'")
                            print(f"        字体: {run.font.name}")
                            print(f"        字号: {run.font.size}")
        else:
            print(f"  有文本框: False")

    # 测试修改
    print("\n=== 测试修改 ===")
import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    from agent.template_based_ppt import TemplateBasedPPTGenerator

    generator = TemplateBasedPPTGenerator(template_path)

    # 复制目录页
    new_prs = Presentation()
    new_prs.slide_width = prs.slide_width
    new_prs.slide_height = prs.slide_height

    new_slide = generator._copy_slide(toc_slide, new_prs)

    print("\n复制后的形状:")
    for i, shape in enumerate(new_slide.shapes):
        print(f"\n形状{i+1}:")
        print(f"  类型: {shape.shape_type}")
        print(f"  位置: ({shape.left / 914400:.3f}, {shape.top / 914400:.3f})")

        if shape.has_text_frame:
            print(f"  有文本框: True")
            for j, para in enumerate(shape.text_frame.paragraphs):
                if para.text.strip():
                    print(f"    段落{j+1}: '{para.text}'")

    # 修改内容
    print("\n修改内容...")
    generator._add_content_to_slide(new_slide, "测试标题", ["要点1", "要点2", "要点3"])

    print("\n修改后的形状:")
    for i, shape in enumerate(new_slide.shapes):
        print(f"\n形状{i+1}:")
        if shape.has_text_frame:
            for j, para in enumerate(shape.text_frame.paragraphs):
                if para.text.strip():
                    print(f"  段落{j+1}: '{para.text}'")


if __name__ == "__main__":
    test_modify()
