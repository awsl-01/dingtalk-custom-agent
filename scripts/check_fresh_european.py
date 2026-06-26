"""
检查清新欧美风模板生成的PPT
"""
import os
import sys
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')
sys.stderr.reconfigure(encoding='utf-8')


def analyze_ppt(file_path):
    """详细分析PPT文件"""
    print(f"=== 分析文件: {file_path} ===")

    if not os.path.exists(file_path):
        print(f"文件不存在!")
        return

    prs = Presentation(file_path)

    print(f"\n基本信息:")
    print(f"  幻灯片数量: {len(prs.slides)}")
    print(f"  幻灯片宽度: {prs.slide_width / 914400:.3f} 英寸")
    print(f"  幻灯片高度: {prs.slide_height / 914400:.3f} 英寸")

    for i, slide in enumerate(prs.slides):
        print(f"\n{'='*50}")
        print(f"第{i+1}页:")
        print(f"{'='*50}")

        # 背景
        try:
            bg = slide.background
            fill = bg.fill
            if fill.type is not None:
                if fill.type == 1:  # SOLID_FILL
                    print(f"背景颜色: {fill.fore_color.rgb}")
        except:
            pass

        # 形状
        print(f"形状数量: {len(slide.shapes)}")
        for j, shape in enumerate(slide.shapes):
            print(f"\n  形状{j+1}:")
            print(f"    类型: {shape.shape_type}")
            print(f"    位置: ({shape.left / 914400:.3f}, {shape.top / 914400:.3f})")
            print(f"    尺寸: ({shape.width / 914400:.3f} x {shape.height / 914400:.3f})")

            if shape.has_text_frame:
                print(f"    文本框:")
                for k, para in enumerate(shape.text_frame.paragraphs):
                    if para.text.strip():
                        text = para.text
                        if len(text) > 60:
                            text = text[:60] + "..."
                        print(f"      段落{k+1}: {text}")

                        if para.runs:
                            run = para.runs[0]
                            if run.font.size:
                                print(f"        字号: {run.font.size.pt}pt")


def main():
    """主函数"""
    # 分析模板
    template_path = "templates/清新欧美风.pptx"
    print("=== 模板文件 ===")
    analyze_ppt(template_path)

    # 分析生成的文件
    generated_path = "test_output/从百草园到三味书屋课文解析.pptx"
    print("\n\n=== 生成的文件 ===")
    analyze_ppt(generated_path)


if __name__ == "__main__":
    main()
