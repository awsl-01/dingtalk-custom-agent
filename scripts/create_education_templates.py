"""
创建教育行业PPT模板
为14种风格创建真实的PPTX模板文件
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 模板输出目录
TEMPLATE_DIR = "ppt-master/skills/ppt-master/templates/education"

def hex_to_rgb(hex_color):
    """将十六进制颜色转换为RGB"""
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def create_ink_painting_template():
    """创建国风水墨模板"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 颜色定义
    INK_BLACK = hex_to_rgb('#2C3E50')
    BROWN = hex_to_rgb('#8B4513')
    VERMILION = hex_to_rgb('#C0392B')
    RICE_PAPER = hex_to_rgb('#F5F5DC')
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    # ====== 封面页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RICE_PAPER

    # 添加水墨装饰条（顶部）
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.8)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = INK_BLACK
    top_bar.line.fill.background()

    # 添加水墨装饰条（底部）
    bottom_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(6.7), Inches(13.333), Inches(0.8)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = INK_BLACK
    bottom_bar.line.fill.background()

    # 添加印章装饰（右侧）
    seal = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(11), Inches(2), Inches(1.5), Inches(1.5)
    )
    seal.fill.solid()
    seal.fill.fore_color.rgb = VERMILION
    seal.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(8), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "国风水墨教学课件"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = INK_BLACK
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(8), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "古典雅致 · 文化传承"
    p.font.size = Pt(28)
    p.font.color.rgb = BROWN
    p.alignment = PP_ALIGN.CENTER

    # ====== 目录页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RICE_PAPER

    # 左侧装饰
    left_deco = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.5), Inches(7.5)
    )
    left_deco.fill.solid()
    left_deco.fill.fore_color.rgb = INK_BLACK
    left_deco.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.5), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "目录"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = INK_BLACK
    p.alignment = PP_ALIGN.LEFT

    # 目录项
    toc_items = ["壹 · 课程导入", "贰 · 知识讲解", "叁 · 课堂练习", "肆 · 课堂小结"]
    for i, item in enumerate(toc_items):
        y = 2 + i * 1.2
        # 竖线装饰
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.5), Inches(y), Inches(0.08), Inches(0.8)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = VERMILION
        line.line.fill.background()

        # 文字
        text_box = slide.shapes.add_textbox(Inches(2), Inches(y), Inches(8), Inches(0.8))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(28)
        p.font.color.rgb = INK_BLACK

    # ====== 内容页模板 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RICE_PAPER

    # 顶部装饰
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.08)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = INK_BLACK
    top_bar.line.fill.background()

    # 左侧装饰条
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(0.8), Inches(0.08), Inches(5.9)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = VERMILION
    left_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.5), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "知识讲解"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = INK_BLACK
    p.alignment = PP_ALIGN.LEFT

    # 内容区
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    points = [
        "• 知识点一：古典文学的基础概念",
        "• 知识点二：诗词格律与韵脚分析",
        "• 知识点三：意境营造与修辞手法"
    ]

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(22)
        p.font.color.rgb = INK_BLACK
        p.space_after = Pt(20)

    # ====== 练习页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RICE_PAPER

    # 标题区域
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(1.5)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = INK_BLACK
    title_bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "课堂练习"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 练习内容
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    exercises = [
        "一、填空题",
        "1. 《________》是唐代诗人李白的代表作之一。",
        "",
        "二、简答题",
        "2. 请简述这首诗表达的主要情感。",
        "",
        "三、思考题",
        "3. 结合诗句，分析诗中运用的修辞手法。"
    ]

    for i, text in enumerate(exercises):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(20)
        p.font.color.rgb = INK_BLACK
        p.space_after = Pt(8)

    # ====== 结束页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = INK_BLACK

    # 感谢文字
    thank_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]
    p.text = "谢谢观看"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "学无止境 · 精益求精"
    p.font.size = Pt(24)
    p.font.color.rgb = hex_to_rgb('#CCCCCC')
    p.alignment = PP_ALIGN.CENTER

    # 保存
    output_path = os.path.join(TEMPLATE_DIR, "ink_painting.pptx")
    prs.save(output_path)
    print(f"已创建: {output_path}")
    return output_path


def create_macaron_cartoon_template():
    """创建马卡龙卡通模板"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 颜色定义
    PINK = hex_to_rgb('#FF6B9D')
    LAVENDER = hex_to_rgb('#C9B1FF')
    LEMON = hex_to_rgb('#FFD93D')
    LIGHT_PINK = hex_to_rgb('#FFF5F5')
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    # ====== 封面页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_PINK

    # 顶部装饰条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.3)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PINK
    top_bar.line.fill.background()

    # 底部装饰条
    bottom_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.2), Inches(13.333), Inches(0.3)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = LAVENDER
    bottom_bar.line.fill.background()

    # 装饰圆形
    circle1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(1), Inches(1), Inches(1.5), Inches(1.5)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = LEMON
    circle1.line.fill.background()

    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(11), Inches(5.5), Inches(1.2), Inches(1.2)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = LAVENDER
    circle2.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "趣味学习乐园"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = PINK
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "快乐学习 · 健康成长"
    p.font.size = Pt(26)
    p.font.color.rgb = LAVENDER
    p.alignment = PP_ALIGN.CENTER

    # ====== 内容页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # 左侧装饰
    left_deco = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.4), Inches(7.5)
    )
    left_deco.fill.solid()
    left_deco.fill.fore_color.rgb = PINK
    left_deco.line.fill.background()

    # 标题背景
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(0.5), Inches(4), Inches(1)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = PINK
    title_bg.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.6), Inches(3.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "今日学习"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 内容卡片
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(2), Inches(10), Inches(4.5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_PINK
    card.line.color.rgb = PINK

    # 内容
    content_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(3.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    points = [
        "🌈 知识点一：基础概念学习",
        "⭐ 知识点二：趣味互动练习",
        "🎨 知识点三：创意手工制作",
        "🎵 知识点四：儿歌律动时间"
    ]

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(24)
        p.font.color.rgb = hex_to_rgb('#333333')
        p.space_after = Pt(16)

    # ====== 练习页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_PINK

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🎮 趣味练习"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PINK
    p.alignment = PP_ALIGN.CENTER

    # 练习卡片
    for i in range(3):
        x = 1 + i * 4
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(2), Inches(3.5), Inches(4.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LAVENDER

        # 题号
        num_box = slide.shapes.add_textbox(Inches(x + 0.5), Inches(2.5), Inches(2.5), Inches(0.8))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"第{i+1}题"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = PINK
        p.alignment = PP_ALIGN.CENTER

        # 题目
        q_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(3.5), Inches(3), Inches(2.5))
        tf = q_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "请完成以下练习题目..."
        p.font.size = Pt(18)
        p.font.color.rgb = hex_to_rgb('#666666')

    # ====== 结束页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PINK

    # 装饰圆形
    for pos in [(1, 1), (11, 5), (2, 6), (10, 1)]:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(pos[0]), Inches(pos[1]), Inches(0.8), Inches(0.8)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = LEMON
        circle.line.fill.background()

    # 感谢文字
    thank_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]
    p.text = "太棒了！"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "继续加油哦~"
    p.font.size = Pt(28)
    p.font.color.rgb = LEMON
    p.alignment = PP_ALIGN.CENTER

    # 保存
    output_path = os.path.join(TEMPLATE_DIR, "macaron_cartoon.pptx")
    prs.save(output_path)
    print(f"已创建: {output_path}")
    return output_path


def create_fresh_forest_template():
    """创建清新森系模板"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 颜色定义
    FOREST_GREEN = hex_to_rgb('#2ECC71')
    LIGHT_GREEN = hex_to_rgb('#87D37C')
    SUNSHINE = hex_to_rgb('#F39C12')
    LIGHT_GRAY = hex_to_rgb('#F8F9FA')
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = hex_to_rgb('#2C3E50')

    # ====== 封面页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_GRAY

    # 顶部绿色装饰
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(1.5)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = FOREST_GREEN
    top_bar.line.fill.background()

    # 叶子装饰（圆形模拟）
    leaf1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(0.5), Inches(0.3), Inches(0.6), Inches(0.6)
    )
    leaf1.fill.solid()
    leaf1.fill.fore_color.rgb = LIGHT_GREEN
    leaf1.line.fill.background()

    leaf2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(12), Inches(0.5), Inches(0.8), Inches(0.8)
    )
    leaf2.fill.solid()
    leaf2.fill.fore_color.rgb = SUNSHINE
    leaf2.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "自然科学探索"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = FOREST_GREEN
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4.8), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "探索自然 · 发现奥秘"
    p.font.size = Pt(26)
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER

    # ====== 内容页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # 左侧绿色装饰
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.3), Inches(7.5)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = FOREST_GREEN
    left_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🌿 知识讲解"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = FOREST_GREEN
    p.alignment = PP_ALIGN.LEFT

    # 内容卡片
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.8), Inches(11), Inches(5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = hex_to_rgb('#F0FFF0')
    card.line.color.rgb = LIGHT_GREEN

    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(10), Inches(4))
    tf = content_box.text_frame
    tf.word_wrap = True

    points = [
        "🌱 生命的起源与演化",
        "🌳 生态系统的平衡",
        "🍃 植物的光合作用",
        "🦋 动物的生存智慧"
    ]

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(24)
        p.font.color.rgb = DARK
        p.space_after = Pt(20)

    # ====== 结束页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = FOREST_GREEN

    thank_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢聆听"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "保护自然 · 从我做起"
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_GREEN
    p.alignment = PP_ALIGN.CENTER

    # 保存
    output_path = os.path.join(TEMPLATE_DIR, "fresh_forest.pptx")
    prs.save(output_path)
    print(f"已创建: {output_path}")
    return output_path


def create_tech_neon_template():
    """创建科技霓虹模板"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 颜色定义
    DARK_BG = hex_to_rgb('#0D0D0D')
    NEON_GREEN = hex_to_rgb('#00FF88')
    NEON_BLUE = hex_to_rgb('#00D4FF')
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    # ====== 封面页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # 霓虹装饰线
    line1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.05)
    )
    line1.fill.solid()
    line1.fill.fore_color.rgb = NEON_GREEN
    line1.line.fill.background()

    line2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.45), Inches(13.333), Inches(0.05)
    )
    line2.fill.solid()
    line2.fill.fore_color.rgb = NEON_BLUE
    line2.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "科技未来"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = NEON_GREEN
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(2), Inches(4.8), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "INNOVATION · TECHNOLOGY"
    p.font.size = Pt(24)
    p.font.color.rgb = NEON_BLUE
    p.alignment = PP_ALIGN.CENTER

    # ====== 内容页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "// 知识讲解"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NEON_GREEN
    p.alignment = PP_ALIGN.LEFT

    # 内容
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    points = [
        ">>> 核心概念解析",
        ">>> 技术原理演示",
        ">>> 实践应用案例",
        ">>> 未来发展趋势"
    ]

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.space_after = Pt(20)

    # ====== 结束页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    thank_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]
    p.text = "THANKS"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = NEON_GREEN
    p.alignment = PP_ALIGN.CENTER

    # 保存
    output_path = os.path.join(TEMPLATE_DIR, "tech_neon.pptx")
    prs.save(output_path)
    print(f"已创建: {output_path}")
    return output_path


def create_warm_healing_template():
    """创建暖黄治愈风模板"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 颜色定义
    WARM_YELLOW = hex_to_rgb('#F39C12')
    ORANGE = hex_to_rgb('#E67E22')
    GREEN = hex_to_rgb('#27AE60')
    LIGHT_YELLOW = hex_to_rgb('#FFF9E6')
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = hex_to_rgb('#2C3E50')

    # ====== 封面页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_YELLOW

    # 顶部装饰
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.4)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = WARM_YELLOW
    top_bar.line.fill.background()

    # 装饰圆形
    circle1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(1), Inches(1), Inches(1), Inches(1)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = WARM_YELLOW
    circle1.line.fill.background()

    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(11.5), Inches(6), Inches(0.8), Inches(0.8)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = GREEN
    circle2.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "温暖成长课堂"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WARM_YELLOW
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(2), Inches(4.8), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "用心陪伴 · 快乐成长"
    p.font.size = Pt(26)
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER

    # ====== 内容页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # 左侧装饰
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.3), Inches(7.5)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = WARM_YELLOW
    left_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "今日主题"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WARM_YELLOW
    p.alignment = PP_ALIGN.LEFT

    # 内容
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    points = [
        "☀ 分享与倾听",
        "☀ 团队合作",
        "☀ 积极心态",
        "☀ 感恩与成长"
    ]

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(24)
        p.font.color.rgb = DARK
        p.space_after = Pt(20)

    # ====== 结束页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WARM_YELLOW

    thank_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]
    p.text = "一起加油！"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 保存
    output_path = os.path.join(TEMPLATE_DIR, "warm_healing.pptx")
    prs.save(output_path)
    print(f"已创建: {output_path}")
    return output_path


def create_hand_drawn_template():
    """创建手绘插画模板"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 颜色定义
    SKY_BLUE = hex_to_rgb('#4A90E2')
    GRASS_GREEN = hex_to_rgb('#7ED321')
    ORANGE = hex_to_rgb('#F5A623')
    CREAM = hex_to_rgb('#FFFEF7')
    DARK = hex_to_rgb('#333333')
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    # ====== 封面页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = CREAM

    # 手绘风格装饰（不规则圆形）
    shapes_data = [
        (0.5, 0.5, 1.2, 1.2, SKY_BLUE),
        (11.5, 6, 1, 1, GRASS_GREEN),
        (1, 6.5, 0.8, 0.8, ORANGE),
    ]
    for x, y, w, h, color in shapes_data:
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "创意手绘课堂"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = SKY_BLUE
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(2), Inches(4.8), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "自由创作 · 快乐表达"
    p.font.size = Pt(26)
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER

    # ====== 内容页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = CREAM

    # 便签纸效果
    note = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1), Inches(11), Inches(5.5)
    )
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xE0)
    note.line.color.rgb = ORANGE

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "今日创作"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = SKY_BLUE

    # 内容
    content_box = slide.shapes.add_textbox(Inches(2), Inches(2.8), Inches(9), Inches(3))
    tf = content_box.text_frame
    tf.word_wrap = True

    points = [
        "涂鸦时间：自由绘画",
        "手工制作：创意折纸",
        "故事创作：看图说话",
        "作品展示：分享交流"
    ]

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"✦ {point}"
        p.font.size = Pt(22)
        p.font.color.rgb = DARK
        p.space_after = Pt(16)

    # ====== 结束页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = SKY_BLUE

    thank_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]
    p.text = "创作快乐！"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    output_path = os.path.join(TEMPLATE_DIR, "hand_drawn.pptx")
    prs.save(output_path)
    print(f"已创建: {output_path}")
    return output_path


def create_minimalist_ins_template():
    """创建极简ins风模板"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 颜色定义
    DARK_GRAY = hex_to_rgb('#2C3E50')
    LIGHT_GRAY = hex_to_rgb('#ECF0F1')
    RED_ACCENT = hex_to_rgb('#E74C3C')
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    # ====== 封面页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # 极简装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3), Inches(3.2), Inches(7.333), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RED_ACCENT
    line.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "专业教学"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "简约 · 高效"
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # ====== 内容页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "核心要点"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY

    # 装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(1.7), Inches(2), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RED_ACCENT
    line.line.fill.background()

    # 内容 - 大留白设计
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(4))
    tf = content_box.text_frame
    tf.word_wrap = True

    points = [
        "01  基础概念",
        "02  核心原理",
        "03  实践应用"
    ]

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(28)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(24)

    # ====== 结束页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_GRAY

    thank_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]
    p.text = "THANKS"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    output_path = os.path.join(TEMPLATE_DIR, "minimalist_ins.pptx")
    prs.save(output_path)
    print(f"已创建: {output_path}")
    return output_path


def create_lab_tech_template():
    """创建实验室科技风模板"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 颜色定义
    TEAL = hex_to_rgb('#1ABC9C')
    BLUE = hex_to_rgb('#3498DB')
    RED = hex_to_rgb('#E74C3C')
    LIGHT_GRAY = hex_to_rgb('#ECF0F1')
    DARK = hex_to_rgb('#2C3E50')
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    # ====== 封面页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_GRAY

    # 顶部装饰
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.8)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = TEAL
    top_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "科学实验室"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(2), Inches(4.8), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "探索 · 实验 · 发现"
    p.font.size = Pt(26)
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER

    # ====== 内容页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # 左侧装饰
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.3), Inches(7.5)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = TEAL
    left_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "实验步骤"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEAL

    # 内容
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    points = [
        "1. 准备实验器材",
        "2. 记录初始数据",
        "3. 进行实验操作",
        "4. 观察实验现象",
        "5. 分析实验结果"
    ]

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(22)
        p.font.color.rgb = DARK
        p.space_after = Pt(16)

    # ====== 结束页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = TEAL

    thank_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]
    p.text = "实验完成！"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    output_path = os.path.join(TEMPLATE_DIR, "lab_tech.pptx")
    prs.save(output_path)
    print(f"已创建: {output_path}")
    return output_path


def create_sports_energy_template():
    """创建活力运动风模板"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 颜色定义
    ORANGE_RED = hex_to_rgb('#FF5722')
    GREEN = hex_to_rgb('#4CAF50')
    BLUE = hex_to_rgb('#2196F3')
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = hex_to_rgb('#333333')

    # ====== 封面页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = ORANGE_RED

    # 装饰
    shape1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(0.5), Inches(0.5), Inches(1.5), Inches(1.5)
    )
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = GREEN
    shape1.line.fill.background()

    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(11), Inches(5.5), Inches(1.5), Inches(1.5)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = BLUE
    shape2.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "活力运动"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(2), Inches(4.8), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "健康 · 活力 · 快乐"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # ====== 内容页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # 顶部装饰
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.3)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = ORANGE_RED
    top_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "今日运动"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ORANGE_RED

    # 内容
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    points = [
        "热身运动：5分钟",
        "主要活动：20分钟",
        "团队游戏：10分钟",
        "放松拉伸：5分钟"
    ]

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"● {point}"
        p.font.size = Pt(24)
        p.font.color.rgb = DARK
        p.space_after = Pt(20)

    # ====== 结束页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = GREEN

    thank_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]
    p.text = "运动快乐！"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    output_path = os.path.join(TEMPLATE_DIR, "sports_energy.pptx")
    prs.save(output_path)
    print(f"已创建: {output_path}")
    return output_path


def create_music_rhythm_template():
    """创建音乐律动风模板"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 颜色定义
    PURPLE = hex_to_rgb('#9C27B0')
    PINK = hex_to_rgb('#E91E63')
    YELLOW = hex_to_rgb('#FFEB3B')
    LIGHT_PURPLE = hex_to_rgb('#F3E5F5')
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    # ====== 封面页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_PURPLE

    # 音符装饰（圆形模拟）
    notes = [(1, 1, YELLOW), (11.5, 0.5, PINK), (0.5, 6, PURPLE), (12, 6, YELLOW)]
    for x, y, color in notes:
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.6), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "音乐课堂"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(2), Inches(4.8), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "感受音乐 · 享受艺术"
    p.font.size = Pt(26)
    p.font.color.rgb = PINK
    p.alignment = PP_ALIGN.CENTER

    # ====== 内容页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "音乐知识"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    # 内容
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    points = [
        "音符与节拍",
        "旋律与和声",
        "乐器介绍",
        "音乐欣赏"
    ]

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"♫ {point}"
        p.font.size = Pt(24)
        p.font.color.rgb = PURPLE
        p.space_after = Pt(20)

    # ====== 结束页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PURPLE

    thank_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]
    p.text = "音乐相伴！"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    output_path = os.path.join(TEMPLATE_DIR, "music_rhythm.pptx")
    prs.save(output_path)
    print(f"已创建: {output_path}")
    return output_path


def create_all_templates():
    """创建所有模板"""
    os.makedirs(TEMPLATE_DIR, exist_ok=True)

    print("开始创建教育PPT模板...")
    print("=" * 50)

    templates = [
        ("ink_painting", create_ink_painting_template),
        ("macaron_cartoon", create_macaron_cartoon_template),
        ("fresh_forest", create_fresh_forest_template),
        ("tech_neon", create_tech_neon_template),
        ("warm_healing", create_warm_healing_template),
        ("hand_drawn", create_hand_drawn_template),
        ("minimalist_ins", create_minimalist_ins_template),
        ("lab_tech", create_lab_tech_template),
        ("sports_energy", create_sports_energy_template),
        ("music_rhythm", create_music_rhythm_template),
    ]

    created = []
    for name, func in templates:
        try:
            path = func()
            created.append((name, path))
            print(f"[OK] {name} 模板创建成功")
        except Exception as e:
            print(f"[FAIL] {name} 模板创建失败: {e}")

    print("=" * 50)
    print(f"成功创建 {len(created)} 个模板")

    return created


if __name__ == "__main__":
    create_all_templates()
