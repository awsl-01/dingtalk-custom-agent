"""
创建示例模板文件
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


def create_template(output_path: str):
    """创建包含封面页和目录页的模板"""
    prs = Presentation()

    # 设置16:9宽屏比例
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 定义设计基因
    PRIMARY_COLOR = RGBColor(0x1A, 0x56, 0xDB)  # 主色调 - 深蓝
    ACCENT_COLOR = RGBColor(0xFF, 0x6B, 0x35)    # 强调色 - 橙色
    TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)       # 文字颜色 - 深灰
    LIGHT_TEXT = RGBColor(0xFF, 0xFF, 0xFF)       # 浅色文字 - 白色
    SUBTITLE_COLOR = RGBColor(0xCC, 0xDD, 0xEE)  # 副标题颜色

    # ====== 封面页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR

    # 装饰线条
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(3.2), Inches(2), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "演示文稿标题"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = LIGHT_TEXT
    p.font.name = "微软雅黑"
    p.alignment = PP_ALIGN.LEFT

    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(1))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "副标题说明文字"
    p.font.size = Pt(24)
    p.font.color.rgb = SUBTITLE_COLOR
    p.font.name = "微软雅黑"
    p.alignment = PP_ALIGN.LEFT

    # ====== 目录页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 浅色背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)

    # 顶部装饰条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.08)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PRIMARY_COLOR
    top_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "目录"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.font.name = "微软雅黑"
    p.alignment = PP_ALIGN.LEFT

    # 左侧装饰线
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(1.4), Inches(0.08), Inches(0.5)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = ACCENT_COLOR
    accent_line.line.fill.background()

    # 目录项
    toc_items = ["第一章 概述", "第二章 核心内容", "第三章 详细分析", "第四章 总结与展望"]
    content_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.5), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(toc_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"•  {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = "微软雅黑"
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(16)

    # 保存模板
    prs.save(output_path)
    print(f"模板已创建: {output_path}")


if __name__ == "__main__":
    create_template("templates/default_template.pptx")
