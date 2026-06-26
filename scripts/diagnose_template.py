"""
诊断模板PPT生成问题
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 强制使用UTF-8编码
sys.stdout.reconfigure(encoding='utf-8')
sys.stderr.reconfigure(encoding='utf-8')


def diagnose_template(template_path: str):
    """诊断模板文件"""
    print(f"=== 诊断模板: {template_path} ===")

    if not os.path.exists(template_path):
        print(f"错误：模板文件不存在")
        return

    prs = Presentation(template_path)

    print(f"幻灯片数量: {len(prs.slides)}")
    print(f"幻灯片宽度: {prs.slide_width / 914400:.3f} 英寸")
    print(f"幻灯片高度: {prs.slide_height / 914400:.3f} 英寸")
    print()

    for i, slide in enumerate(prs.slides):
        print(f"--- 第 {i+1} 张幻灯片 ---")

        # 背景
        try:
            bg = slide.background
            fill = bg.fill
            if fill.type is not None:
                print(f"  背景类型: {fill.type}")
                if fill.type == 1:  # SOLID_FILL
                    print(f"  背景颜色: {fill.fore_color.rgb}")
        except Exception as e:
            print(f"  背景提取失败: {e}")

        # 形状
        print(f"  形状数量: {len(slide.shapes)}")
        for j, shape in enumerate(slide.shapes):
            print(f"    形状 {j+1}:")
            print(f"      类型: {shape.shape_type}")
            print(f"      位置: ({shape.left / 914400:.3f}, {shape.top / 914400:.3f}) 英寸")
            print(f"      尺寸: ({shape.width / 914400:.3f} x {shape.height / 914400:.3f}) 英寸")

            if shape.has_text_frame:
                tf = shape.text_frame
                print(f"      文本框:")
                for k, para in enumerate(tf.paragraphs):
                    if para.text.strip():
                        print(f"        段落 {k+1}: '{para.text[:50]}...'")
                        if para.runs:
                            run = para.runs[0]
                            if run.font.name:
                                print(f"          字体: {run.font.name}")
                            if run.font.size:
                                print(f"          字号: {run.font.size.pt}pt")
                            try:
                                if run.font.color and run.font.color.type is not None:
                                    print(f"          颜色: {run.font.color.rgb}")
                            except:
                                pass
                            if run.font.bold is not None:
                                print(f"          加粗: {run.font.bold}")

        print()


def diagnose_generated_ppt(ppt_path: str):
    """诊断生成的PPT"""
    print(f"=== 诊断生成的PPT: {ppt_path} ===")

    if not os.path.exists(ppt_path):
        print(f"错误：文件不存在")
        return

    prs = Presentation(ppt_path)

    print(f"幻灯片数量: {len(prs.slides)}")

    layout_types = []
    overflow_issues = []

    for i, slide in enumerate(prs.slides):
        print(f"\n--- 第 {i+1} 张幻灯片 ---")

        # 检查背景
        bg_color = None
        try:
            bg = slide.background
            fill = bg.fill
            if fill.type == 1:  # SOLID_FILL
                bg_color = fill.fore_color.rgb
                print(f"  背景颜色: {bg_color}")
        except:
            pass

        # 推断版式类型
        layout_type = "未知"
        if i == 0:
            layout_type = "封面"
        elif i == len(prs.slides) - 1:
            layout_type = "结束页"
        else:
            # 检查是否有半透明遮罩（章节页特征）
            has_overlay = False
            for shape in slide.shapes:
                if shape.shape_type == 1:  # RECTANGLE
                    width = shape.width / 914400
                    height = shape.height / 914400
                    if width > 10 and height > 5:  # 大矩形可能是遮罩
                        has_overlay = True
                        break

            if has_overlay:
                layout_type = "章节页"
            else:
                # 检查是否有图片占位符
                has_image_placeholder = False
                for shape in slide.shapes:
                    if shape.shape_type == 1:  # RECTANGLE
                        width = shape.width / 914400
                        if width > 4 and width < 6:  # 可能是图片占位符
                            has_image_placeholder = True
                            break

                if has_image_placeholder:
                    layout_type = "图文页"
                else:
                    layout_type = "正文页"

        layout_types.append(layout_type)
        print(f"  版式类型: {layout_type}")

        # 检查文字溢出
        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                if tf.paragraphs:
                    # 计算文本总高度
                    total_text_height = 0
                    for para in tf.paragraphs:
                        if para.text.strip():
                            # 粗略估算：每行约0.3英寸
                            lines = max(1, len(para.text) // 40 + 1)
                            total_text_height += lines * 0.3

                    box_height = shape.height / 914400
                    if total_text_height > box_height * 0.9:  # 超过90%视为可能溢出
                        overflow_issues.append(i + 1)
                        print(f"  ⚠️ 可能存在文字溢出 (文本高度: {total_text_height:.2f}英寸, 框高度: {box_height:.2f}英寸)")

    # 检查版式多样性
    print("\n=== 版式多样性检查 ===")
    consecutive_same = 0
    prev_type = None
    for i, lt in enumerate(layout_types):
        if lt == prev_type:
            consecutive_same += 1
            if consecutive_same >= 2:
                print(f"⚠️ 第 {i-1}-{i+1} 页连续使用相同版式: {lt}")
        else:
            consecutive_same = 0
        prev_type = lt

    if overflow_issues:
        print(f"\n⚠️ 存在文字溢出问题的页: {overflow_issues}")
    else:
        print("\n✓ 未发现明显文字溢出问题")

    print(f"\n版式分布: {dict((x, layout_types.count(x)) for x in set(layout_types))}")


def test_simple_generation():
    """测试简单的PPT生成"""
import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    from agent.template_design_ppt import TemplateDesignPPTGenerator

    template_path = "templates/default_template.pptx"
    output_dir = "test_output"

    if not os.path.exists(template_path):
        print(f"模板文件不存在: {template_path}")
        return

    print("\n=== 测试简单PPT生成 ===")

    # 创建一个简单的测试内容
    test_data = {
        "title": "测试PPT",
        "subtitle": "用于诊断问题",
        "slides": [
            {
                "title": "第一章 概述",
                "content": ["要点1", "要点2", "要点3"],
                "layout_type": "chapter"
            },
            {
                "title": "核心内容",
                "content": ["详细说明1", "详细说明2", "详细说明3", "详细说明4"],
                "layout_type": "content"
            },
            {
                "title": "图文展示",
                "content": ["左侧内容1", "左侧内容2"],
                "layout_type": "image_text"
            },
            {
                "title": "金句",
                "content": ["这是一句很重要的话", "出处"],
                "layout_type": "quote"
            },
            {
                "title": "总结",
                "content": ["总结要点1", "总结要点2", "总结要点3"],
                "layout_type": "content"
            }
        ]
    }

    # 手动创建PPT来测试
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 加载设计基因
    from agent.template_design_ppt import TemplateDesignExtractor, DesignGene, SlideBuilder
    extractor = TemplateDesignExtractor(template_path)
    gene = extractor.extract()

    print(f"设计基因提取完成:")
    print(f"  背景颜色: {gene.background_color}")
    print(f"  标题字体: {gene.title_font_name}")
    print(f"  标题字号: {gene.title_font_size}")
    print(f"  正文字体: {gene.body_font_name}")
    print(f"  正文字号: {gene.body_font_size}")
    print(f"  强调色: {gene.accent_color}")

    # 使用SlideBuilder构建
    builder = SlideBuilder(gene)

    from agent.template_design_ppt import SlideContent

    # 创建封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = gene.accent_color

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = test_data["title"]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.name = gene.title_font_name
    p.alignment = PP_ALIGN.LEFT

    # 创建内容页
    for i, slide_data in enumerate(test_data["slides"]):
        content = SlideContent(
            title=slide_data["title"],
            content=slide_data["content"],
            layout_type=slide_data["layout_type"]
        )
        builder.build_slide(prs, content, i + 2)

    # 创建结束页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = gene.accent_color

    thank_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = thank_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "谢谢观看"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.name = gene.title_font_name
    p.alignment = PP_ALIGN.CENTER

    # 保存
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "test_diagnosis.pptx")
    prs.save(output_path)

    print(f"\n测试PPT已生成: {output_path}")
    return output_path


if __name__ == "__main__":
    # 诊断模板
    diagnose_template("templates/default_template.pptx")

    print("\n" + "="*50 + "\n")

    # 测试生成
    output_path = test_simple_generation()

    if output_path:
        print("\n" + "="*50 + "\n")
        # 诊断生成的PPT
        diagnose_generated_ppt(output_path)
