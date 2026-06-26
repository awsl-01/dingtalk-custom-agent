"""
快速测试模板PPT生成器
"""
import os
import sys
import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from agent.template_based_ppt import generate_template_based_ppt

sys.stdout.reconfigure(encoding='utf-8')
sys.stderr.reconfigure(encoding='utf-8')


def main():
    """主函数"""
    print("=== 模板PPT生成器快速测试 ===\n")

    # 检查模板文件
    template_path = "templates/default_template.pptx"
    if not os.path.exists(template_path):
        print(f"错误：模板文件不存在: {template_path}")
        print("请先运行 create_template.py 创建默认模板")
        return

    # 用户输入
    print("请输入PPT主题（或按回车使用默认主题）：")
    user_input = input().strip()

    if not user_input:
        user_input = "生成一个从百草园到三味书屋适合初中的语文课件PPT"
        print(f"使用默认主题: {user_input}")

    # 输出目录
    output_dir = "projects"
    os.makedirs(output_dir, exist_ok=True)

    try:
        print(f"\n正在生成PPT...")
        output_path, title = generate_template_based_ppt(
            user_message=user_input,
            template_path=template_path,
            output_dir=output_dir
        )

        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            print(f"\n✓ 生成成功！")
            print(f"  文件路径: {output_path}")
            print(f"  标题: {title}")
            print(f"  文件大小: {file_size / 1024:.2f} KB")

            # 显示幻灯片内容
            from pptx import Presentation
            prs = Presentation(output_path)
            print(f"  幻灯片数量: {len(prs.slides)}")

            print(f"\n=== 幻灯片内容预览 ===")
            for i, slide in enumerate(prs.slides):
                print(f"\n第{i+1}页:")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                text = para.text[:60]
                                if len(para.text) > 60:
                                    text += "..."
                                print(f"  - {text}")
                                break
        else:
            print(f"\n✗ 文件不存在: {output_path}")

    except Exception as e:
        print(f"\n✗ 生成失败: {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main()
