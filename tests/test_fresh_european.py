"""
测试清新欧美风模板
"""
import os
import sys
import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from agent.template_based_ppt import generate_template_based_ppt

sys.stdout.reconfigure(encoding='utf-8')
sys.stderr.reconfigure(encoding='utf-8')


def main():
    """主函数"""
    print("=== 测试清新欧美风模板 ===\n")

    # 模板路径
    template_path = "templates/清新欧美风.pptx"

    if not os.path.exists(template_path):
        print(f"错误：模板文件不存在: {template_path}")
        return

    # 先分析模板结构
    from pptx import Presentation
    prs = Presentation(template_path)

    print(f"模板信息：")
    print(f"  幻灯片数量: {len(prs.slides)}")
    print(f"  幻灯片宽度: {prs.slide_width / 914400:.3f} 英寸")
    print(f"  幻灯片高度: {prs.slide_height / 914400:.3f} 英寸")

    for i, slide in enumerate(prs.slides):
        print(f"\n  第{i+1}页:")
        try:
            bg = slide.background
            fill = bg.fill
            if fill.type is not None:
                if fill.type == 1:  # SOLID_FILL
                    print(f"    背景颜色: {fill.fore_color.rgb}")
        except:
            print(f"    背景: 无法读取")

        print(f"    形状数量: {len(slide.shapes)}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        print(f"    - {para.text[:30]}...")
                        break

    # 生成PPT
    print("\n\n=== 生成PPT ===")
    user_message = "生成一个从百草园到三味书屋适合初中的语文课件PPT，包含：1. 课文背景 2. 生字词 3. 段落分析 4. 主题思想 5. 课后练习"

    output_dir = "test_output"
    os.makedirs(output_dir, exist_ok=True)

    try:
        output_path, title = generate_template_based_ppt(
            user_message=user_message,
            template_path=template_path,
            output_dir=output_dir
        )

        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            print(f"\n✓ 生成成功！")
            print(f"  文件路径: {output_path}")
            print(f"  标题: {title}")
            print(f"  文件大小: {file_size / 1024:.2f} KB")

            # 验证生成的PPT
            prs = Presentation(output_path)
            print(f"  幻灯片数量: {len(prs.slides)}")

            print(f"\n=== 幻灯片内容预览 ===")
            for i, slide in enumerate(prs.slides):
                print(f"\n第{i+1}页:")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                text = para.text[:50]
                                if len(para.text) > 50:
                                    text += "..."
                                print(f"  - {text}")
                                break
        else:
            print(f"\n✗ 文件不存在: {output_path}")

    except Exception as e:
        print(f"\n✗ 生成失败: {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main()
