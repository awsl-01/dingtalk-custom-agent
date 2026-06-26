"""
测试字体大小检测
"""
import os
import sys
import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout.reconfigure(encoding='utf-8')
sys.stderr.reconfigure(encoding='utf-8')

# 创建一个简单的PPT
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 添加一张幻灯片
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 添加标题
title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(1))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "测试标题"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# 添加正文
content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
tf = content_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "测试正文内容"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

# 保存
output_path = "test_output/font_size_test.pptx"
os.makedirs("test_output", exist_ok=True)
prs.save(output_path)

print(f"已生成: {output_path}")

# 重新读取并检查
prs2 = Presentation(output_path)
slide = prs2.slides[0]

print("\n=== 检查字体大小 ===")
for shape in slide.shapes:
    if shape.has_text_frame:
        print(f"\n形状: {shape.shape_type}")
        for para in shape.text_frame.paragraphs:
            if para.text.strip():
                print(f"  文本: {para.text}")
                if para.runs:
                    run = para.runs[0]
                    print(f"    字体: {run.font.name}")
                    print(f"    字号: {run.font.size}")
                    if run.font.size:
                        print(f"    字号(pt): {run.font.size.pt}")
                else:
                    print(f"    没有runs")
