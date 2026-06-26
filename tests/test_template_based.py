"""
测试基于模板的PPT生成器（直接复制模板）
"""
import os
import sys
import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from agent.template_based_ppt import generate_template_based_ppt

sys.stdout.reconfigure(encoding='utf-8')
sys.stderr.reconfigure(encoding='utf-8')


def test_template_based():
    """测试基于模板的PPT生成"""
    template_path = "templates/default_template.pptx"
    output_dir = "test_output"

    if not os.path.exists(template_path):
        print(f"错误：模板文件不存在: {template_path}")
        return

    print("=== 测试基于模板的PPT生成 ===")
    print("特点：直接复制模板幻灯片，保留所有格式\n")

    # 测试内容
    test_message = "生成一个从百草园到三味书屋适合初中的语文课件PPT，包含：1. 课文背景 2. 生字词 3. 段落分析 4. 主题思想 5. 课后练习"

    try:
        output_path, title = generate_template_based_ppt(
            user_message=test_message,
            template_path=template_path,
            output_dir=output_dir
        )

        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            print(f"✓ 生成成功: {output_path}")
            print(f"  标题: {title}")
            print(f"  文件大小: {file_size / 1024:.2f} KB")

            # 验证生成的PPT
            from pptx import Presentation
            prs = Presentation(output_path)
            print(f"  幻灯片数量: {len(prs.slides)}")

            # 检查每页
            for i, slide in enumerate(prs.slides):
                print(f"\n  第{i+1}页:")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                print(f"    - {para.text[:50]}...")
                                break
        else:
            print(f"✗ 文件不存在: {output_path}")

    except Exception as e:
        print(f"✗ 生成失败: {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    os.makedirs("test_output", exist_ok=True)
    test_template_based()
