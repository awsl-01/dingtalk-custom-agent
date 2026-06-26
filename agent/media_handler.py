"""
媒体文件处理模块
负责从钉钉下载图片/文件，提取文本内容
"""
import os
import json
import hashlib
import logging
import tempfile
from datetime import datetime
from typing import Optional, Tuple

import httpx
import config

logger = logging.getLogger(__name__)


async def download_media_from_dingtalk(download_code: str, save_dir: str,
                                       file_name: str = "") -> Optional[str]:
    """
    通过钉钉 downloadCode 下载媒体文件

    参数:
        download_code: 钉钉消息中的 downloadCode
        save_dir: 保存目录
        file_name: 文件名（可选）

    返回:
        保存的文件路径，失败返回 None
    """
    try:
        from dingtalk.bot import get_access_token
        token = await get_access_token()

        os.makedirs(save_dir, exist_ok=True)

        # 使用钉钉 SDK 的文件下载 API
        # 先获取下载链接
        download_url = await _get_download_url(token, download_code)
        if not download_url:
            # 备用方案：直接使用旧 API
            download_url = f"https://oapi.dingtalk.com/media/download?access_token={token}&downloadCode={download_code}"

        logger.info(f"下载文件: {download_url[:100]}...")

        async with httpx.AsyncClient() as client:
            resp = await client.get(download_url, follow_redirects=True, timeout=60)

            if resp.status_code != 200:
                logger.error(f"下载媒体文件失败: HTTP {resp.status_code}")
                # 尝试备用方案
                backup_url = f"https://oapi.dingtalk.com/media/download?access_token={token}&downloadCode={download_code}"
                resp = await client.get(backup_url, follow_redirects=True, timeout=60)
                if resp.status_code != 200:
                    logger.error(f"备用下载也失败: HTTP {resp.status_code}")
                    return None

            # 检查响应内容是否是错误信息
            content_type = resp.headers.get("content-type", "")
            if "application/json" in content_type:
                try:
                    error_data = resp.json()
                    if error_data.get("errcode", 0) != 0:
                        logger.error(f"下载失败: {error_data}")
                        return None
                except Exception:
                    pass

            # 从Content-Disposition或响应内容推断文件名
            if not file_name:
                content_disp = resp.headers.get("content-disposition", "")
                if "filename=" in content_disp:
                    file_name = content_disp.split("filename=")[-1].strip('"')
                    # 处理 URL 编码的文件名
                    try:
                        import urllib.parse
                        file_name = urllib.parse.unquote(file_name)
                    except Exception:
                        pass
                else:
                    # 根据Content-Type推断
                    ext = _content_type_to_ext(content_type)
                    file_name = f"{hashlib.md5(download_code.encode()).hexdigest()[:12]}{ext}"

            # 确保文件名有扩展名
            if '.' not in file_name:
                ext = _content_type_to_ext(content_type)
                file_name += ext

            file_path = os.path.join(save_dir, file_name)
            with open(file_path, "wb") as f:
                f.write(resp.content)

            logger.info(f"媒体文件已保存: {file_path} ({len(resp.content)} bytes)")
            return file_path

    except Exception as e:
        logger.error(f"下载媒体文件异常: {e}", exc_info=True)
        return None


async def _get_download_url(access_token: str, download_code: str) -> Optional[str]:
    """
    通过钉钉 API 获取文件下载链接

    参数:
        access_token: 访问令牌
        download_code: 下载码

    返回:
        下载链接，失败返回 None
    """
    try:
        import config
        url = "https://api.dingtalk.com/v1.0/robot/messageFiles/download"
        headers = {
            "Content-Type": "application/json",
            "x-acs-dingtalk-access-token": access_token,
        }
        data = {
            "robotCode": config.DINGTALK_APP_KEY,
            "downloadCode": download_code,
        }

        async with httpx.AsyncClient() as client:
            resp = await client.post(url, json=data, headers=headers, timeout=30)

            if resp.status_code == 200:
                result = resp.json()
                download_url = result.get("downloadUrl", "")
                if download_url:
                    logger.info(f"获取下载链接成功: {download_url[:100]}...")
                    return download_url

            logger.warning(f"获取下载链接失败: HTTP {resp.status_code}, {resp.text}")
            return None

    except Exception as e:
        logger.warning(f"获取下载链接异常: {e}")
        return None


def _content_type_to_ext(content_type: str) -> str:
    """根据Content-Type推断文件扩展名"""
    mapping = {
        "image/jpeg": ".jpg",
        "image/png": ".png",
        "image/gif": ".gif",
        "image/webp": ".webp",
        "application/pdf": ".pdf",
        "application/msword": ".doc",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
        "application/vnd.ms-excel": ".xls",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
        "application/vnd.ms-powerpoint": ".ppt",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
        "text/plain": ".txt",
    }
    for ct, ext in mapping.items():
        if ct in content_type:
            return ext
    return ".bin"


async def extract_text_from_image(image_path: str) -> str:
    """
    使用多模态LLM提取图片中的文字内容

    参数:
        image_path: 图片文件路径

    返回:
        提取的文字内容
    """
    import base64

    try:
        with open(image_path, "rb") as f:
            image_data = base64.b64encode(f.read()).decode("utf-8")

        # 根据文件扩展名确定MIME类型
        ext = os.path.splitext(image_path)[1].lower()
        mime_map = {
            ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
            ".png": "image/png", ".gif": "image/gif",
            ".webp": "image/webp", ".bmp": "image/bmp",
        }
        mime_type = mime_map.get(ext, "image/png")

        from openai import OpenAI
        client = OpenAI(
            api_key=config.OPENAI_API_KEY,
            base_url=config.OPENAI_BASE_URL,
        )

        response = client.chat.completions.create(
            model=config.OPENAI_MODEL,
            max_tokens=4096,
            messages=[
                {
                    "role": "system",
                    "content": "你是一个专业的文字识别助手。请仔细识别图片中的所有文字内容，保持原始格式和结构。"
                               "如果是课表、表格等结构化数据，请按表格格式输出。"
                               "如果是纯文本，请完整输出所有文字。不要添加任何解释或总结。"
                },
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{mime_type};base64,{image_data}"
                            }
                        },
                        {
                            "type": "text",
                            "text": "请提取图片中的所有文字内容，保持原始格式。"
                        }
                    ]
                }
            ],
        )

        text = response.choices[0].message.content
        logger.info(f"图片OCR提取完成，长度: {len(text)}")
        return text

    except Exception as e:
        logger.error(f"图片OCR提取失败: {e}")
        return ""


def extract_text_from_pdf(file_path: str) -> str:
    """从PDF文件中提取文字"""
    text_parts = []

    # 尝试 PyMuPDF
    try:
        import fitz
        doc = fitz.open(file_path)
        for page in doc:
            page_text = page.get_text()
            if page_text.strip():
                text_parts.append(page_text.strip())
        doc.close()
        if text_parts:
            return "\n\n".join(text_parts)
    except ImportError:
        pass
    except Exception as e:
        logger.warning(f"PyMuPDF提取PDF失败: {e}")

    # 备用：PyPDF2
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text and page_text.strip():
                text_parts.append(page_text.strip())
    except Exception as e:
        logger.warning(f"PyPDF2提取PDF失败: {e}")

    return "\n\n".join(text_parts)


def extract_text_from_docx(file_path: str) -> str:
    """从Word文档中提取文字"""
    try:
        from docx import Document
        doc = Document(file_path)
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)

        # 也提取表格内容
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    paragraphs.append(row_text)

        return "\n".join(paragraphs)
    except Exception as e:
        logger.error(f"提取Word文档失败: {e}")
        return ""


def extract_text_from_excel(file_path: str) -> str:
    """从Excel文件中提取文字"""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, data_only=True)
        all_text = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            all_text.append(f"【{sheet_name}】")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(
                    str(cell) if cell is not None else ""
                    for cell in row
                )
                if row_text.strip(" |"):
                    all_text.append(row_text)

        return "\n".join(all_text)
    except Exception as e:
        logger.error(f"提取Excel文件失败: {e}")
        return ""


def extract_text_from_pptx(file_path: str) -> str:
    """从PPT文件中提取文字"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        all_text = []

        for i, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- 第{i}页 ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text.append(text)
            if len(slide_text) > 1:
                all_text.append("\n".join(slide_text))

        return "\n\n".join(all_text)
    except ImportError:
        logger.warning("python-pptx未安装，无法提取PPT文本")
        return ""
    except Exception as e:
        logger.error(f"提取PPT文件失败: {e}")
        return ""


def extract_text_from_markdown(file_path: str) -> str:
    """从Markdown文件中提取文字"""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
        # 移除Markdown标记，保留纯文本
        import re
        # 移除标题标记
        content = re.sub(r'^#{1,6}\s+', '', content, flags=re.MULTILINE)
        # 移除粗体/斜体标记
        content = re.sub(r'\*{1,3}(.+?)\*{1,3}', r'\1', content)
        content = re.sub(r'_{1,3}(.+?)_{1,3}', r'\1', content)
        # 移除链接标记
        content = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', content)
        # 移除图片标记
        content = re.sub(r'!\[([^\]]*)\]\([^\)]+\)', r'\1', content)
        # 移除代码块
        content = re.sub(r'```[\s\S]*?```', '', content)
        # 移除行内代码
        content = re.sub(r'`([^`]+)`', r'\1', content)
        # 移除引用标记
        content = re.sub(r'^>\s+', '', content, flags=re.MULTILINE)
        # 移除列表标记
        content = re.sub(r'^[\s]*[-*+]\s+', '', content, flags=re.MULTILINE)
        content = re.sub(r'^[\s]*\d+\.\s+', '', content, flags=re.MULTILINE)
        # 移除水平线
        content = re.sub(r'^[-*_]{3,}\s*$', '', content, flags=re.MULTILINE)
        # 清理多余空行
        content = re.sub(r'\n{3,}', '\n\n', content)
        return content.strip()
    except UnicodeDecodeError:
        try:
            with open(file_path, "r", encoding="gbk", errors="ignore") as f:
                return f.read()
        except Exception as e:
            logger.error(f"提取Markdown文件失败: {e}")
            return ""
    except Exception as e:
        logger.error(f"提取Markdown文件失败: {e}")
        return ""


def extract_text_from_csv(file_path: str) -> str:
    """从CSV文件中提取文字"""
    try:
        import csv
        rows = []
        # 尝试不同编码
        for encoding in ['utf-8', 'gbk', 'utf-8-sig']:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    reader = csv.reader(f)
                    for row in reader:
                        row_text = " | ".join(cell.strip() for cell in row if cell.strip())
                        if row_text:
                            rows.append(row_text)
                break
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.warning(f"CSV编码 {encoding} 读取失败: {e}")
                continue

        return "\n".join(rows)
    except Exception as e:
        logger.error(f"提取CSV文件失败: {e}")
        return ""


def extract_text_from_json(file_path: str) -> str:
    """从JSON文件中提取文字"""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)

        # 递归提取所有字符串值
        def extract_strings(obj, prefix=""):
            results = []
            if isinstance(obj, dict):
                for key, value in obj.items():
                    results.extend(extract_strings(value, f"{prefix}{key}: "))
            elif isinstance(obj, list):
                for i, item in enumerate(obj):
                    results.extend(extract_strings(item, f"{prefix}[{i}]: "))
            elif isinstance(obj, str) and obj.strip():
                results.append(f"{prefix}{obj}")
            elif obj is not None:
                results.append(f"{prefix}{obj}")
            return results

        strings = extract_strings(data)
        return "\n".join(strings)
    except Exception as e:
        logger.error(f"提取JSON文件失败: {e}")
        return ""


def extract_text_from_html(file_path: str) -> str:
    """从HTML文件中提取文字"""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()

        # 尝试使用 BeautifulSoup
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(content, 'html.parser')
            # 移除脚本和样式
            for tag in soup(['script', 'style', 'nav', 'footer', 'header']):
                tag.decompose()
            text = soup.get_text(separator='\n', strip=True)
            # 清理多余空行
            import re
            text = re.sub(r'\n{3,}', '\n\n', text)
            return text.strip()
        except ImportError:
            # 简单正则提取
            import re
            # 移除HTML标签
            text = re.sub(r'<[^>]+>', ' ', content)
            # 移除多余空白
            text = re.sub(r'\s+', ' ', text)
            # 解码HTML实体
            text = text.replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>')
            text = text.replace('&quot;', '"').replace('&#39;', "'").replace('&nbsp;', ' ')
            return text.strip()
    except Exception as e:
        logger.error(f"提取HTML文件失败: {e}")
        return ""


def extract_text_from_xml(file_path: str) -> str:
    """从XML文件中提取文字"""
    try:
        import xml.etree.ElementTree as ET
        tree = ET.parse(file_path)
        root = tree.getroot()

        def extract_text(element):
            texts = []
            if element.text and element.text.strip():
                texts.append(element.text.strip())
            for child in element:
                texts.extend(extract_text(child))
            if element.tail and element.tail.strip():
                texts.append(element.tail.strip())
            return texts

        all_text = extract_text(root)
        return "\n".join(all_text)
    except Exception as e:
        logger.error(f"提取XML文件失败: {e}")
        return ""


def extract_text_from_rtf(file_path: str) -> str:
    """从RTF文件中提取文字"""
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()

        # 简单RTF解析：移除控制字
        import re
        # 移除RTF头部
        content = re.sub(r'^\{\\rtf1.*?\\pard', '', content, flags=re.DOTALL)
        # 移除控制字
        content = re.sub(r'\\[a-z]+\d*\s?', '', content)
        # 移除大括号
        content = content.replace('{', '').replace('}', '')
        # 清理空白
        content = re.sub(r'\s+', ' ', content)
        return content.strip()
    except Exception as e:
        logger.error(f"提取RTF文件失败: {e}")
        return ""


def extract_text_from_ini(file_path: str) -> str:
    """从INI/配置文件中提取文字"""
    try:
        import configparser
        config = configparser.ConfigParser()
        config.read(file_path, encoding='utf-8')

        lines = []
        for section in config.sections():
            lines.append(f"[{section}]")
            for key, value in config.items(section):
                lines.append(f"{key} = {value}")
            lines.append("")

        return "\n".join(lines)
    except Exception as e:
        logger.error(f"提取INI文件失败: {e}")
        return ""


async def extract_text_from_file(file_path: str) -> Tuple[str, str]:
    """
    从文件中提取文字（自动识别文件类型）

    参数:
        file_path: 文件路径

    返回:
        (提取的文字内容, 文件类型)
    """
    ext = os.path.splitext(file_path)[1].lower()

    # 图片文件
    if ext in (".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".tiff", ".tif"):
        text = await extract_text_from_image(file_path)
        return text, "image"

    # PDF文件
    elif ext == ".pdf":
        text = extract_text_from_pdf(file_path)
        return text, "pdf"

    # Word文档
    elif ext in (".doc", ".docx"):
        text = extract_text_from_docx(file_path)
        return text, "word"

    # Excel文件
    elif ext in (".xls", ".xlsx"):
        text = extract_text_from_excel(file_path)
        return text, "excel"

    # PPT文件
    elif ext in (".ppt", ".pptx"):
        text = extract_text_from_pptx(file_path)
        return text, "ppt"

    # 文本文件
    elif ext == ".txt":
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read(), "text"
        except UnicodeDecodeError:
            with open(file_path, "r", encoding="gbk", errors="ignore") as f:
                return f.read(), "text"

    # Markdown文件
    elif ext in (".md", ".markdown"):
        text = extract_text_from_markdown(file_path)
        return text, "markdown"

    # CSV文件
    elif ext == ".csv":
        text = extract_text_from_csv(file_path)
        return text, "csv"

    # JSON文件
    elif ext == ".json":
        text = extract_text_from_json(file_path)
        return text, "json"

    # HTML文件
    elif ext in (".html", ".htm"):
        text = extract_text_from_html(file_path)
        return text, "html"

    # XML文件
    elif ext == ".xml":
        text = extract_text_from_xml(file_path)
        return text, "xml"

    # RTF文件
    elif ext == ".rtf":
        text = extract_text_from_rtf(file_path)
        return text, "rtf"

    # 配置文件
    elif ext in (".ini", ".cfg", ".conf", ".properties"):
        text = extract_text_from_ini(file_path)
        return text, "config"

    # 代码文件
    elif ext in (".py", ".js", ".java", ".c", ".cpp", ".h", ".hpp", ".cs",
                 ".go", ".rs", ".rb", ".php", ".swift", ".kt", ".scala",
                 ".r", ".m", ".sql", ".sh", ".bash", ".bat", ".ps1"):
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read(), "code"
        except UnicodeDecodeError:
            with open(file_path, "r", encoding="gbk", errors="ignore") as f:
                return f.read(), "code"

    else:
        logger.warning(f"不支持的文件类型: {ext}")
        # 尝试作为文本读取
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read(1000)  # 只读取前1000字符
                if content.strip():
                    return content, "text"
        except Exception:
            pass
        return "", "unknown"


def get_today_dir() -> str:
    """获取今天的日期目录名"""
    return datetime.now().strftime("%Y-%m-%d")
