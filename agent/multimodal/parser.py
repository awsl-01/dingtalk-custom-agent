"""
文件深度解析器

对 PDF/Word/PPT 中的表格、图表、页眉页脚进行结构化抽取
"""
import os
import logging
import json
from typing import Dict, Optional, List, Any
from dataclasses import dataclass, field

logger = logging.getLogger(__name__)


@dataclass
class ParseResult:
    """解析结果"""
    text: str                    # 全文文本
    tables: list = field(default_factory=list)      # 表格数据
    images: list = field(default_factory=list)       # 图片信息
    headers: list = field(default_factory=list)      # 页眉
    footers: list = field(default_factory=list)      # 页脚
    metadata: dict = field(default_factory=dict)     # 元数据
    structure: dict = field(default_factory=dict)    # 文档结构
    error: str = ""                                  # 错误信息


@dataclass
class TableData:
    """表格数据"""
    headers: list        # 表头
    rows: list           # 数据行
    page: int = 0        # 所在页码
    position: dict = field(default_factory=dict)  # 位置信息


class DeepFileParser:
    """
    文件深度解析器

    支持：
    - PDF：文本、表格、图片、页眉页脚
    - Word：文本、表格、样式
    - PPT：文本、幻灯片结构
    - Excel：多工作表、公式
    """

    def __init__(self):
        self._parsers = {
            ".pdf": self._parse_pdf,
            ".doc": self._parse_word,
            ".docx": self._parse_word,
            ".ppt": self._parse_ppt,
            ".pptx": self._parse_ppt,
            ".xls": self._parse_excel,
            ".xlsx": self._parse_excel,
        }

    async def parse(self, file_path: str,
                   extract_tables: bool = True,
                   extract_images: bool = False) -> ParseResult:
        """
        解析文件

        参数:
            file_path: 文件路径
            extract_tables: 是否提取表格
            extract_images: 是否提取图片

        返回:
            解析结果
        """
        if not os.path.exists(file_path):
            return ParseResult(text="", error=f"文件不存在: {file_path}")

        ext = os.path.splitext(file_path)[1].lower()
        parser = self._parsers.get(ext)

        if not parser:
            return ParseResult(text="", error=f"不支持的文件格式: {ext}")

        try:
            result = await parser(file_path, extract_tables, extract_images)
            return result
        except Exception as e:
            logger.error(f"解析文件失败: {e}")
            return ParseResult(text="", error=str(e))

    async def _parse_pdf(self, file_path: str,
                        extract_tables: bool = True,
                        extract_images: bool = False) -> ParseResult:
        """解析 PDF 文件"""
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(file_path)

            texts = []
            tables = []
            images = []
            headers = []
            footers = []

            for page_num in range(len(doc)):
                page = doc[page_num]

                # 提取文本
                text = page.get_text()
                texts.append(text)

                # 提取表格（简化实现）
                if extract_tables:
                    page_tables = self._extract_tables_from_text(text)
                    for table in page_tables:
                        tables.append({
                            "headers": table.get("headers", []),
                            "rows": table.get("rows", []),
                            "page": page_num + 1
                        })

                # 提取图片
                if extract_images:
                    image_list = page.get_images()
                    for img_idx, img in enumerate(image_list):
                        images.append({
                            "page": page_num + 1,
                            "index": img_idx,
                            "xref": img[0]
                        })

                # 提取页眉页脚（简化实现）
                header_text = self._extract_header_footer(text, "header")
                footer_text = self._extract_header_footer(text, "footer")
                if header_text:
                    headers.append({"page": page_num + 1, "text": header_text})
                if footer_text:
                    footers.append({"page": page_num + 1, "text": footer_text})

            # 提取元数据
            metadata = doc.metadata or {}

            doc.close()

            return ParseResult(
                text="\n\n".join(texts),
                tables=tables,
                images=images,
                headers=headers,
                footers=footers,
                metadata={
                    "title": metadata.get("title", ""),
                    "author": metadata.get("author", ""),
                    "subject": metadata.get("subject", ""),
                    "pages": len(texts),
                }
            )

        except ImportError:
            return ParseResult(
                text="",
                error="需要安装 PyMuPDF: pip install PyMuPDF"
            )

    async def _parse_word(self, file_path: str,
                         extract_tables: bool = True,
                         extract_images: bool = False) -> ParseResult:
        """解析 Word 文件"""
        try:
            from docx import Document

            doc = Document(file_path)

            texts = []
            tables = []

            # 提取段落文本
            for para in doc.paragraphs:
                if para.text.strip():
                    texts.append(para.text)

            # 提取表格
            if extract_tables:
                for table_idx, table in enumerate(doc.tables):
                    headers = []
                    rows = []

                    for row_idx, row in enumerate(table.rows):
                        row_data = [cell.text.strip() for cell in row.cells]

                        if row_idx == 0:
                            headers = row_data
                        else:
                            rows.append(row_data)

                    if headers or rows:
                        tables.append({
                            "headers": headers,
                            "rows": rows,
                            "index": table_idx
                        })

            # 提取元数据
            metadata = {}
            if doc.core_properties:
                metadata = {
                    "title": doc.core_properties.title or "",
                    "author": doc.core_properties.author or "",
                    "created": str(doc.core_properties.created) if doc.core_properties.created else "",
                    "modified": str(doc.core_properties.modified) if doc.core_properties.modified else "",
                }

            return ParseResult(
                text="\n".join(texts),
                tables=tables,
                metadata=metadata
            )

        except ImportError:
            return ParseResult(
                text="",
                error="需要安装 python-docx: pip install python-docx"
            )

    async def _parse_ppt(self, file_path: str,
                        extract_tables: bool = True,
                        extract_images: bool = False) -> ParseResult:
        """解析 PPT 文件"""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)

            texts = []
            tables = []
            structure = {"slides": []}

            for slide_idx, slide in enumerate(prs.slides):
                slide_texts = []
                slide_tables = []

                for shape in slide.shapes:
                    # 提取文本
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)

                    # 提取表格
                    if extract_tables and shape.has_table:
                        table = shape.table
                        headers = []
                        rows = []

                        for row_idx, row in enumerate(table.rows):
                            row_data = [cell.text.strip() for cell in row.cells]

                            if row_idx == 0:
                                headers = row_data
                            else:
                                rows.append(row_data)

                        if headers or rows:
                            slide_tables.append({
                                "headers": headers,
                                "rows": rows
                            })

                # 记录幻灯片结构
                structure["slides"].append({
                    "index": slide_idx + 1,
                    "title": slide_texts[0] if slide_texts else "",
                    "content_count": len(slide_texts)
                })

                texts.extend(slide_texts)
                tables.extend(slide_tables)

            # 提取元数据
            metadata = {}
            if prs.core_properties:
                metadata = {
                    "title": prs.core_properties.title or "",
                    "author": prs.core_properties.author or "",
                    "slides": len(prs.slides),
                }

            return ParseResult(
                text="\n\n".join(texts),
                tables=tables,
                metadata=metadata,
                structure=structure
            )

        except ImportError:
            return ParseResult(
                text="",
                error="需要安装 python-pptx: pip install python-pptx"
            )

    async def _parse_excel(self, file_path: str,
                          extract_tables: bool = True,
                          extract_images: bool = False) -> ParseResult:
        """解析 Excel 文件"""
        try:
            import openpyxl

            wb = openpyxl.load_workbook(file_path, data_only=True)

            texts = []
            tables = []
            structure = {"sheets": []}

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]

                # 提取数据
                data = []
                for row in ws.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    data.append(row_data)

                if data:
                    # 第一行作为表头
                    headers = data[0] if data else []
                    rows = data[1:] if len(data) > 1 else []

                    tables.append({
                        "sheet": sheet_name,
                        "headers": headers,
                        "rows": rows
                    })

                    # 生成文本表示
                    sheet_texts = [f"【{sheet_name}】"]
                    for row in data[:10]:  # 只取前10行
                        sheet_texts.append("\t".join(row))
                    texts.append("\n".join(sheet_texts))

                # 记录工作表结构
                structure["sheets"].append({
                    "name": sheet_name,
                    "rows": ws.max_row,
                    "columns": ws.max_column
                })

            wb.close()

            # 提取元数据
            metadata = {
                "sheets": len(wb.sheetnames),
                "file_size": os.path.getsize(file_path)
            }

            return ParseResult(
                text="\n\n".join(texts),
                tables=tables,
                metadata=metadata,
                structure=structure
            )

        except ImportError:
            return ParseResult(
                text="",
                error="需要安装 openpyxl: pip install openpyxl"
            )

    def _extract_tables_from_text(self, text: str) -> List[dict]:
        """从文本中提取表格（简化实现）"""
        tables = []
        lines = text.split("\n")

        # 简单的表格检测：查找连续的分隔线
        current_table = None
        for line in lines:
            line = line.strip()
            if not line:
                if current_table and len(current_table.get("rows", [])) > 0:
                    tables.append(current_table)
                    current_table = None
                continue

            # 检测表格行（包含 | 或制表符）
            if "|" in line or "\t" in line:
                cells = [c.strip() for c in line.split("|") if c.strip()]
                if len(cells) > 1:
                    if current_table is None:
                        current_table = {"headers": cells, "rows": []}
                    else:
                        current_table["rows"].append(cells)

        if current_table and len(current_table.get("rows", [])) > 0:
            tables.append(current_table)

        return tables

    def _extract_header_footer(self, text: str, position: str) -> str:
        """提取页眉页脚（简化实现）"""
        lines = text.split("\n")
        if not lines:
            return ""

        if position == "header":
            # 取前几行作为页眉
            header_lines = lines[:3]
            return "\n".join(header_lines).strip()
        else:
            # 取后几行作为页脚
            footer_lines = lines[-3:]
            return "\n".join(footer_lines).strip()

    async def parse_batch(self, file_paths: List[str],
                         extract_tables: bool = True) -> List[ParseResult]:
        """
        批量解析文件

        参数:
            file_paths: 文件路径列表
            extract_tables: 是否提取表格

        返回:
            解析结果列表
        """
        results = []
        for path in file_paths:
            result = await self.parse(path, extract_tables)
            results.append(result)
        return results

    def get_supported_formats(self) -> List[str]:
        """获取支持的文件格式"""
        return list(self._parsers.keys())


# 全局文件解析器实例
_parser: Optional[DeepFileParser] = None


def get_file_parser() -> DeepFileParser:
    """获取全局文件解析器实例"""
    global _parser
    if _parser is None:
        _parser = DeepFileParser()
    return _parser
