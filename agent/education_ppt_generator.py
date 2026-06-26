"""
教育PPT生成模块
支持：教案生成、课件大纲、说课稿、教学反思
根据学情自动调整教案难度
"""

import os
import json
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from openai import OpenAI
import config


# ==================== 教育内容生成提示词 ====================

LESSON_PLAN_PROMPT = """你是一位资深教育专家，擅长根据教材章节生成专业的教案。

请根据用户提供的教材章节信息，生成一份完整的教案，严格按照以下JSON格式输出：

{
    "title": "教案标题",
    "chapter": "章节名称",
    "subject": "学科",
    "grade": "年级",
    "duration": "课时安排（如：2课时）",
    "difficulty_level": "难度等级（基础/中等/提高/拓展）",
    "lesson_objectives": {
        "knowledge": ["知识与技能目标1", "知识与技能目标2"],
        "process": ["过程与方法目标1", "过程与方法目标2"],
        "emotion": ["情感态度价值观目标1"]
    },
    "key_points": ["教学重点1", "教学重点2"],
    "difficult_points": ["教学难点1", "教学难点2"],
    "teaching_methods": ["教学方法1", "教学方法2"],
    "student_analysis": {
        "level": "学情水平描述",
        "prerequisites": ["前置知识1", "前置知识2"],
        "interests": ["学生兴趣点1", "学生兴趣点2"]
    },
    "teaching_process": {
        "lead_in": {
            "duration": "5分钟",
            "activities": ["导入活动1", "导入活动2"],
            "design_intent": "设计意图说明"
        },
        "new_lesson": {
            "duration": "25分钟",
            "activities": [
                {
                    "title": "环节1标题",
                    "duration": "10分钟",
                    "content": "教学内容说明",
                    "methods": ["教学方法"],
                    "student_activities": ["学生活动1"]
                }
            ]
        },
        "practice": {
            "duration": "10分钟",
            "exercises": ["练习1", "练习2"]
        },
        "summary": {
            "duration": "3分钟",
            "content": "课堂小结内容"
        },
        "homework": {
            "required": ["必做作业1"],
            "optional": ["选做作业1"]
        }
    },
    "board_design": {
        "title": "板书标题",
        "structure": ["板书内容层次1", "板书内容层次2"]
    },
    "teaching_reflection": {
        "expected_highlights": ["预期亮点1", "预期亮点2"],
        "potential_issues": ["可能问题1", "可能问题2"],
        "improvement_suggestions": ["改进方向1"]
    }
}

要求：
1. 教学目标要具体、可衡量、可达成
2. 教学过程要详细，包含时间分配
3. 学情分析要结合实际教学场景
4. 根据难度等级调整内容深度
5. 板书设计要简洁明了"""

SLIDE_OUTLINE_PROMPT = """你是一位专业的课件设计专家，擅长将教案转化为精美的课件大纲。

请根据教案内容，生成课件大纲，严格按照以下JSON格式输出：

{
    "title": "课件标题",
    "subtitle": "副标题（含学科、年级、教师信息）",
    "style": "课件风格（简约/活泼/学术/科技）",
    "slides": [
        {
            "type": "封面",
            "title": "封面标题",
            "subtitle": "副标题",
            "elements": ["标题", "副标题", "教师信息", "日期"]
        },
        {
            "type": "目录",
            "title": "目录",
            "items": ["目录项1", "目录项2", "目录项3"]
        },
        {
            "type": "教学目标",
            "title": "学习目标",
            "objectives": ["目标1", "目标2", "目标3"],
            "icon_suggestion": "目标图标建议"
        },
        {
            "type": "导入",
            "title": "课堂导入",
            "content": "导入内容描述",
            "visual_suggestion": "视觉元素建议（如：图片、视频、动画）"
        },
        {
            "type": "知识讲解",
            "title": "知识点标题",
            "content": ["要点1", "要点2", "要点3"],
            "visual_suggestion": "可视化建议",
            "interaction": "互动设计建议"
        },
        {
            "type": "例题演示",
            "title": "例题标题",
            "problem": "题目内容",
            "solution_steps": ["解题步骤1", "解题步骤2"],
            "tips": "解题技巧提示"
        },
        {
            "type": "课堂练习",
            "title": "课堂练习",
            "exercises": ["练习题1", "练习题2"],
            "hints": ["提示1", "提示2"]
        },
        {
            "type": "知识总结",
            "title": "课堂小结",
            "summary_points": ["总结点1", "总结点2"],
            "mind_map_structure": "思维导图结构建议"
        },
        {
            "type": "作业布置",
            "title": "课后作业",
            "required": ["必做作业"],
            "optional": ["选做作业"],
            "deadline": "截止时间建议"
        },
        {
            "type": "结束页",
            "title": "谢谢观看",
            "elements": ["Q&A提示", "联系方式"]
        }
    ]
}

要求：
1. 每页幻灯片内容简洁，适合课堂展示
2. 注重视觉效果和互动设计
3. 知识点要分层次展示
4. 配色和风格要符合学科特点
5. 总页数控制在12-20页"""

LESSON_SPEECH_PROMPT = """你是一位经验丰富的教研专家，擅长撰写说课稿。

请根据教案内容，生成一份完整的说课稿，严格按照以下JSON格式输出：

{
    "title": "说课稿标题",
    "presenter": "说课人",
    "school": "所在学校",
    "date": "说课日期",
    "sections": [
        {
            "section": "说教材",
            "title": "一、说教材",
            "content": "教材分析内容...",
            "key_points": ["教材地位", "内容分析", "编排意图"]
        },
        {
            "section": "说学情",
            "title": "二、说学情",
            "content": "学情分析内容...",
            "key_points": ["认知水平", "学习特点", "已有基础"]
        },
        {
            "section": "说教学目标",
            "title": "三、说教学目标",
            "content": "教学目标阐述...",
            "key_points": ["知识目标", "能力目标", "情感目标"]
        },
        {
            "section": "说教学重难点",
            "title": "四、说教学重难点",
            "content": "重难点分析...",
            "key_points": ["重点确定依据", "难点突破策略"]
        },
        {
            "section": "说教法与学法",
            "title": "五、说教法与学法",
            "content": "教法学法阐述...",
            "key_points": ["教学方法", "学习方法", "教学手段"]
        },
        {
            "section": "说教学过程",
            "title": "六、说教学过程",
            "content": "教学过程详细阐述...",
            "sub_sections": [
                {
                    "title": "（一）导入新课",
                    "content": "导入环节设计意图..."
                },
                {
                    "title": "（二）探究新知",
                    "content": "新授环节设计意图..."
                },
                {
                    "title": "（三）巩固练习",
                    "content": "练习环节设计意图..."
                },
                {
                    "title": "（四）课堂小结",
                    "content": "小结环节设计意图..."
                },
                {
                    "title": "（五）布置作业",
                    "content": "作业设计意图..."
                }
            ]
        },
        {
            "section": "说板书设计",
            "title": "七、说板书设计",
            "content": "板书设计说明...",
            "key_points": ["板书内容", "设计意图"]
        },
        {
            "section": "说教学反思",
            "title": "八、说教学反思",
            "content": "教学反思内容...",
            "key_points": ["预期效果", "改进方向"]
        }
    ],
    "conclusion": "结束语"
}

要求：
1. 说课稿要体现教学理念和设计思路
2. 语言要专业、流畅，适合口头表达
3. 每个部分要有明确的逻辑关系
4. 突出教学设计的创新点
5. 总时长控制在15-20分钟"""

TEACHING_REFLECTION_PROMPT = """你是一位注重教学反思的优秀教师，擅长从多个维度进行教学反思。

请根据教案和实际教学情况，生成一份教学反思，严格按照以下JSON格式输出：

{
    "title": "教学反思标题",
    "lesson_info": {
        "topic": "课题",
        "date": "教学日期",
        "class": "授课班级",
        "duration": "课时"
    },
    "reflection_dimensions": [
        {
            "dimension": "教学目标达成",
            "title": "一、教学目标达成情况",
            "content": "目标达成分析...",
            "achievements": ["达成的目标1", "达成的目标2"],
            "gaps": ["未完全达成的目标"],
            "rating": "优秀/良好/一般/需改进"
        },
        {
            "dimension": "教学内容处理",
            "title": "二、教学内容处理",
            "content": "内容处理分析...",
            "highlights": ["亮点1", "亮点2"],
            "issues": ["不足1"],
            "suggestions": ["改进建议1"]
        },
        {
            "dimension": "教学方法运用",
            "title": "三、教学方法运用",
            "content": "方法运用分析...",
            "effective_methods": ["有效的方法"],
            "ineffective_methods": ["效果不佳的方法"],
            "alternative_suggestions": ["替代方案"]
        },
        {
            "dimension": "学生参与度",
            "title": "四、学生参与度分析",
            "content": "参与度分析...",
            "positive_signs": ["积极表现"],
            "concerns": ["需要关注的问题"],
            "strategies": ["提升策略"]
        },
        {
            "dimension": "课堂生成",
            "title": "五、课堂生成处理",
            "content": "课堂生成分析...",
            "unexpected_events": ["突发事件"],
            "handling_effectiveness": ["处理效果"],
            "improvement": ["改进方向"]
        },
        {
            "dimension": "教学效果",
            "title": "六、教学效果评估",
            "content": "效果评估...",
            "test_results": "测试/作业情况",
            "student_feedback": "学生反馈",
            "overall_rating": "整体评价"
        }
    ],
    "summary": {
        "main_achievements": ["主要成就1", "主要成就2"],
        "main_issues": ["主要问题1"],
        "action_items": ["后续改进措施1", "后续改进措施2"]
    },
    "next_steps": {
        "short_term": ["短期改进计划"],
        "long_term": ["长期发展规划"]
    }
}

要求：
1. 反思要具体、真实，有数据支撑
2. 分析要有深度，不能流于表面
3. 改进措施要可操作、可衡量
4. 语言要客观、专业
5. 体现教师的专业成长意识"""


# ==================== 学情难度调整 ====================

DIFFICULTY_LEVELS = {
    "基础": {
        "description": "面向学习困难学生，注重基础概念和基本技能",
        "adjustments": {
            "content_depth": "浅显易懂，多用生活实例",
            "exercise_difficulty": "以基础题为主，适当降低难度",
            "scaffolding": "提供更多脚手架和提示",
            "pace": "放慢教学节奏，增加练习时间",
            "support": "加强个别辅导，分层作业设计"
        }
    },
    "中等": {
        "description": "面向中等水平学生，兼顾基础与提高",
        "adjustments": {
            "content_depth": "适中深度，理论与实践结合",
            "exercise_difficulty": "基础题+提高题，比例7:3",
            "scaffolding": "适度引导，培养独立思考",
            "pace": "正常教学节奏",
            "support": "小组合作学习，同伴互助"
        }
    },
    "提高": {
        "description": "面向学有余力学生，注重思维拓展",
        "adjustments": {
            "content_depth": "深入挖掘，拓展延伸",
            "exercise_difficulty": "提高题+拓展题，比例6:4",
            "scaffolding": "减少支架，鼓励自主探究",
            "pace": "适当加快，增加拓展内容",
            "support": "项目式学习，自主探究"
        }
    },
    "拓展": {
        "description": "面向学科特长生，注重创新思维培养",
        "adjustments": {
            "content_depth": "深度探究，跨学科融合",
            "exercise_difficulty": "综合题+创新题",
            "scaffolding": "最小化支架，开放性任务",
            "pace": "灵活安排，专题研究",
            "support": "导师制，竞赛辅导"
        }
    }
}


# ==================== AI内容生成 ====================

def generate_with_ai(system_prompt: str, user_message: str) -> dict:
    """调用AI生成结构化内容"""
    client = OpenAI(
        api_key=config.OPENAI_API_KEY,
        base_url=config.OPENAI_BASE_URL,
    )

    response = client.chat.completions.create(
        model=config.OPENAI_MODEL,
        max_tokens=8192,
        temperature=0.7,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_message},
        ],
    )

    content = response.choices[0].message.content
    return parse_json_response(content)


def parse_json_response(content: str) -> dict:
    """解析AI返回的JSON内容"""
    try:
        return json.loads(content)
    except json.JSONDecodeError:
        if "```json" in content:
            json_str = content.split("```json")[1].split("```")[0].strip()
            return json.loads(json_str)
        elif "```" in content:
            json_str = content.split("```")[1].split("```")[0].strip()
            return json.loads(json_str)
        raise ValueError("无法解析AI返回的内容")


def generate_lesson_plan(chapter_info: str, difficulty: str = "中等", student_info: str = "") -> dict:
    """生成教案"""
    difficulty_desc = DIFFICULTY_LEVELS.get(difficulty, DIFFICULTY_LEVELS["中等"])

    user_message = f"""请根据以下信息生成教案：

教材章节：{chapter_info}
难度等级：{difficulty}（{difficulty_desc['description']}）
学情调整要求：{json.dumps(difficulty_desc['adjustments'], ensure_ascii=False)}
{f'学生情况：{student_info}' if student_info else ''}

请生成完整的教案内容。"""

    return generate_with_ai(LESSON_PLAN_PROMPT, user_message)


def generate_slide_outline(lesson_plan: dict, style: str = "简约") -> dict:
    """根据教案生成课件大纲"""
    user_message = f"""请根据以下教案生成课件大纲：

教案内容：
{json.dumps(lesson_plan, ensure_ascii=False, indent=2)}

课件风格：{style}

请生成详细的课件大纲。"""

    return generate_with_ai(SLIDE_OUTLINE_PROMPT, user_message)


def generate_lesson_speech(lesson_plan: dict) -> dict:
    """根据教案生成说课稿"""
    user_message = f"""请根据以下教案生成说课稿：

教案内容：
{json.dumps(lesson_plan, ensure_ascii=False, indent=2)}

请生成完整的说课稿。"""

    return generate_with_ai(LESSON_SPEECH_PROMPT, user_message)


def generate_teaching_reflection(lesson_plan: dict, actual_situation: str = "") -> dict:
    """根据教案生成教学反思"""
    user_message = f"""请根据以下教案生成教学反思：

教案内容：
{json.dumps(lesson_plan, ensure_ascii=False, indent=2)}
{f'实际教学情况：{actual_situation}' if actual_situation else ''}

请生成详细的教学反思。"""

    return generate_with_ai(TEACHING_REFLECTION_PROMPT, user_message)


# ==================== PPT生成 ====================

# 教育PPT配色方案
EDUCATION_COLORS = {
    "primary": RGBColor(0x15, 0x65, 0xC0),      # 深蓝 - 主色
    "secondary": RGBColor(0x42, 0xA5, 0xF5),     # 亮蓝 - 辅助色
    "accent": RGBColor(0xFF, 0x98, 0x00),        # 橙色 - 强调色
    "success": RGBColor(0x4C, 0xAF, 0x50),       # 绿色 - 成功/正向
    "warning": RGBColor(0xF4, 0x43, 0x36),       # 红色 - 警告
    "text_dark": RGBColor(0x1A, 0x1A, 0x2E),     # 深色文字
    "text_light": RGBColor(0x6B, 0x72, 0x80),    # 浅色文字
    "bg_white": RGBColor(0xFF, 0xFF, 0xFF),      # 白色背景
    "bg_light": RGBColor(0xF5, 0xF7, 0xFA),      # 浅灰背景
    "bg_blue": RGBColor(0xE3, 0xF2, 0xFD),       # 浅蓝背景
}


def add_background(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title_text, page_num=None):
    """添加标题栏"""
    # 顶部装饰条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.08)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = EDUCATION_COLORS["secondary"]
    top_bar.line.fill.background()

    # 标题背景
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.08), Inches(13.333), Inches(1.2)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = EDUCATION_COLORS["primary"]
    title_bg.line.fill.background()

    # 标题文字
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11), Inches(0.9))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.LEFT

    # 页码
    if page_num:
        page_num_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(12.3), Inches(0.35), Inches(0.6), Inches(0.6)
        )
        page_num_shape.fill.solid()
        page_num_shape.fill.fore_color.rgb = EDUCATION_COLORS["accent"]
        page_num_shape.line.fill.background()
        tf = page_num_shape.text_frame
        p = tf.paragraphs[0]
        p.text = str(page_num)
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER


def add_section_header(slide, text, y_position=1.5):
    """添加章节标题"""
    # 左侧装饰线
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(y_position), Inches(0.08), Inches(0.5)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = EDUCATION_COLORS["accent"]
    accent_line.line.fill.background()

    # 章节标题
    header_box = slide.shapes.add_textbox(Inches(1.1), Inches(y_position - 0.1), Inches(10), Inches(0.7))
    tf = header_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = EDUCATION_COLORS["primary"]


def add_content_box(slide, title, items, x=0.8, y=2.3, width=11.5, height=4.5):
    """添加内容框"""
    # 内容背景
    content_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    content_bg.fill.solid()
    content_bg.fill.fore_color.rgb = EDUCATION_COLORS["bg_light"]
    content_bg.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    content_bg.line.width = Pt(1)

    # 内容标题
    if title:
        title_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.2), Inches(width - 0.6), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = EDUCATION_COLORS["primary"]
        content_y = y + 0.8
    else:
        content_y = y + 0.3

    # 内容列表
    content_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(content_y), Inches(width - 0.6), Inches(height - (content_y - y) - 0.3))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = EDUCATION_COLORS["text_dark"]
        p.space_after = Pt(12)


def add_two_column_layout(slide, left_title, left_items, right_title, right_items, y_start=2.3):
    """添加双栏布局"""
    # 左栏
    left_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(y_start), Inches(5.5), Inches(4.5)
    )
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = EDUCATION_COLORS["bg_blue"]
    left_bg.line.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)

    # 左栏标题
    left_title_box = slide.shapes.add_textbox(Inches(1.1), Inches(y_start + 0.2), Inches(5), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = EDUCATION_COLORS["primary"]

    # 左栏内容
    left_content = slide.shapes.add_textbox(Inches(1.1), Inches(y_start + 0.8), Inches(5), Inches(3.5))
    tf = left_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(15)
        p.font.color.rgb = EDUCATION_COLORS["text_dark"]
        p.space_after = Pt(10)

    # 右栏
    right_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(y_start), Inches(5.5), Inches(4.5)
    )
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = RGBColor(0xFFF3E0)
    right_bg.line.color.rgb = RGBColor(0xFFE0B2)

    # 右栏标题
    right_title_box = slide.shapes.add_textbox(Inches(7.1), Inches(y_start + 0.2), Inches(5), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = EDUCATION_COLORS["accent"]

    # 右栏内容
    right_content = slide.shapes.add_textbox(Inches(7.1), Inches(y_start + 0.8), Inches(5), Inches(3.5))
    tf = right_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(15)
        p.font.color.rgb = EDUCATION_COLORS["text_dark"]
        p.space_after = Pt(10)


def add_notes(slide, notes_text):
    """添加演讲者备注"""
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = notes_text


def create_lesson_plan_ppt(lesson_plan: dict, output_path: str) -> str:
    """创建教案PPT"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ====== 封面页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, EDUCATION_COLORS["primary"])

    # 主标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = lesson_plan.get("title", "教案")
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{lesson_plan.get('subject', '')} · {lesson_plan.get('grade', '')} · {lesson_plan.get('chapter', '')}"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    p.alignment = PP_ALIGN.CENTER

    # 信息栏
    info_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(1))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"课时：{lesson_plan.get('duration', '')} | 难度：{lesson_plan.get('difficulty_level', '中等')} | {datetime.now().strftime('%Y年%m月%d日')}"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
    p.alignment = PP_ALIGN.CENTER

    # ====== 教学目标页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, EDUCATION_COLORS["bg_white"])
    add_title_bar(slide, "教学目标", 2)

    objectives = lesson_plan.get("lesson_objectives", {})
    y_pos = 1.5

    for obj_type, obj_title in [("knowledge", "知识与技能"), ("process", "过程与方法"), ("emotion", "情感态度价值观")]:
        if objectives.get(obj_type):
            add_section_header(slide, obj_title, y_pos)
            content_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.6), Inches(11), Inches(1))
            tf = content_box.text_frame
            tf.word_wrap = True
            for i, obj in enumerate(objectives[obj_type]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"•  {obj}"
                p.font.size = Pt(16)
                p.font.color.rgb = EDUCATION_COLORS["text_dark"]
            y_pos += 1.8

    # ====== 教学重难点页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, EDUCATION_COLORS["bg_white"])
    add_title_bar(slide, "教学重难点", 3)

    add_two_column_layout(
        slide,
        "教学重点",
        lesson_plan.get("key_points", []),
        "教学难点",
        lesson_plan.get("difficult_points", [])
    )

    # ====== 学情分析页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, EDUCATION_COLORS["bg_white"])
    add_title_bar(slide, "学情分析", 4)

    student_analysis = lesson_plan.get("student_analysis", {})

    # 学情水平
    add_section_header(slide, "学情水平", 1.5)
    level_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11), Inches(0.8))
    tf = level_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = student_analysis.get("level", "")
    p.font.size = Pt(16)
    p.font.color.rgb = EDUCATION_COLORS["text_dark"]

    # 前置知识
    add_section_header(slide, "前置知识", 3.2)
    prereq_box = slide.shapes.add_textbox(Inches(1.2), Inches(3.9), Inches(11), Inches(1.2))
    tf = prereq_box.text_frame
    tf.word_wrap = True
    for i, prereq in enumerate(student_analysis.get("prerequisites", [])):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"•  {prereq}"
        p.font.size = Pt(15)
        p.font.color.rgb = EDUCATION_COLORS["text_dark"]

    # 学生兴趣
    add_section_header(slide, "学生兴趣点", 5.3)
    interests_box = slide.shapes.add_textbox(Inches(1.2), Inches(6.0), Inches(11), Inches(1))
    tf = interests_box.text_frame
    tf.word_wrap = True
    for i, interest in enumerate(student_analysis.get("interests", [])):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"•  {interest}"
        p.font.size = Pt(15)
        p.font.color.rgb = EDUCATION_COLORS["text_dark"]

    # ====== 教学过程页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, EDUCATION_COLORS["bg_white"])
    add_title_bar(slide, "教学过程", 5)

    teaching_process = lesson_plan.get("teaching_process", {})

    # 教学环节列表
    phases = [
        ("lead_in", "课堂导入"),
        ("new_lesson", "新课讲授"),
        ("practice", "巩固练习"),
        ("summary", "课堂小结"),
        ("homework", "作业布置")
    ]

    y_pos = 1.5
    for phase_key, phase_name in phases:
        phase_data = teaching_process.get(phase_key, {})
        if phase_data:
            # 环节标题和时间
            duration = phase_data.get("duration", "")
            header_text = f"{phase_name}（{duration}）" if duration else phase_name
            add_section_header(slide, header_text, y_pos)

            # 环节内容
            content_items = []
            if phase_key == "lead_in":
                content_items = phase_data.get("activities", [])
            elif phase_key == "new_lesson":
                for activity in phase_data.get("activities", []):
                    content_items.append(f"{activity.get('title', '')}: {activity.get('content', '')[:50]}...")
            elif phase_key == "practice":
                content_items = phase_data.get("exercises", [])
            elif phase_key == "summary":
                content_items = [phase_data.get("content", "")]
            elif phase_key == "homework":
                content_items = [f"必做: {', '.join(phase_data.get('required', []))}"]
                if phase_data.get("optional"):
                    content_items.append(f"选做: {', '.join(phase_data['optional'])}")

            content_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.6), Inches(11), Inches(0.8))
            tf = content_box.text_frame
            tf.word_wrap = True
            for i, item in enumerate(content_items[:2]):  # 只显示前2条
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"•  {item}"
                p.font.size = Pt(14)
                p.font.color.rgb = EDUCATION_COLORS["text_dark"]

            y_pos += 1.1

    # ====== 板书设计页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, EDUCATION_COLORS["bg_white"])
    add_title_bar(slide, "板书设计", 6)

    board_design = lesson_plan.get("board_design", {})

    # 板书标题
    board_title_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(0.8))
    tf = board_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = board_design.get("title", lesson_plan.get("title", ""))
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = EDUCATION_COLORS["primary"]
    p.alignment = PP_ALIGN.CENTER

    # 板书内容
    board_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(2.8), Inches(9), Inches(3.5)
    )
    board_bg.fill.solid()
    board_bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    board_bg.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

    board_content = slide.shapes.add_textbox(Inches(2.5), Inches(3), Inches(8), Inches(3))
    tf = board_content.text_frame
    tf.word_wrap = True

    for i, item in enumerate(board_design.get("structure", [])):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = EDUCATION_COLORS["text_dark"]
        p.space_after = Pt(16)
        p.alignment = PP_ALIGN.CENTER

    # ====== 教学反思页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, EDUCATION_COLORS["bg_white"])
    add_title_bar(slide, "教学反思", 7)

    reflection = lesson_plan.get("teaching_reflection", {})

    # 预期亮点
    add_section_header(slide, "预期亮点", 1.5)
    highlights_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11), Inches(1.2))
    tf = highlights_box.text_frame
    tf.word_wrap = True
    for i, highlight in enumerate(reflection.get("expected_highlights", [])):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"✓  {highlight}"
        p.font.size = Pt(15)
        p.font.color.rgb = EDUCATION_COLORS["success"]

    # 可能问题
    add_section_header(slide, "可能问题", 3.6)
    issues_box = slide.shapes.add_textbox(Inches(1.2), Inches(4.3), Inches(11), Inches(1.2))
    tf = issues_box.text_frame
    tf.word_wrap = True
    for i, issue in enumerate(reflection.get("potential_issues", [])):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"⚠  {issue}"
        p.font.size = Pt(15)
        p.font.color.rgb = EDUCATION_COLORS["warning"]

    # 改进方向
    add_section_header(slide, "改进方向", 5.7)
    improvements_box = slide.shapes.add_textbox(Inches(1.2), Inches(6.4), Inches(11), Inches(0.8))
    tf = improvements_box.text_frame
    tf.word_wrap = True
    for i, suggestion in enumerate(reflection.get("improvement_suggestions", [])):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"→  {suggestion}"
        p.font.size = Pt(15)
        p.font.color.rgb = EDUCATION_COLORS["primary"]

    # ====== 结束页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, EDUCATION_COLORS["primary"])

    thank_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]
    p.text = "教案设计完成"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "祝教学顺利！"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    p.alignment = PP_ALIGN.CENTER

    # 保存文件
    prs.save(output_path)
    return output_path


def create_slide_outline_ppt(slide_outline: dict, output_path: str) -> str:
    """创建课件大纲PPT"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_data = slide_outline.get("slides", [])

    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(slide, EDUCATION_COLORS["bg_white"])

        slide_type = slide_data.get("type", "")
        title = slide_data.get("title", f"第{i+1}页")

        # 根据类型设置不同的样式
        if slide_type == "封面":
            add_background(slide, EDUCATION_COLORS["primary"])

            title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER

            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(1))
            tf = subtitle_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_data.get("subtitle", "")
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
            p.alignment = PP_ALIGN.CENTER

        elif slide_type == "目录":
            add_title_bar(slide, title, i + 1)

            items = slide_data.get("items", [])
            content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(5))
            tf = content_box.text_frame
            tf.word_wrap = True

            for j, item in enumerate(items):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"  {j+1:02d}  {item}"
                p.font.size = Pt(22)
                p.font.color.rgb = EDUCATION_COLORS["text_dark"]
                p.space_after = Pt(20)

        elif slide_type == "教学目标":
            add_title_bar(slide, title, i + 1)

            objectives = slide_data.get("objectives", [])
            content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(5))
            tf = content_box.text_frame
            tf.word_wrap = True

            for j, obj in enumerate(objectives):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"🎯  {obj}"
                p.font.size = Pt(20)
                p.font.color.rgb = EDUCATION_COLORS["text_dark"]
                p.space_after = Pt(16)

        elif slide_type == "结束页":
            add_background(slide, EDUCATION_COLORS["primary"])

            title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER

        else:
            # 默认布局
            add_title_bar(slide, title, i + 1)

            content = slide_data.get("content", [])
            if isinstance(content, list):
                add_content_box(slide, None, content)
            elif isinstance(content, str):
                add_content_box(slide, None, [content])

    prs.save(output_path)
    return output_path


def create_lesson_speech_ppt(speech_data: dict, output_path: str) -> str:
    """创建说课稿PPT"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 封面页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, EDUCATION_COLORS["primary"])

    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = speech_data.get("title", "说课稿")
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    info_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(2))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"说课人：{speech_data.get('presenter', '')}"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = f"{speech_data.get('school', '')}  {speech_data.get('date', '')}"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
    p.alignment = PP_ALIGN.CENTER

    # 各个说课环节
    sections = speech_data.get("sections", [])
    for i, section in enumerate(sections):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(slide, EDUCATION_COLORS["bg_white"])
        add_title_bar(slide, section.get("title", ""), i + 2)

        # 内容
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = section.get("content", "")
        p.font.size = Pt(18)
        p.font.color.rgb = EDUCATION_COLORS["text_dark"]
        p.line_spacing = Pt(28)

        # 关键点
        key_points = section.get("key_points", [])
        if key_points:
            points_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5))
            tf = points_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "关键点：" + " | ".join(key_points)
            p.font.size = Pt(14)
            p.font.color.rgb = EDUCATION_COLORS["primary"]
            p.font.bold = True

    # 结束页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, EDUCATION_COLORS["primary"])

    thank_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]
    p.text = "谢谢各位专家指导"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    conclusion_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1))
    tf = conclusion_box.text_frame
    p = tf.paragraphs[0]
    p.text = speech_data.get("conclusion", "")
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    p.alignment = PP_ALIGN.CENTER

    prs.save(output_path)
    return output_path


def create_teaching_reflection_ppt(reflection_data: dict, output_path: str) -> str:
    """创建教学反思PPT"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 封面页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, EDUCATION_COLORS["primary"])

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = reflection_data.get("title", "教学反思")
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    lesson_info = reflection_data.get("lesson_info", {})
    info_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(1.5))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"课题：{lesson_info.get('topic', '')}"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = f"班级：{lesson_info.get('class', '')}  日期：{lesson_info.get('date', '')}"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
    p.alignment = PP_ALIGN.CENTER

    # 反思维度页
    dimensions = reflection_data.get("reflection_dimensions", [])
    for i, dim in enumerate(dimensions):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(slide, EDUCATION_COLORS["bg_white"])
        add_title_bar(slide, dim.get("title", ""), i + 2)

        # 内容
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(2.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = dim.get("content", "")
        p.font.size = Pt(16)
        p.font.color.rgb = EDUCATION_COLORS["text_dark"]
        p.line_spacing = Pt(24)

        # 评级
        rating = dim.get("rating", "")
        if rating:
            rating_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.8), Inches(4.2), Inches(2), Inches(0.6)
            )
            rating_box.fill.solid()
            if "优秀" in rating:
                rating_box.fill.fore_color.rgb = EDUCATION_COLORS["success"]
            elif "良好" in rating:
                rating_box.fill.fore_color.rgb = EDUCATION_COLORS["secondary"]
            else:
                rating_box.fill.fore_color.rgb = EDUCATION_COLORS["accent"]
            rating_box.line.fill.background()

            tf = rating_box.text_frame
            p = tf.paragraphs[0]
            p.text = rating
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER

        # 要点列表
        y_pos = 5.0
        for key, label in [("achievements", "成就"), ("highlights", "亮点"), ("effective_methods", "有效方法")]:
            items = dim.get(key, [])
            if items:
                label_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(11.5), Inches(0.4))
                tf = label_box.text_frame
                p = tf.paragraphs[0]
                p.text = f"{label}："
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = EDUCATION_COLORS["primary"]

                items_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.3), Inches(11), Inches(0.8))
                tf = items_box.text_frame
                tf.word_wrap = True
                for j, item in enumerate(items[:3]):
                    if j == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"•  {item}"
                    p.font.size = Pt(13)
                    p.font.color.rgb = EDUCATION_COLORS["text_dark"]

                y_pos += 1.2

    # 总结页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, EDUCATION_COLORS["bg_white"])
    add_title_bar(slide, "反思总结与改进计划", len(dimensions) + 2)

    summary = reflection_data.get("summary", {})

    # 主要成就
    add_section_header(slide, "主要成就", 1.5)
    achievements_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11), Inches(1.2))
    tf = achievements_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(summary.get("main_achievements", [])):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"✓  {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = EDUCATION_COLORS["success"]

    # 主要问题
    add_section_header(slide, "主要问题", 3.6)
    issues_box = slide.shapes.add_textbox(Inches(1.2), Inches(4.3), Inches(11), Inches(1))
    tf = issues_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(summary.get("main_issues", [])):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"⚠  {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = EDUCATION_COLORS["warning"]

    # 改进措施
    add_section_header(slide, "改进措施", 5.5)
    actions_box = slide.shapes.add_textbox(Inches(1.2), Inches(6.2), Inches(11), Inches(1))
    tf = actions_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(summary.get("action_items", [])):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"→  {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = EDUCATION_COLORS["primary"]

    # 结束页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, EDUCATION_COLORS["primary"])

    thank_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]
    p.text = "教学反思完成"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    next_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1))
    tf = next_box.text_frame
    p = tf.paragraphs[0]
    p.text = "持续改进，追求卓越"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    p.alignment = PP_ALIGN.CENTER

    prs.save(output_path)
    return output_path


# ==================== 主入口函数 ====================

def generate_education_ppt(user_message: str, output_dir: str = "projects") -> tuple:
    """
    教育PPT生成主入口

    支持的请求类型：
    1. 教案生成：包含"教案"关键词
    2. 课件大纲：包含"课件"或"大纲"关键词
    3. 说课稿：包含"说课"关键词
    4. 教学反思：包含"反思"关键词
    5. 完整生成：包含"教学设计"或"全套"关键词
    """
    os.makedirs(output_dir, exist_ok=True)

    # 解析用户需求
    user_message_lower = user_message.lower()

    # 提取难度信息
    difficulty = "中等"
    for level in ["基础", "中等", "提高", "拓展"]:
        if level in user_message:
            difficulty = level
            break

    # 提取学情信息
    student_info = ""
    if "学情" in user_message or "学生" in user_message:
        # 尝试提取学情描述
        for keyword in ["学情：", "学情:", "学生情况：", "学生情况:"]:
            if keyword in user_message:
                student_info = user_message.split(keyword)[1].split("。")[0].split("；")[0]
                break

    # 提取章节信息
    chapter_info = user_message

    # 判断生成类型
    if "说课" in user_message_lower:
        # 生成说课稿
        lesson_plan = generate_lesson_plan(chapter_info, difficulty, student_info)
        speech_data = generate_lesson_speech(lesson_plan)

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"说课稿_{timestamp}.pptx"
        output_path = os.path.join(output_dir, filename)

        create_lesson_speech_ppt(speech_data, output_path)
        return output_path, speech_data.get("title", "说课稿")

    elif "反思" in user_message_lower:
        # 生成教学反思
        lesson_plan = generate_lesson_plan(chapter_info, difficulty, student_info)
        reflection_data = generate_teaching_reflection(lesson_plan)

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"教学反思_{timestamp}.pptx"
        output_path = os.path.join(output_dir, filename)

        create_teaching_reflection_ppt(reflection_data, output_path)
        return output_path, reflection_data.get("title", "教学反思")

    elif "课件" in user_message_lower or "大纲" in user_message_lower:
        # 生成课件大纲
        lesson_plan = generate_lesson_plan(chapter_info, difficulty, student_info)
        slide_outline = generate_slide_outline(lesson_plan)

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"课件大纲_{timestamp}.pptx"
        output_path = os.path.join(output_dir, filename)

        create_slide_outline_ppt(slide_outline, output_path)
        return output_path, slide_outline.get("title", "课件大纲")

    elif "教案" in user_message_lower or "教学设计" in user_message_lower or "全套" in user_message_lower:
        # 生成完整教案
        lesson_plan = generate_lesson_plan(chapter_info, difficulty, student_info)

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"教案_{timestamp}.pptx"
        output_path = os.path.join(output_dir, filename)

        create_lesson_plan_ppt(lesson_plan, output_path)
        return output_path, lesson_plan.get("title", "教案")

    else:
        # 默认生成教案
        lesson_plan = generate_lesson_plan(chapter_info, difficulty, student_info)

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"教案_{timestamp}.pptx"
        output_path = os.path.join(output_dir, filename)

        create_lesson_plan_ppt(lesson_plan, output_path)
        return output_path, lesson_plan.get("title", "教案")


if __name__ == "__main__":
    # 测试用例
    test_message = "请为高中数学《导数的概念》生成一份教案，难度中等，学生是理科班学生，基础较好"
    path, title = generate_education_ppt(test_message)
    print(f"生成完成: {path}")
    print(f"标题: {title}")
