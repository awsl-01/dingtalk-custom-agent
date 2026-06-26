"""
基于模板的PPT生成器 - 直接复制模板幻灯片并修改内容
这种方式可以完全保留模板的格式和样式
"""
import os
import json
import copy
from typing import Dict, List, Tuple, Optional, Any
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from openai import OpenAI
import config


class TemplateBasedPPTGenerator:
    """基于模板的PPT生成器 - 直接复制模板幻灯片"""

    def __init__(self, template_path: str):
        self.template_path = template_path
        self.template_prs = Presentation(template_path)

    def _copy_slide(self, source_slide, target_prs):
        """复制幻灯片（保留所有格式）"""
        # 获取源幻灯片的布局
        source_layout = source_slide.slide_layout

        # 在目标演示文稿中查找相同名称的布局
        target_layout = None
        for layout in target_prs.slide_layouts:
            if layout.name == source_layout.name:
                target_layout = layout
                break

        # 如果找不到相同布局，使用空白布局
        if target_layout is None:
            target_layout = target_prs.slide_layouts[6]

        # 添加新幻灯片
        new_slide = target_prs.slides.add_slide(target_layout)

        # 复制背景
        self._copy_background(source_slide, new_slide)

        # 复制所有形状
        for shape in source_slide.shapes:
            self._copy_shape(shape, new_slide)

        return new_slide

    def _copy_background(self, source_slide, target_slide):
        """复制背景"""
        try:
            source_bg = source_slide.background
            target_bg = target_slide.background

            # 复制填充
            source_fill = source_bg.fill
            target_fill = target_bg.fill

            if source_fill.type is not None:
                if source_fill.type == 1:  # SOLID_FILL
                    target_fill.solid()
                    target_fill.fore_color.rgb = source_fill.fore_color.rgb
        except Exception as e:
            print(f"复制背景失败: {e}")

    def _copy_shape(self, source_shape, target_slide):
        """复制形状"""
        try:
            # 获取形状的位置和大小
            left = source_shape.left
            top = source_shape.top
            width = source_shape.width
            height = source_shape.height

            # 根据形状类型复制
            if source_shape.shape_type == MSO_SHAPE.RECTANGLE:
                # 矩形
                new_shape = target_slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left, top, width, height
                )
                self._copy_fill(source_shape, new_shape)
                self._copy_line(source_shape, new_shape)

                # 如果有文本框，也要复制文本
                if source_shape.has_text_frame:
                    self._copy_text_frame(source_shape.text_frame, new_shape.text_frame)

            elif source_shape.shape_type == MSO_SHAPE.OVAL:
                # 椭圆
                new_shape = target_slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    left, top, width, height
                )
                self._copy_fill(source_shape, new_shape)
                self._copy_line(source_shape, new_shape)

                # 如果有文本框，也要复制文本
                if source_shape.has_text_frame:
                    self._copy_text_frame(source_shape.text_frame, new_shape.text_frame)

            elif source_shape.has_text_frame:
                # 文本框
                new_shape = target_slide.shapes.add_textbox(
                    left, top, width, height
                )
                self._copy_text_frame(source_shape.text_frame, new_shape.text_frame)

        except Exception as e:
            print(f"复制形状失败: {e}")

    def _copy_fill(self, source_shape, target_shape):
        """复制填充"""
        try:
            source_fill = source_shape.fill
            target_fill = target_shape.fill

            if source_fill.type is not None:
                if source_fill.type == 1:  # SOLID_FILL
                    target_fill.solid()
                    target_fill.fore_color.rgb = source_fill.fore_color.rgb
        except Exception as e:
            print(f"复制填充失败: {e}")

    def _copy_line(self, source_shape, target_shape):
        """复制边框"""
        try:
            source_line = source_shape.line
            target_line = target_shape.line

            if source_line.fill.type is not None:
                if source_line.fill.type == 1:  # SOLID_FILL
                    target_line.fill.solid()
                    target_line.fill.fore_color.rgb = source_line.fill.fore_color.rgb
            else:
                target_line.fill.background()
        except Exception as e:
            print(f"复制边框失败: {e}")

    def _copy_text_frame(self, source_tf, target_tf):
        """复制文本框"""
        try:
            # 复制文本框属性
            target_tf.word_wrap = source_tf.word_wrap

            # 清空目标文本框
            for i in range(len(target_tf.paragraphs) - 1, 0, -1):
                target_tf._element.remove(target_tf.paragraphs[i]._element)

            # 复制段落
            for i, source_para in enumerate(source_tf.paragraphs):
                if i == 0:
                    target_para = target_tf.paragraphs[0]
                else:
                    target_para = target_tf.add_paragraph()

                # 复制段落属性
                target_para.alignment = source_para.alignment

                # 复制runs
                for source_run in source_para.runs:
                    target_run = target_para.add_run()
                    target_run.text = source_run.text

                    # 复制字体属性
                    if source_run.font.name:
                        target_run.font.name = source_run.font.name
                    if source_run.font.size:
                        target_run.font.size = source_run.font.size
                    if source_run.font.bold is not None:
                        target_run.font.bold = source_run.font.bold
                    if source_run.font.italic is not None:
                        target_run.font.italic = source_run.font.italic
                    try:
                        if source_run.font.color and source_run.font.color.type is not None:
                            target_run.font.color.rgb = source_run.font.color.rgb
                    except:
                        pass

        except Exception as e:
            print(f"复制文本框失败: {e}")

    def _modify_text(self, text_frame, new_text):
        """修改文本框内容（保留格式）"""
        try:
            if not text_frame.paragraphs:
                return

            # 获取第一个段落的格式
            first_para = text_frame.paragraphs[0]
            if not first_para.runs:
                return

            # 保留第一个run的格式
            first_run = first_para.runs[0]
            font_name = first_run.font.name
            font_size = first_run.font.size
            font_bold = first_run.font.bold
            font_color = None
            try:
                if first_run.font.color and first_run.font.color.type is not None:
                    font_color = first_run.font.color.rgb
            except:
                pass

            # 清空所有runs
            for para in text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""

            # 设置新文本（使用原有格式）
            first_run.text = new_text
            if font_name:
                first_run.font.name = font_name
            if font_size:
                first_run.font.size = font_size
            if font_bold is not None:
                first_run.font.bold = font_bold
            if font_color:
                first_run.font.color.rgb = font_color

        except Exception as e:
            print(f"修改文本失败: {e}")

    def _add_content_to_slide(self, slide, title, content_items):
        """向幻灯片添加内容（修改现有文本框）"""
        try:
            # 获取所有形状并按位置排序
            all_shapes = list(slide.shapes)
            all_shapes.sort(key=lambda s: (s.top / 914400, s.left / 914400))

            # 找到主标题形状（页面最上方的大标题）
            title_shape = None
            remaining_shapes = []

            for shape in all_shapes:
                if not shape.has_text_frame:
                    remaining_shapes.append(shape)
                    continue

                top_inches = shape.top / 914400
                height_inches = shape.height / 914400
                text = ""
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        text = para.text.strip()
                        break

                # 主标题：在页面顶部（y < 1.5），高度较小（< 1.5），有文本
                if top_inches < 1.5 and height_inches < 1.5 and text:
                    if title_shape is None:
                        title_shape = shape
                    else:
                        remaining_shapes.append(shape)
                else:
                    remaining_shapes.append(shape)

            # 修改主标题
            if title_shape and title:
                self._modify_text(title_shape.text_frame, title)

            # 检测是否有多卡片布局（查找带编号的卡片结构）
            cards = self._detect_content_cards(remaining_shapes)

            if cards and content_items:
                # 多卡片布局：分配内容到各个卡片
                self._fill_content_cards(cards, content_items)
            elif content_items:
                # 单内容区域：查找最大的文本框
                text_shapes = [s for s in remaining_shapes if s.has_text_frame]
                if text_shapes:
                    # 按面积排序，取最大的
                    text_shapes.sort(key=lambda s: (s.width * s.height), reverse=True)
                    self._modify_content_text(text_shapes[0].text_frame, content_items)

        except Exception as e:
            print(f"添加内容失败: {e}")
            import traceback
            traceback.print_exc()

    def _detect_content_cards(self, shapes):
        """检测多卡片布局，返回卡片列表。使用卡片背景区域分组，避免形状被多个卡片抢夺。"""
        cards = []

        # 查找所有形状（包括无文本的背景形状），按位置排序
        all_sorted = sorted(shapes, key=lambda s: (s.top / 914400, s.left / 914400))

        # 查找大矩形背景形状（卡片背景通常面积较大且没有文本）
        bg_shapes = []
        for shape in all_sorted:
            if shape.has_text_frame:
                has_text = any(p.text.strip() for p in shape.text_frame.paragraphs)
                if has_text:
                    continue
            w = shape.width / 914400
            h = shape.height / 914400
            # 卡片背景通常是较大的矩形（宽>2, 高>2）
            if w > 2.0 and h > 2.0:
                bg_shapes.append(shape)

        if len(bg_shapes) < 2:
            return cards  # 不是多卡片布局

        # 为每个背景区域找到属于它的文本形状
        for bg in bg_shapes:
            bg_left = bg.left / 914400
            bg_top = bg.top / 914400
            bg_right = bg_left + bg.width / 914400
            bg_bottom = bg_top + bg.height / 914400

            # 容差：0.3英寸
            margin = 0.3

            number = None
            number_shape = None
            title_shape = None
            desc_shape = None

            for shape in all_sorted:
                if not shape.has_text_frame:
                    continue
                text = ""
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        text = para.text.strip()
                        break
                if not text:
                    continue

                s_left = shape.left / 914400
                s_top = shape.top / 914400
                s_right = s_left + shape.width / 914400
                s_bottom = s_top + shape.height / 914400

                # 检查形状中心是否在卡片背景区域内
                s_cx = (s_left + s_right) / 2
                s_cy = (s_top + s_bottom) / 2
                if not (bg_left - margin <= s_cx <= bg_right + margin and
                        bg_top - margin <= s_cy <= bg_bottom + margin):
                    continue

                # 分类：编号、标题、描述
                if text.isdigit() and len(text) <= 2 and number is None:
                    number = text
                    number_shape = shape
                elif len(text) < 30 and title_shape is None:
                    title_shape = shape
                elif desc_shape is None:
                    desc_shape = shape

            if number:
                cards.append({
                    'number': number,
                    'num_shape': number_shape,
                    'title_shape': title_shape,
                    'desc_shape': desc_shape,
                    'bg_shape': bg,
                })

        # 按编号排序
        cards.sort(key=lambda c: c['number'])
        return cards

    def _hide_card(self, card):
        """隐藏整个卡片（将所有形状移到幻灯片外）"""
        try:
            # 隐藏背景
            if card.get('bg_shape'):
                self._move_shape_offscreen(card['bg_shape'])
            # 隐藏编号
            if card.get('num_shape'):
                self._move_shape_offscreen(card['num_shape'])
            # 隐藏标题
            if card.get('title_shape'):
                self._move_shape_offscreen(card['title_shape'])
            # 隐藏描述
            if card.get('desc_shape'):
                self._move_shape_offscreen(card['desc_shape'])
        except Exception as e:
            print(f"隐藏卡片失败: {e}")

    def _move_shape_offscreen(self, shape):
        """将形状移到幻灯片外（隐藏）"""
        try:
            # 将形状移到左侧很远的位置（幻灯片宽度的10倍）
            shape.left = -shape.width * 10
            shape.top = 0
        except Exception as e:
            print(f"移动形状失败: {e}")

    def _fill_content_cards(self, cards, content_items):
        """将内容分配到各个卡片"""
        try:
            card_count = len(cards)
            if card_count == 0:
                return

            for i, card in enumerate(cards):
                if i < len(content_items):
                    content = content_items[i]

                    # 修改卡片标题（取内容的前15个字符作为标题）
                    if card['title_shape']:
                        title_text = content
                        if len(title_text) > 15:
                            title_text = title_text[:15] + "..."
                        self._modify_text(card['title_shape'].text_frame, title_text)

                    # 修改卡片描述（使用完整内容）
                    if card['desc_shape']:
                        self._modify_text(card['desc_shape'].text_frame, content)
                else:
                    # 没有对应内容的卡片，隐藏整个卡片
                    self._hide_card(card)

        except Exception as e:
            print(f"填充内容卡片失败: {e}")

    def _modify_content_text(self, text_frame, content_items):
        """修改内容文本（保留项目符号格式）"""
        try:
            if not text_frame.paragraphs:
                return

            # 获取第一个段落的格式
            first_para = text_frame.paragraphs[0]
            if not first_para.runs:
                return

            # 保留格式
            first_run = first_para.runs[0]
            font_name = first_run.font.name
            font_size = first_run.font.size
            font_color = None
            try:
                if first_run.font.color and first_run.font.color.type is not None:
                    font_color = first_run.font.color.rgb
            except:
                pass

            # 检测项目符号
            bullet_char = "•"
            if first_run.text and first_run.text[0] in "•·-–—*▪▸►":
                bullet_char = first_run.text[0]

            # 清空所有段落
            for i in range(len(text_frame.paragraphs) - 1, 0, -1):
                text_frame._element.remove(text_frame.paragraphs[i]._element)

            # 添加新内容
            for i, item in enumerate(content_items):
                if i == 0:
                    para = text_frame.paragraphs[0]
                    # 清空原有runs
                    for run in para.runs:
                        run.text = ""
                else:
                    para = text_frame.add_paragraph()

                run = para.add_run()
                run.text = f"{bullet_char}  {item}"

                # 应用格式
                if font_name:
                    run.font.name = font_name
                if font_size:
                    run.font.size = font_size
                if font_color:
                    run.font.color.rgb = font_color

        except Exception as e:
            print(f"修改内容文本失败: {e}")

    def generate_content(self, user_message: str) -> Dict[str, Any]:
        """调用AI生成PPT内容"""
        system_prompt = """你是一个专业的PPT内容生成助手。根据用户的需求，生成结构化的PPT内容。

请严格按照以下JSON格式输出，不要输出其他内容：
{
    "title": "PPT标题",
    "subtitle": "副标题（可选）",
    "slides": [
        {
            "title": "幻灯片标题",
            "content": ["要点1", "要点2", "要点3"],
            "notes": "演讲者备注（可选）"
        }
    ]
}

要求：
1. 内容简洁精炼，适合PPT展示
2. 每页幻灯片3-5个要点为宜
3. 使用专业的表达方式
4. 总页数控制在5-15页"""

        client = OpenAI(
            api_key=config.OPENAI_API_KEY,
            base_url=config.OPENAI_BASE_URL,
        )

        response = client.chat.completions.create(
            model=config.OPENAI_MODEL,
            max_tokens=4096,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_message},
            ],
        )

        content = response.choices[0].message.content

        try:
            return json.loads(content)
        except json.JSONDecodeError:
            if "```json" in content:
                json_str = content.split("```json")[1].split("```")[0].strip()
                return json.loads(json_str)
            elif "```" in content:
                json_str = content.split("```")[1].split("```")[0].strip()
                return json.loads(json_str)
            raise ValueError("无法解析AI返回的内容")

    def generate(self, user_message: str, output_dir: str = "projects") -> Tuple[str, str]:
        """生成PPT（基于模板）"""
        # 生成内容
        ppt_data = self.generate_content(user_message)

        # 创建新的演示文稿
        new_prs = Presentation()

        # 设置幻灯片尺寸（与模板相同）
        new_prs.slide_width = self.template_prs.slide_width
        new_prs.slide_height = self.template_prs.slide_height

        # 获取模板幻灯片
        template_slides = list(self.template_prs.slides)

        # 复制封面页并修改标题
        if len(template_slides) >= 1:
            cover_slide = self._copy_slide(template_slides[0], new_prs)
            self._modify_cover_content(cover_slide, ppt_data)

        # 为每个内容创建页面（基于目录页模板）
        if len(template_slides) >= 2:
            toc_template = template_slides[1]

            for i, slide_data in enumerate(ppt_data.get("slides", [])):
                # 复制目录页模板
                content_slide = self._copy_slide(toc_template, new_prs)

                # 修改内容
                title = slide_data.get("title", f"第{i+2}页")
                content = slide_data.get("content", [])
                self._add_content_to_slide(content_slide, title, content)

        # 复制封面页作为结束页
        if len(template_slides) >= 1:
            ending_slide = self._copy_slide(template_slides[0], new_prs)
            self._modify_ending_content(ending_slide)

        # 保存文件
        os.makedirs(output_dir, exist_ok=True)
        title = ppt_data.get("title", "presentation")
        safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).strip()
        safe_title = safe_title[:50]
        output_path = os.path.join(output_dir, f"{safe_title}.pptx")
        new_prs.save(output_path)

        return output_path, title

    def _modify_cover_content(self, slide, ppt_data):
        """修改封面页内容"""
        try:
            # 查找标题和副标题文本框
            for shape in slide.shapes:
                if shape.has_text_frame:
                    top_inches = shape.top / 914400

                    # 标题（通常在 y=2 英寸左右）
                    if 1.5 < top_inches < 3.0:
                        self._modify_text(shape.text_frame, ppt_data.get("title", "演示文稿"))

                    # 副标题（通常在 y=3.5 英寸左右）
                    elif 3.0 < top_inches < 4.5:
                        if ppt_data.get("subtitle"):
                            self._modify_text(shape.text_frame, ppt_data["subtitle"])
                        else:
                            # 如果没有副标题，清空
                            self._modify_text(shape.text_frame, "")

        except Exception as e:
            print(f"修改封面内容失败: {e}")

    def _modify_ending_content(self, slide):
        """修改结束页内容"""
        try:
            # 查找并修改所有文本框
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # 将标题改为"谢谢观看"
                    if shape.text_frame.paragraphs:
                        text = shape.text_frame.paragraphs[0].text
                        if text and len(text) > 5:  # 原来的标题
                            self._modify_text(shape.text_frame, "谢谢观看")

        except Exception as e:
            print(f"修改结束页内容失败: {e}")


def generate_template_based_ppt(user_message: str, template_path: str, output_dir: str = "projects") -> Tuple[str, str]:
    """使用模板生成PPT（直接复制模板）"""
    generator = TemplateBasedPPTGenerator(template_path)
    return generator.generate(user_message, output_dir)
