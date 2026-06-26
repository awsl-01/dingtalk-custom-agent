import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from openai import OpenAI
import config

SYSTEM_PROMPT = """你是一个专业的PPT内容生成助手。根据用户的需求，生成结构化的PPT内容。

请严格按照以下JSON格式输出，不要输出其他内容：
{
    "title": "PPT标题",
    "subtitle": "副标题（可选）",
    "slides": [
        {
            "title": "幻灯片标题",
            "content": ["要点1", "要点2", "要点3"],
            "notes": "演讲者备注（可选）"
        }
    ]
}

要求：
1. 内容简洁精炼，适合PPT展示
2. 每页幻灯片3-5个要点为宜
3. 使用专业的表达方式
4. 总页数控制在5-15页"""


def generate_ppt_content(user_message: str) -> dict:
    """调用AI生成PPT内容结构"""
    client = OpenAI(
        api_key=config.OPENAI_API_KEY,
        base_url=config.OPENAI_BASE_URL,
    )

    response = client.chat.completions.create(
        model=config.OPENAI_MODEL,
        max_tokens=4096,
        messages=[
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": user_message},
        ],
    )

    content = response.choices[0].message.content

    # 提取JSON内容
    try:
        # 尝试直接解析
        return json.loads(content)
    except json.JSONDecodeError:
        # 尝试从markdown代码块中提取
        if "```json" in content:
            json_str = content.split("```json")[1].split("```")[0].strip()
            return json.loads(json_str)
        elif "```" in content:
            json_str = content.split("```")[1].split("```")[0].strip()
            return json.loads(json_str)
        raise ValueError("无法解析AI返回的内容")


def create_pptx(ppt_data: dict, output_path: str) -> str:
    """根据结构化数据生成PPTX文件"""
    prs = Presentation()

    # 设置16:9宽屏比例
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 颜色主题
    PRIMARY_COLOR = RGBColor(0x2B, 0x57, 0x9A)  # 深蓝色
    SECONDARY_COLOR = RGBColor(0x42, 0x85, 0xF4)  # 亮蓝色
    ACCENT_COLOR = RGBColor(0xFF, 0x6B, 0x35)  # 橙色强调色
    TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)  # 深灰色文字
    LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)  # 浅灰背景

    # ====== 封面页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR

    # 装饰线条
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(3.2), Inches(2), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ppt_data.get("title", "演示文稿")
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.LEFT

    # 副标题
    if ppt_data.get("subtitle"):
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(1))
        tf = subtitle_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ppt_data["subtitle"]
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
        p.alignment = PP_ALIGN.LEFT

    # ====== 内容页 ======
    for i, slide_data in enumerate(ppt_data.get("slides", [])):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

        # 浅色背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = LIGHT_BG if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)

        # 顶部装饰条
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(13.333), Inches(0.08)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = SECONDARY_COLOR
        top_bar.line.fill.background()

        # 页码指示器
        page_num = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(12.3), Inches(0.3), Inches(0.6), Inches(0.6)
        )
        page_num.fill.solid()
        page_num.fill.fore_color.rgb = PRIMARY_COLOR
        page_num.line.fill.background()
        tf = page_num.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(i + 2)
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", f"第{i+2}页")
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.alignment = PP_ALIGN.LEFT

        # 左侧装饰线
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(1.4), Inches(0.08), Inches(0.5)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = ACCENT_COLOR
        accent_line.line.fill.background()

        # 内容区域
        content_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.5), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for j, point in enumerate(slide_data.get("content", [])):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # 添加项目符号
            p.text = f"•  {point}"
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_COLOR
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(16)

            # 段落缩进
            p.level = 0

        # 添加演讲者备注
        if slide_data.get("notes"):
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = slide_data["notes"]

    # ====== 结束页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR

    # 感谢文字
    thank_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = thank_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "谢谢观看"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # Q&A文字
    qa_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1))
    tf = qa_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Questions & Answers"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    p.alignment = PP_ALIGN.CENTER

    # 保存文件
    prs.save(output_path)
    return output_path


def generate_ppt(user_message: str, output_dir: str = "projects") -> str:
    """完整的PPT生成流程"""
    # 确保输出目录存在
    os.makedirs(output_dir, exist_ok=True)

    # 生成PPT内容
    ppt_data = generate_ppt_content(user_message)

    # 生成文件名
    title = ppt_data.get("title", "presentation")
    safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).strip()
    safe_title = safe_title[:50]  # 限制长度
    output_path = os.path.join(output_dir, f"{safe_title}.pptx")

    # 生成PPTX
    create_pptx(ppt_data, output_path)

    return output_path, ppt_data.get("title", "演示文稿")
