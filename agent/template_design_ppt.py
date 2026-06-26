"""
基于模板设计基因的PPT生成器
从模板中提取设计基因，生成风格一致的完整PPT
"""
import os
import json
import copy
from typing import Dict, List, Tuple, Optional, Any
from dataclasses import dataclass, field
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from openai import OpenAI
import config

# 版式类型定义
class SlideLayoutType:
    COVER = "cover"           # 封面页
    TOC = "toc"               # 目录页
    CHAPTER = "chapter"       # 章节页
    CONTENT = "content"       # 正文页
    IMAGE_TEXT = "image_text" # 图文页
    QUOTE = "quote"           # 引用/金句页
    ENDING = "ending"         # 结束页


@dataclass
class DesignGene:
    """设计基因数据类"""
    # 背景相关
    background_color: Optional[RGBColor] = None  # 正文页背景色（从目录页提取）
    background_image: Optional[str] = None
    background_type: str = "solid"  # solid, image, gradient

    # 封面/章节页背景色（从封面页提取）
    accent_color: RGBColor = RGBColor(0x1A, 0x56, 0xDB)

    # 标题样式
    title_font_name: str = "微软雅黑"
    title_font_size: Pt = Pt(32)
    title_font_color: RGBColor = RGBColor(0x33, 0x33, 0x33)
    title_font_bold: bool = True
    title_font_italic: bool = False
    title_alignment: PP_ALIGN = PP_ALIGN.LEFT

    # 副标题样式
    subtitle_font_name: str = "微软雅黑"
    subtitle_font_size: Pt = Pt(20)
    subtitle_font_color: RGBColor = RGBColor(0xCC, 0xDD, 0xEE)
    subtitle_font_bold: bool = False

    # 正文样式
    body_font_name: str = "微软雅黑"
    body_font_size: Pt = Pt(18)
    body_font_color: RGBColor = RGBColor(0x33, 0x33, 0x33)
    body_font_bold: bool = False
    body_line_spacing: float = 1.5

    # 项目符号样式
    bullet_style: str = "•"
    bullet_indent: float = 0.3

    # 装饰元素
    has_accent_line: bool = True
    accent_line_color: RGBColor = RGBColor(0xFF, 0x6B, 0x35)  # 装饰线颜色
    accent_line_height: float = 0.05
    has_accent_block: bool = False

    # 页脚/页码
    has_footer: bool = False
    has_page_number: bool = True
    page_number_position: str = "bottom_right"  # top_right, bottom_right, bottom_center

    # 布局参数
    slide_width: float = 13.333
    slide_height: float = 7.5
    margin_left: float = 0.8
    margin_right: float = 0.8
    margin_top: float = 0.5
    margin_bottom: float = 0.5

    # 标题区域
    title_top: float = 0.5
    title_height: float = 1.0
    title_width: float = 11.0

    # 正文区域
    content_top: float = 1.8
    content_width: float = 10.5
    content_height: float = 5.0

    # 装饰线位置
    accent_line_top: float = 1.4
    accent_line_left: float = 0.8


@dataclass
class SlideContent:
    """幻灯片内容数据"""
    title: str
    content: List[str] = field(default_factory=list)
    notes: str = ""
    image_path: Optional[str] = None
    layout_type: str = SlideLayoutType.CONTENT


class TemplateDesignExtractor:
    """模板设计基因提取器"""

    def __init__(self, template_path: str):
        self.template_path = template_path
        self.prs = Presentation(template_path)
        self.gene = DesignGene()

    def extract(self) -> DesignGene:
        """提取设计基因"""
        # 更新幻灯片尺寸
        self.gene.slide_width = self.prs.slide_width / 914400  # EMU to inches
        self.gene.slide_height = self.prs.slide_height / 914400

        # 分析每张幻灯片，分别处理封面页和目录页
        slides = list(self.prs.slides)
        if len(slides) >= 1:
            self._extract_cover_design(slides[0])
        if len(slides) >= 2:
            self._extract_toc_design(slides[1])

        return self.gene

    def _extract_cover_design(self, slide):
        """从封面页提取设计基因"""
        # 提取封面背景颜色（用于章节页和结束页）
        background = slide.background
        fill = background.fill
        if fill.type is not None:
            try:
                if fill.type == 1:  # SOLID_FILL
                    self.gene.accent_color = fill.fore_color.rgb
            except:
                pass

        # 提取封面标题样式
        for shape in slide.shapes:
            if shape.has_text_frame:
                self._extract_cover_text_style(shape)

    def _extract_toc_design(self, slide):
        """从目录页提取设计基因"""
        # 提取目录页背景颜色（用于正文页）
        background = slide.background
        fill = background.fill
        if fill.type is not None:
            try:
                if fill.type == 1:  # SOLID_FILL
                    self.gene.background_color = fill.fore_color.rgb
            except:
                pass

        # 提取装饰元素和文本样式
        for shape in slide.shapes:
            if shape.has_text_frame:
                self._extract_toc_text_style(shape)

            # 提取装饰元素
            if shape.shape_type == MSO_SHAPE.RECTANGLE:
                self._extract_decorative_element(shape)

    def _extract_cover_text_style(self, shape):
        """从封面页提取文本样式"""
        tf = shape.text_frame
        if not tf.paragraphs:
            return

        for para in tf.paragraphs:
            if not para.text.strip():
                continue

            # 封面页的大文本是标题
            if para.runs and para.runs[0].font.size:
                font_size = para.runs[0].font.size
                if font_size >= Pt(36):  # 封面标题通常较大
                    self._extract_title_style(para, is_cover=True)
                elif font_size >= Pt(20):  # 副标题
                    self._extract_subtitle_style(para)

    def _extract_toc_text_style(self, shape):
        """从目录页提取文本样式"""
        tf = shape.text_frame
        if not tf.paragraphs:
            return

        for para in tf.paragraphs:
            if not para.text.strip():
                continue

            # 检查是否是标题（通常是"目录"二字）
            if para.runs and para.runs[0].font.size:
                font_size = para.runs[0].font.size
                if font_size >= Pt(28):
                    # 目录页标题样式，但不覆盖封面标题
                    pass
                else:
                    # 正文样式（项目符号列表）
                    self._extract_body_style(para)

    def _extract_title_style(self, paragraph, is_cover=False):
        """提取标题样式"""
        if not paragraph.runs:
            return

        run = paragraph.runs[0]
        if run.font.name:
            self.gene.title_font_name = run.font.name
        if run.font.size:
            self.gene.title_font_size = run.font.size
        try:
            if run.font.color and run.font.color.type is not None:
                if is_cover:
                    # 封面标题通常是白色
                    pass
                else:
                    self.gene.title_font_color = run.font.color.rgb
        except (AttributeError, TypeError):
            pass
        if run.font.bold is not None:
            self.gene.title_font_bold = run.font.bold
        if run.font.italic is not None:
            self.gene.title_font_italic = run.font.italic

        if paragraph.alignment:
            self.gene.title_alignment = paragraph.alignment

    def _extract_subtitle_style(self, paragraph):
        """提取副标题样式"""
        if not paragraph.runs:
            return

        run = paragraph.runs[0]
        if run.font.name:
            self.gene.subtitle_font_name = run.font.name
        if run.font.size:
            self.gene.subtitle_font_size = run.font.size
        try:
            if run.font.color and run.font.color.type is not None:
                self.gene.subtitle_font_color = run.font.color.rgb
        except (AttributeError, TypeError):
            pass
        if run.font.bold is not None:
            self.gene.subtitle_font_bold = run.font.bold

    def _extract_body_style(self, paragraph):
        """提取正文样式"""
        if not paragraph.runs:
            return

        run = paragraph.runs[0]
        if run.font.name:
            self.gene.body_font_name = run.font.name
        if run.font.size:
            self.gene.body_font_size = run.font.size
        try:
            if run.font.color and run.font.color.type is not None:
                self.gene.body_font_color = run.font.color.rgb
        except (AttributeError, TypeError):
            pass
        if run.font.bold is not None:
            self.gene.body_font_bold = run.font.bold

        # 检测项目符号
        text = paragraph.text.strip()
        if text and text[0] in "•·-–—*▪▸►":
            self.gene.bullet_style = text[0]

        # 检测缩进
        if paragraph.level > 0:
            self.gene.bullet_indent = 0.3 * paragraph.level

    def _extract_decorative_element(self, shape):
        """提取装饰元素"""
        # 检测细线（宽度或高度很小的矩形）
        width_inches = shape.width / 914400
        height_inches = shape.height / 914400

        if width_inches > 1.0 and height_inches < 0.1:
            self.gene.has_accent_line = True
            try:
                if shape.fill.type and shape.fill.fore_color.rgb:
                    self.gene.accent_line_color = shape.fill.fore_color.rgb
            except:
                pass
            self.gene.accent_line_top = shape.top / 914400
            self.gene.accent_line_left = shape.left / 914400

        # 检测色块
        if width_inches > 0.5 and height_inches > 0.5:
            self.gene.has_accent_block = True


class SlideBuilder:
    """幻灯片构建器"""

    def __init__(self, gene: DesignGene):
        self.gene = gene

    def build_slide(self, prs: Presentation, content: SlideContent, page_num: int) -> None:
        """构建单张幻灯片"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

        # 设置背景
        self._set_background(slide, content.layout_type)

        # 根据版式类型构建
        if content.layout_type == SlideLayoutType.CHAPTER:
            self._build_chapter_slide(slide, content, page_num)
        elif content.layout_type == SlideLayoutType.IMAGE_TEXT:
            self._build_image_text_slide(slide, content, page_num)
        elif content.layout_type == SlideLayoutType.QUOTE:
            self._build_quote_slide(slide, content, page_num)
        else:
            self._build_content_slide(slide, content, page_num)

        # 添加装饰元素
        self._add_decorative_elements(slide, content.layout_type)

        # 添加页码
        if self.gene.has_page_number and content.layout_type not in [
            SlideLayoutType.COVER, SlideLayoutType.ENDING
        ]:
            self._add_page_number(slide, page_num)

        # 添加备注
        if content.notes:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = content.notes

    def _set_background(self, slide, layout_type: str):
        """设置幻灯片背景"""
        background = slide.background
        fill = background.fill

        # 章节页和结束页使用强调色
        if layout_type in [SlideLayoutType.CHAPTER, SlideLayoutType.ENDING]:
            fill.solid()
            fill.fore_color.rgb = self.gene.accent_color
        elif self.gene.background_color:
            fill.solid()
            fill.fore_color.rgb = self.gene.background_color
        else:
            fill.solid()
            fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    def _build_content_slide(self, slide, content: SlideContent, page_num: int):
        """构建正文页"""
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(self.gene.margin_left),
            Inches(self.gene.title_top),
            Inches(self.gene.title_width),
            Inches(self.gene.title_height)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = content.title
        run.font.size = self.gene.title_font_size
        run.font.bold = self.gene.title_font_bold
        run.font.color.rgb = self.gene.title_font_color
        run.font.name = self.gene.title_font_name
        p.alignment = self.gene.title_alignment

        # 正文内容
        if content.content:
            content_box = slide.shapes.add_textbox(
                Inches(self.gene.margin_left + 0.4),
                Inches(self.gene.content_top),
                Inches(self.gene.content_width),
                Inches(self.gene.content_height)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE

            for i, point in enumerate(content.content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                run = p.add_run()
                run.text = f"{self.gene.bullet_style}  {point}"
                run.font.size = self.gene.body_font_size
                run.font.color.rgb = self.gene.body_font_color
                run.font.name = self.gene.body_font_name
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(12)

    def _build_chapter_slide(self, slide, content: SlideContent, page_num: int):
        """构建章节页（全图背景+半透明遮罩+大标题）"""
        # 添加半透明遮罩
        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(self.gene.slide_width), Inches(self.gene.slide_height)
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)

        # 设置透明度（通过修改XML）
        try:
            from pptx.oxml.ns import qn
            fill_elem = overlay.fill._fill
            # 获取XML元素
            xml_elem = fill_elem._element if hasattr(fill_elem, '_element') else fill_elem
            solidFill = xml_elem.find(qn('a:solidFill'))
            if solidFill is not None:
                srgb = solidFill.find(qn('a:srgbClr'))
                if srgb is not None:
                    alpha = srgb.makeelement(qn('a:alpha'), {'val': '50000'})
                    srgb.append(alpha)
        except Exception as e:
            # 如果设置透明度失败，使用半透明灰色作为备选
            overlay.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)

        # 大标题居中
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(self.gene.slide_width - 2), Inches(2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = content.title
        run.font.size = Pt(48)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.name = self.gene.title_font_name
        p.alignment = PP_ALIGN.CENTER

        # 副标题
        if content.content:
            subtitle_box = slide.shapes.add_textbox(
                Inches(2), Inches(4.5),
                Inches(self.gene.slide_width - 4), Inches(1)
            )
            tf = subtitle_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = content.content[0]
            run.font.size = Pt(24)
            run.font.color.rgb = self.gene.subtitle_font_color
            run.font.name = self.gene.body_font_name
            p.alignment = PP_ALIGN.CENTER

    def _build_image_text_slide(self, slide, content: SlideContent, page_num: int):
        """构建图文页（左文右图）"""
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(self.gene.margin_left),
            Inches(self.gene.title_top),
            Inches(self.gene.title_width),
            Inches(self.gene.title_height)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = content.title
        run.font.size = self.gene.title_font_size
        run.font.bold = self.gene.title_font_bold
        run.font.color.rgb = self.gene.title_font_color
        run.font.name = self.gene.title_font_name
        p.alignment = self.gene.title_alignment

        # 左侧内容区域
        content_box = slide.shapes.add_textbox(
            Inches(self.gene.margin_left),
            Inches(self.gene.content_top),
            Inches(5.5),
            Inches(self.gene.content_height)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, point in enumerate(content.content[:3]):  # 最多3个要点
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            run.text = f"{self.gene.bullet_style}  {point}"
            run.font.size = self.gene.body_font_size
            run.font.color.rgb = self.gene.body_font_color
            run.font.name = self.gene.body_font_name
            p.space_after = Pt(12)

        # 右侧图片占位符
        img_placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(7), Inches(self.gene.content_top),
            Inches(5.5), Inches(4)
        )
        img_placeholder.fill.solid()
        img_placeholder.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        img_placeholder.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

        # 占位符文字
        tf = img_placeholder.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "[图片]"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(80)

    def _build_quote_slide(self, slide, content: SlideContent, page_num: int):
        """构建引用/金句页"""
        # 居中大号文本
        quote_box = slide.shapes.add_textbox(
            Inches(2), Inches(2),
            Inches(self.gene.slide_width - 4), Inches(3)
        )
        tf = quote_box.text_frame
        tf.word_wrap = True

        # 引号装饰
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "\""
        run.font.size = Pt(72)
        run.font.color.rgb = self.gene.accent_color
        run.font.name = self.gene.title_font_name
        p.alignment = PP_ALIGN.LEFT

        # 引用文本
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = content.content[0] if content.content else content.title
        run.font.size = Pt(28)
        run.font.italic = True
        run.font.color.rgb = self.gene.title_font_color
        run.font.name = self.gene.title_font_name
        p.alignment = PP_ALIGN.CENTER

        # 结尾引号
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "\""
        run.font.size = Pt(72)
        run.font.color.rgb = self.gene.accent_color
        run.font.name = self.gene.title_font_name
        p.alignment = PP_ALIGN.RIGHT

        # 出处
        if len(content.content) > 1:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = f"—— {content.content[1]}"
            run.font.size = Pt(18)
            run.font.color.rgb = self.gene.body_font_color
            run.font.name = self.gene.body_font_name
            p.alignment = PP_ALIGN.RIGHT

    def _add_decorative_elements(self, slide, layout_type: str):
        """添加装饰元素"""
        if layout_type in [SlideLayoutType.COVER, SlideLayoutType.ENDING, SlideLayoutType.CHAPTER]:
            return

        # 添加装饰线
        if self.gene.has_accent_line:
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(self.gene.accent_line_left),
                Inches(self.gene.accent_line_top),
                Inches(2), Inches(self.gene.accent_line_height)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = self.gene.accent_line_color
            line.line.fill.background()

    def _add_page_number(self, slide, page_num: int):
        """添加页码"""
        if self.gene.page_number_position == "top_right":
            left = Inches(self.gene.slide_width - 1.2)
            top = Inches(0.3)
        elif self.gene.page_number_position == "bottom_center":
            left = Inches(self.gene.slide_width / 2 - 0.3)
            top = Inches(self.gene.slide_height - 0.7)
        else:  # bottom_right
            left = Inches(self.gene.slide_width - 1.2)
            top = Inches(self.gene.slide_height - 0.7)

        page_num_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left, top, Inches(0.5), Inches(0.5)
        )
        page_num_shape.fill.solid()
        page_num_shape.fill.fore_color.rgb = self.gene.accent_color
        page_num_shape.line.fill.background()

        tf = page_num_shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(page_num)
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER


class TextOverflowChecker:
    """文字溢出检查器"""

    def __init__(self, gene: DesignGene):
        self.gene = gene

    def check_overflow(self, content_box, text: str, font_size: Pt) -> Tuple[bool, int]:
        """
        检查文本是否会溢出

        Returns:
            (是否溢出, 预计行数)
        """
        # 获取文本框尺寸
        width_inches = content_box.width / 914400
        height_inches = content_box.height / 914400

        # 计算每行可容纳的字符数（粗略估算）
        # 假设每个字符宽度约为字号的0.6倍
        char_width = font_size.pt * 0.6 / 72  # 转换为英寸
        chars_per_line = int(width_inches / char_width)

        if chars_per_line <= 0:
            return True, 0

        # 计算行数
        lines = len(text) / chars_per_line
        if lines != int(lines):
            lines = int(lines) + 1
        else:
            lines = int(lines)

        # 计算可容纳的行数
        line_height = font_size.pt * self.gene.body_line_spacing / 72  # 转换为英寸
        max_lines = int(height_inches / line_height)

        return lines > max_lines, lines

    def split_content(self, content: List[str], max_lines_per_slide: int) -> List[List[str]]:
        """将内容拆分到多张幻灯片"""
        slides_content = []
        current_slide = []
        current_lines = 0

        for item in content:
            # 估算每项内容的行数
            item_lines = max(1, len(item) // 30 + 1)  # 粗略估算

            if current_lines + item_lines > max_lines_per_slide and current_slide:
                slides_content.append(current_slide)
                current_slide = [item]
                current_lines = item_lines
            else:
                current_slide.append(item)
                current_lines += item_lines

        if current_slide:
            slides_content.append(current_slide)

        return slides_content


class LayoutDiversityChecker:
    """排版多样性检查器"""

    def __init__(self):
        self.layout_history: List[str] = []

    def check_and_fix_layout(self, current_layout: str, content: SlideContent) -> str:
        """检查并修正版式，确保无三页连用同一种结构"""
        # 如果历史少于2页，直接使用当前版式
        if len(self.layout_history) < 2:
            self.layout_history.append(current_layout)
            return current_layout

        # 检查是否连续3页相同
        if (self.layout_history[-1] == self.layout_history[-2] == current_layout):
            # 强制改变版式
            if current_layout == SlideLayoutType.CONTENT:
                # 如果内容有图片路径，使用图文版式
                if content.image_path:
                    new_layout = SlideLayoutType.IMAGE_TEXT
                # 如果内容较少，使用引用版式
                elif len(content.content) <= 2:
                    new_layout = SlideLayoutType.QUOTE
                else:
                    new_layout = SlideLayoutType.IMAGE_TEXT
            elif current_layout == SlideLayoutType.IMAGE_TEXT:
                new_layout = SlideLayoutType.CONTENT
            else:
                new_layout = SlideLayoutType.CONTENT

            self.layout_history.append(new_layout)
            return new_layout

        self.layout_history.append(current_layout)
        return current_layout


class TemplateDesignPPTGenerator:
    """基于模板设计基因的PPT生成器"""

    def __init__(self, template_path: str):
        self.template_path = template_path
        self.extractor = TemplateDesignExtractor(template_path)
        self.gene: Optional[DesignGene] = None
        self.overflow_checker: Optional[TextOverflowChecker] = None
        self.layout_checker = LayoutDiversityChecker()
        self.slide_builder: Optional[SlideBuilder] = None

    def _init_design_gene(self):
        """初始化设计基因"""
        if self.gene is None:
            self.gene = self.extractor.extract()
            self.overflow_checker = TextOverflowChecker(self.gene)
            self.slide_builder = SlideBuilder(self.gene)

    def generate_content(self, user_message: str) -> Dict[str, Any]:
        """调用AI生成PPT内容"""
        system_prompt = """你是一个专业的PPT内容生成助手。根据用户的需求，生成结构化的PPT内容。

请严格按照以下JSON格式输出，不要输出其他内容：
{
    "title": "PPT标题",
    "subtitle": "副标题（可选）",
    "slides": [
        {
            "title": "幻灯片标题",
            "content": ["要点1", "要点2", "要点3"],
            "layout_type": "content/chapter/image_text/quote",
            "notes": "演讲者备注（可选）",
            "image_needed": true/false
        }
    ]
}

版式类型说明：
- content: 标准正文页，标题+要点列表
- chapter: 章节页，大标题居中
- image_text: 图文页，左文右图
- quote: 引用/金句页

要求：
1. 内容简洁精炼，适合PPT展示
2. 每页幻灯片3-5个要点为宜
3. 使用专业的表达方式
4. 总页数控制在5-15页
5. 连续三页不能使用相同的layout_type
6. 第一页和最后一页不需要指定layout_type"""

        client = OpenAI(
            api_key=config.OPENAI_API_KEY,
            base_url=config.OPENAI_BASE_URL,
        )

        response = client.chat.completions.create(
            model=config.OPENAI_MODEL,
            max_tokens=4096,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_message},
            ],
        )

        content = response.choices[0].message.content

        try:
            return json.loads(content)
        except json.JSONDecodeError:
            if "```json" in content:
                json_str = content.split("```json")[1].split("```")[0].strip()
                return json.loads(json_str)
            elif "```" in content:
                json_str = content.split("```")[1].split("```")[0].strip()
                return json.loads(json_str)
            raise ValueError("无法解析AI返回的内容")

    def _create_cover_slide(self, prs: Presentation, data: Dict[str, Any]):
        """创建封面页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 背景（使用从封面页提取的强调色）
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.gene.accent_color

        # 装饰线（使用从模板提取的装饰线颜色）
        if self.gene.has_accent_line:
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(self.gene.accent_line_left),
                Inches(3.2),
                Inches(2), Inches(self.gene.accent_line_height)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = self.gene.accent_line_color
            line.line.fill.background()

        # 标题（白色，因为封面页是深色背景）
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = data.get("title", "演示文稿")
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.name = self.gene.title_font_name
        p.alignment = PP_ALIGN.LEFT

        # 副标题（浅色）
        if data.get("subtitle"):
            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(1))
            tf = subtitle_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = data["subtitle"]
            run.font.size = Pt(24)
            run.font.color.rgb = self.gene.subtitle_font_color
            run.font.name = self.gene.subtitle_font_name
            p.alignment = PP_ALIGN.LEFT

    def _create_ending_slide(self, prs: Presentation):
        """创建结束页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 背景（使用从封面页提取的强调色）
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.gene.accent_color

        # 感谢文字
        thank_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(self.gene.slide_width - 2), Inches(2)
        )
        tf = thank_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "谢谢观看"
        run.font.size = Pt(48)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.name = self.gene.title_font_name
        p.alignment = PP_ALIGN.CENTER

        # Q&A文字
        qa_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.5),
            Inches(self.gene.slide_width - 2), Inches(1)
        )
        tf = qa_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Questions & Answers"
        run.font.size = Pt(28)
        run.font.color.rgb = self.gene.subtitle_font_color
        run.font.name = self.gene.body_font_name
        p.alignment = PP_ALIGN.CENTER

    def _map_layout_type(self, layout_str: str) -> str:
        """映射布局类型"""
        layout_map = {
            "chapter": SlideLayoutType.CHAPTER,
            "content": SlideLayoutType.CONTENT,
            "image_text": SlideLayoutType.IMAGE_TEXT,
            "quote": SlideLayoutType.QUOTE,
        }
        return layout_map.get(layout_str, SlideLayoutType.CONTENT)

    def _check_and_fix_overflow(self, slide_data: Dict, content: SlideContent) -> SlideContent:
        """检查并修复文字溢出"""
        # 计算每张幻灯片可容纳的行数
        line_height = self.gene.body_font_size.pt * self.gene.body_line_spacing / 72
        max_lines = int(self.gene.content_height / line_height)

        # 检查内容是否过多
        total_lines = 0
        for item in content.content:
            total_lines += max(1, len(item) // 30 + 1)

        if total_lines > max_lines:
            # 拆分内容
            split_contents = self.overflow_checker.split_content(content.content, max_lines)

            if len(split_contents) > 1:
                # 返回第一页，后续页需要单独处理
                content.content = split_contents[0]
                # 添加备注说明有续页
                content.notes = f"（续1/{len(split_contents)}）\n" + content.notes

        return content

    def generate(self, user_message: str, output_dir: str = "projects") -> Tuple[str, str]:
        """生成完整PPT"""
        # 初始化设计基因
        self._init_design_gene()

        # 生成内容
        ppt_data = self.generate_content(user_message)

        # 创建PPT
        prs = Presentation()

        # 设置幻灯片尺寸
        prs.slide_width = Inches(self.gene.slide_width)
        prs.slide_height = Inches(self.gene.slide_height)

        # 创建封面页
        self._create_cover_slide(prs, ppt_data)

        # 创建内容页
        page_num = 2
        slides = ppt_data.get("slides", [])

        for i, slide_data in enumerate(slides):
            # 确定布局类型
            layout_type = self._map_layout_type(slide_data.get("layout_type", "content"))

            # 创建内容对象
            content = SlideContent(
                title=slide_data.get("title", f"第{page_num}页"),
                content=slide_data.get("content", []),
                notes=slide_data.get("notes", ""),
                image_path=slide_data.get("image_path"),
                layout_type=layout_type
            )

            # 检查排版多样性
            content.layout_type = self.layout_checker.check_and_fix_layout(
                content.layout_type, content
            )

            # 检查文字溢出
            content = self._check_and_fix_overflow(slide_data, content)

            # 构建幻灯片
            self.slide_builder.build_slide(prs, content, page_num)

            # 如果内容被拆分，创建续页
            if slide_data.get("needs_continuation"):
                # 这里可以添加续页逻辑
                pass

            page_num += 1

        # 创建结束页
        self._create_ending_slide(prs)

        # 最终自查
        self._final_check(prs)

        # 保存文件
        os.makedirs(output_dir, exist_ok=True)
        title = ppt_data.get("title", "presentation")
        safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).strip()
        safe_title = safe_title[:50]
        output_path = os.path.join(output_dir, f"{safe_title}.pptx")
        prs.save(output_path)

        return output_path, title

    def _final_check(self, prs: Presentation):
        """最终自查清单"""
        print("\n=== PPT最终自查清单 ===")

        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            layout_type = "未知"
            title_size = 0
            body_size = 0
            has_overflow = False

            # 检查版式类型
            if i == 0:
                layout_type = "封面"
            elif i == len(prs.slides) - 1:
                layout_type = "结束页"
            else:
                # 根据背景和布局推断
                bg_color = None
                try:
                    bg_color = slide.background.fill.fore_color.rgb
                except:
                    pass

                if bg_color == self.gene.accent_color:
                    layout_type = "章节页"
                else:
                    layout_type = "正文页"

            # 检查字体大小
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.runs:
                            try:
                                font_size = para.runs[0].font.size
                                if font_size is not None:
                                    if font_size >= Pt(28):
                                        title_size = font_size.pt
                                    elif font_size <= Pt(24):
                                        body_size = font_size.pt
                            except:
                                pass

            # 检查是否有文本框溢出（简单检查）
            for shape in slide.shapes:
                if shape.has_text_frame:
                    tf = shape.text_frame
                    # 检查是否有文本但没有自动调整
                    if tf.paragraphs and not tf.auto_size:
                        # 粗略检查是否可能溢出
                        text_height = sum(
                            len(p.text) / 30 * self.gene.body_font_size.pt / 72
                            for p in tf.paragraphs
                        )
                        box_height = shape.height / 914400
                        if text_height > box_height:
                            has_overflow = True

            print(f"第{slide_num}页：版式类型【{layout_type}】，"
                  f"标题字号{title_size}，正文字号{body_size}，"
                  f"是否溢出：{'是' if has_overflow else '否'}")


def generate_template_ppt(user_message: str, template_path: str, output_dir: str = "projects") -> Tuple[str, str]:
    """使用模板设计基因生成PPT"""
    generator = TemplateDesignPPTGenerator(template_path)
    return generator.generate(user_message, output_dir)
