"""
增强版PPT生成器
集成模板风格、设计建议、内容润色
"""

import os
import sys
import json
import subprocess
import logging
from pathlib import Path
from typing import Dict, List, Tuple, Optional

# 设置日志
logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")
logger = logging.getLogger(__name__)

# 添加ppt-master脚本路径
PPT_MASTER_DIR = Path(__file__).parent.parent / "ppt-master"
SKILL_DIR = PPT_MASTER_DIR / "skills" / "ppt-master"
SCRIPTS_DIR = SKILL_DIR / "scripts"
PROJECTS_DIR = PPT_MASTER_DIR / "projects"

sys.path.insert(0, str(SCRIPTS_DIR))


def search_content_for_topic(topic: str, subject: str = "") -> Dict:
    """
    搜索主题相关内容用于润色

    参数:
        topic: 主题
        subject: 学科

    返回:
        搜索结果字典
    """
    try:
        # 导入搜索模块
        from agent.web_search import search_web
        import asyncio

        # 构建搜索查询
        queries = [
            f"{topic} 教学设计 教案 2026",
            f"{topic} {subject} 课件 知识点",
            f"{topic} 教学资源 素材"
        ]

        all_results = {}
        for query in queries[:2]:
            try:
                loop = asyncio.get_event_loop()
                if loop.is_running():
                    import concurrent.futures
                    with concurrent.futures.ThreadPoolExecutor() as executor:
                        future = executor.submit(asyncio.run, search_web(query, 3))
                        results = future.result(timeout=10)
                else:
                    results = loop.run_until_complete(search_web(query, 3))
                all_results[query] = results
            except Exception as e:
                logger.warning(f"搜索失败 [{query}]: {e}")
                all_results[query] = []

        return all_results
    except Exception as e:
        logger.warning(f"搜索功能不可用: {e}")
        return {}


def generate_content_with_style(
    topic: str,
    subject: str,
    grade: str,
    content_type: str,
    template_style: str,
    search_results: Dict
) -> List[Dict]:
    """
    根据模板风格生成优化的PPT内容

    参数:
        topic: 主题
        subject: 学科
        grade: 年级
        content_type: 内容类型
        template_style: 模板风格
        search_results: 搜索结果

    返回:
        优化后的内容列表
    """
    from openai import OpenAI
    import config

    client = OpenAI(
        api_key=config.OPENAI_API_KEY,
        base_url=config.OPENAI_BASE_URL,
    )

    # 格式化搜索结果
    search_context = ""
    if search_results:
        search_context = "以下是网络搜索到的相关教学资料：\n\n"
        for query, results in search_results.items():
            if results:
                search_context += f"【{query}】\n"
                for r in results[:2]:
                    search_context += f"- {r.get('title', '')}: {r.get('snippet', '')}\n"
                search_context += "\n"

    # 根据模板风格调整提示词
    style_prompts = {
        "ink_painting": """国风水墨风格要求：
- 使用古典诗词、成语典故
- 配合毛笔书法元素
- 采用竖排或古典排版
- 每页配水墨插图说明
- 色调：墨色、赭石、朱砂红""",

        "macaron_cartoon": """马卡龙卡通风格要求：
- 使用活泼可爱的语言
- 配合卡通插图和表情
- 采用圆角卡片式布局
- 每页有互动游戏元素
- 色调：粉色、薰衣草紫、柠檬黄""",

        "hand_drawn": """手绘插画风格要求：
- 使用亲切自然的表达
- 配合手绘涂鸦元素
- 采用便签纸式布局
- 每页有思维导图元素
- 色调：天蓝、草绿、橙色""",

        "fresh_forest": """清新森系风格要求：
- 使用清新自然的语言
- 配合植物、自然元素
- 采用卡片式布局
- 每页有知识树元素
- 色调：森林绿、浅绿、阳光橙""",

        "vintage_newspaper": """复古报刊风格要求：
- 使用正式严谨的语言
- 配合报纸排版元素
- 采用分栏式布局
- 每页有时间线元素
- 色调：深棕、巧克力色、深红""",

        "guochao_illustration": """国潮插画风格要求：
- 使用现代国潮语言
- 配合国潮插画元素
- 采用现代中式布局
- 每页有传统纹样装饰
- 色调：中国红、明黄、深蓝""",

        "minimalist_ins": """极简ins风格要求：
- 使用简洁专业的语言
- 配合几何图形元素
- 采用大留白布局
- 每页有重点突出
- 色调：深灰、浅灰、红色点缀""",

        "watercolor": """水彩晕染风格要求：
- 使用艺术浪漫的语言
- 配合水彩插画元素
- 采用柔和渐变布局
- 每页有艺术装饰
- 色调：天蓝、粉色、橙色""",

        "tech_neon": """科技霓虹风格要求：
- 使用科技感语言
- 配合霓虹灯光元素
- 采用深色背景布局
- 每页有数据可视化
- 色调：深黑、深灰、霓虹绿""",

        "warm_healing": """暖黄治愈风格要求：
- 使用温暖治愈的语言
- 配合阳光笑脸元素
- 采用温暖卡片布局
- 每页有正能量语录
- 色调：暖黄、橙色、绿色""",

        "european_fresh": """清新欧美风格要求：
- 使用国际化语言
- 配合字母对话框元素
- 采用图文并茂布局
- 每页有情景对话
- 色调：海蓝、紫红、橙色""",

        "lab_tech": """实验室科技风格要求：
- 使用科学严谨的语言
- 配合实验器材元素
- 采用数据图表布局
- 每页有实验步骤
- 色调：青绿、蓝色、红色""",

        "sports_energy": """活力运动风格要求：
- 使用活力四射的语言
- 配合运动器材元素
- 采用动感布局
- 每页有运动挑战
- 色调：橙红、绿色、蓝色""",

        "music_rhythm": """音乐律动风格要求：
- 使用艺术浪漫的语言
- 配合音符乐器元素
- 采用律动布局
- 每页有音乐元素
- 色调：紫色、粉色、黄色"""
    }

    style_instruction = style_prompts.get(template_style, "请使用专业教育风格")

    system_prompt = f"""你是一位资深教育专家和PPT设计师，擅长设计{subject}学科的教学内容。

当前使用的模板风格：{template_style}
{style_instruction}

请根据以下要求生成{content_type}的PPT内容：

要求：
1. 每页内容简洁精炼，适合课堂演示（每页不超过5个要点）
2. 每个要点控制在15-20字以内
3. 使用项目符号和短句，避免大段文字
4. 每页都要有视觉元素说明（插图、图标、动画等）
5. 内容要生动有趣，适合{grade}学生

请严格按照以下JSON格式输出：
{{
    "title": "PPT标题",
    "subtitle": "副标题",
    "slides": [
        {{
            "title": "幻灯片标题",
            "content": ["要点1", "要点2", "要点3"],
            "visual_elements": "视觉元素说明（如：水墨山水背景、毛笔字体标题）",
            "design_notes": "设计建议（如：采用竖排布局，配合印章元素）",
            "notes": "演讲者备注"
        }}
    ]
}}"""

    user_message = f"""请为以下主题生成{content_type}的PPT内容：

主题：{topic}
学科：{subject}
年级：{grade}
模板风格：{template_style}

{search_context}

请根据搜索结果润色内容，确保内容准确、生动、适合课堂演示。"""

    try:
        response = client.chat.completions.create(
            model=config.OPENAI_MODEL,
            max_tokens=4096,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_message},
            ],
        )

        content = response.choices[0].message.content

        # 解析JSON
        try:
            # 尝试直接解析
            ppt_data = json.loads(content)
        except json.JSONDecodeError:
            # 尝试从markdown代码块中提取
            if "```json" in content:
                json_str = content.split("```json")[1].split("```")[0].strip()
                ppt_data = json.loads(json_str)
            elif "```" in content:
                json_str = content.split("```")[1].split("```")[0].strip()
                ppt_data = json.loads(json_str)
            else:
                raise ValueError("无法解析AI返回的内容")

        return ppt_data

    except Exception as e:
        logger.error(f"生成内容失败: {e}")
        # 返回默认内容
        return {
            "title": topic,
            "subtitle": f"{subject} {grade} {content_type}",
            "slides": [
                {
                    "title": "导入",
                    "content": ["引入话题", "激发兴趣", "明确目标"],
                    "visual_elements": "标题页背景",
                    "design_notes": "简洁大方",
                    "notes": ""
                }
            ]
        }


def apply_template_to_project(project_path: str, template_name: str) -> bool:
    """
    将模板应用到项目

    参数:
        project_path: 项目路径
        template_name: 模板名称（中文）

    返回:
        是否成功
    """
    import shutil

    project_path = Path(project_path)

    # 教育模板目录
    education_dir = SKILL_DIR / "templates" / "education"

    if not education_dir.exists():
        logger.warning(f"教育模板目录不存在: {education_dir}")
        return False

    # 遍历目录查找匹配的模板
    for template_file in education_dir.glob("*.pptx"):
        # 精确匹配（去掉.pptx后缀）
        file_name = template_file.stem
        if file_name == template_name:
            dst = project_path / "template.pptx"
            shutil.copy2(template_file, dst)
            logger.info(f"应用模板: {template_name} (精确匹配)")
            return True

    # 模糊匹配
    for template_file in education_dir.glob("*.pptx"):
        file_name = template_file.stem
        if template_name in file_name or file_name in template_name:
            dst = project_path / "template.pptx"
            shutil.copy2(template_file, dst)
            logger.info(f"应用模板: {file_name} (模糊匹配)")
            return True

    logger.warning(f"未找到模板: {template_name}")
    return False


def replace_text_in_template(template_path: str, replacements: Dict[str, str], output_path: str) -> str:
    """
    替换模板中的文字内容

    参数:
        template_path: 模板文件路径
        replacements: 替换映射 {旧文字: 新文字}
        output_path: 输出文件路径

    返回:
        输出文件路径
    """
    from pptx import Presentation

    prs = Presentation(template_path)

    # 遍历所有幻灯片
    for slide in prs.slides:
        # 遍历所有形状
        for shape in slide.shapes:
            # 如果有文本框
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # 替换文字
                        for old_text, new_text in replacements.items():
                            if old_text in run.text:
                                run.text = run.text.replace(old_text, new_text)

            # 如果有表格
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                for old_text, new_text in replacements.items():
                                    if old_text in run.text:
                                        run.text = run.text.replace(old_text, new_text)

    # 保存
    prs.save(output_path)
    logger.info(f"模板文字替换完成: {output_path}")

    return output_path


def generate_ppt_from_template(
    template_name: str,
    topic: str,
    subject: str,
    grade: str,
    content_type: str,
    outline_markdown: str = ""
) -> Tuple[str, str]:
    """
    直接套用模板生成PPT，按幻灯片位置替换文字内容

    参数:
        template_name: 模板名称
        topic: 主题
        subject: 学科
        grade: 年级
        content_type: 内容类型
        outline_markdown: 大纲内容

    返回:
        (pptx_path, title)
    """
    from pptx import Presentation
    import shutil

    logger.info(f"=== 开始模板套用PPT生成 ===")
    logger.info(f"模板名称: {template_name}")
    logger.info(f"主题: {topic}")
    logger.info(f"学科: {subject}")
    logger.info(f"年级: {grade}")
    logger.info(f"大纲长度: {len(outline_markdown) if outline_markdown else 0}")

    # 教育模板目录
    education_dir = SKILL_DIR / "templates" / "education"
    logger.info(f"模板目录: {education_dir}")
    logger.info(f"模板目录存在: {education_dir.exists()}")

    # 查找模板文件
    template_file = None
    for f in education_dir.glob("*.pptx"):
        logger.info(f"检查模板文件: {f.stem}")
        if f.stem == template_name:
            template_file = f
            break

    if not template_file:
        # 模糊匹配
        for f in education_dir.glob("*.pptx"):
            if template_name in f.stem or f.stem in template_name:
                template_file = f
                break

    if not template_file:
        raise FileNotFoundError(f"未找到模板: {template_name}")

    logger.info(f"使用模板: {template_file.name}")

    # 创建项目目录
    project_name = f"{subject}_{grade}_{topic}" if subject and grade else topic
    project_name = f"{content_type}_{project_name}"
    safe_name = "".join(c for c in project_name if c.isalnum() or c in (' ', '-', '_')).strip()[:50]

    project_path = PROJECTS_DIR / safe_name
    project_path.mkdir(parents=True, exist_ok=True)
    logger.info(f"项目目录: {project_path}")

    # 生成PPT内容
    if outline_markdown:
        # 根据大纲生成内容
        logger.info("根据大纲生成PPT内容...")
        ppt_data = generate_content_from_outline(
            outline_markdown=outline_markdown,
            topic=topic,
            subject=subject,
            grade=grade,
            content_type=content_type
        )
    else:
        # 生成默认内容
        logger.info("使用默认内容...")
        ppt_data = {
            "title": topic,
            "subtitle": f"{subject} {grade} {content_type}",
            "slides": [
                {"title": "导入", "content": ["引入话题", "激发兴趣", "明确目标"]},
                {"title": "知识讲解", "content": ["核心概念", "重点难点", "案例分析"]},
                {"title": "课堂练习", "content": ["练习题目", "互动问答", "小组讨论"]},
                {"title": "课堂小结", "content": ["重点回顾", "知识梳理", "拓展延伸"]}
            ]
        }

    logger.info(f"PPT数据生成完成，标题: {ppt_data.get('title', '无')}")
    logger.info(f"幻灯片数量: {len(ppt_data.get('slides', []))}")

    # 直接在模板上替换文字（按幻灯片位置）
    output_path = str(project_path / f"{topic}.pptx")
    replace_template_content_by_position(str(template_file), ppt_data, output_path)

    logger.info(f"PPT生成完成: {output_path}")

    return output_path, topic


def replace_template_content_by_position(template_path: str, ppt_data: Dict, output_path: str) -> str:
    """
    基于模板风格生成完整PPT，保留模板原有设计

    参数:
        template_path: 模板文件路径
        ppt_data: PPT数据 {title, subtitle, slides: [{title, content: [...]}]}
        output_path: 输出文件路径

    返回:
        输出文件路径
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    logger.info(f"=== 开始生成PPT ===")

    prs = Presentation(template_path)
    slides_data = ppt_data.get("slides", [])
    title_text = ppt_data.get("title", "")
    subtitle_text = ppt_data.get("subtitle", "")

    # 从模板提取样式信息
    style = extract_template_style(prs)
    logger.info(f"模板样式: {style}")

    # 替换封面页文字
    slide0 = prs.slides[0]
    for shape in slide0.shapes:
        if not shape.has_text_frame:
            continue
        full_text = shape.text_frame.text.strip()
        if "模板" in full_text:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if "模板" in run.text:
                        run.text = title_text
        if "从百草园到三味书屋" in full_text or "趣味导入" in full_text:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if "从百草园到三味书屋" in run.text or "趣味导入" in run.text:
                        run.text = subtitle_text

    # 替换目录页文字
    slide1 = prs.slides[1]
    text_shapes = [s for s in slide1.shapes if s.has_text_frame]
    arabic_numbers = ["01", "02", "03", "04", "05"]
    section_map = {}

    for i, shape in enumerate(text_shapes):
        text = shape.text_frame.text.strip()
        if text in arabic_numbers:
            section_idx = int(text) - 1
            if i + 2 < len(text_shapes):
                section_map[section_idx] = (text_shapes[i + 1], text_shapes[i + 2])

    for section_idx, (title_shape, content_shape) in section_map.items():
        if section_idx >= len(slides_data):
            break
        slide_data = slides_data[section_idx]
        new_title = slide_data.get("title", "")
        new_content = slide_data.get("content", [])
        content_text = "；".join(new_content[:3]) if new_content else ""

        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = new_title
        for paragraph in content_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = content_text

    # 创建内容页 - 使用模板的背景和颜色
    for slide_idx, slide_data in enumerate(slides_data):
        new_slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 使用模板背景色
        background = new_slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = style["bg"]

        # 标题 - 使用模板主色
        title_box = new_slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.9))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "")
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = style["primary"]

        # 内容区域
        content_box = new_slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2))
        tf = content_box.text_frame
        tf.word_wrap = True

        content_items = slide_data.get("content", [])
        for i, point in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(16)
            p.font.color.rgb = style["text"]
            p.space_after = Pt(14)

    # 结束页
    end_slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = end_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = style["primary"]

    thank_box = end_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]
    p.text = "谢谢观看"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # 保存
    prs.save(output_path)
    logger.info(f"PPT生成完成，共{len(prs.slides)}页")

    return output_path


def extract_template_style(prs) -> dict:
    """从模板提取样式信息"""
    from pptx.dml.color import RGBColor

    style = {
        "bg": RGBColor(0xFF, 0xF8, 0xDC),
        "primary": RGBColor(0x8B, 0x45, 0x13),
        "text": RGBColor(0x33, 0x33, 0x33),
        "accent": RGBColor(0xC0, 0x39, 0x2B),
    }

    try:
        # 从封面页提取背景色
        if len(prs.slides) > 0:
            slide = prs.slides[0]
            bg = slide.background
            if bg.fill and bg.fill.fore_color and bg.fill.fore_color.rgb:
                style["bg"] = bg.fill.fore_color.rgb

            # 从文本提取主色
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.color and run.font.color.rgb:
                                style["primary"] = run.font.color.rgb
                                break
    except:
        pass

    return style


def generate_content_from_outline(
    outline_markdown: str,
    topic: str,
    subject: str,
    grade: str,
    content_type: str
) -> Dict:
    """
    根据大纲生成润色后的PPT内容数据

    参数:
        outline_markdown: 大纲内容
        topic: 主题
        subject: 学科
        grade: 年级
        content_type: 内容类型

    返回:
        PPT数据字典
    """
    from openai import OpenAI
    import sys
    sys.path.insert(0, str(Path(__file__).parent.parent))
    import config

    client = OpenAI(
        api_key=config.OPENAI_API_KEY,
        base_url=config.OPENAI_BASE_URL,
    )

    # 搜索相关资料用于润色
    search_context = ""
    try:
        from agent.web_search import search_web
        import asyncio
        queries = [f"{topic} 教学设计 教案", f"{topic} 课文赏析 知识点"]
        for q in queries[:2]:
            try:
                results = asyncio.run(search_web(q, 3))
                if results:
                    for r in results[:2]:
                        if r.get('snippet'):
                            search_context += f"- {r['snippet'][:100]}\n"
            except:
                pass
    except:
        pass

    system_prompt = f"""你是一位资深{subject}教师和教学内容设计师，擅长为{grade}学生编写生动有趣的教学PPT内容。

你的任务是根据大纲框架，**润色并充实**每一页PPT的具体文字内容。

要求：
1. 大纲只是方向指引，你需要根据大纲**展开详细描写**
2. 每个要点要具体、生动、有细节，不能只是大纲的简单重复
3. 适当引用原文、举例说明、加入教学引导语
4. 语言要适合{grade}学生理解，生动有趣
5. 每个幻灯片4-5个要点，每个要点30-60字
6. 包含具体的分析、例子、引文等，让内容有深度

{f"参考资料：{search_context}" if search_context else ""}

请严格按照以下JSON格式输出：
{{
    "title": "PPT标题",
    "subtitle": "副标题",
    "slides": [
        {{
            "title": "章节标题",
            "content": ["润色后的详细要点1", "润色后的详细要点2", "润色后的详细要点3", "润色后的详细要点4"]
        }}
    ]
}}"""

    user_message = f"""请根据以下大纲，为{grade}{subject}课《{topic}》生成润色后的PPT内容。

大纲框架（仅供参考，需要展开详写）：
{outline_markdown}

请注意：
- 大纲中的每个章节都要展开为详细的讲解内容
- 要有具体的例子、引文、分析
- 语言要生动，适合课堂使用
- 生成10-15个幻灯片"""

    try:
        response = client.chat.completions.create(
            model=config.OPENAI_MODEL,
            max_tokens=4096,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_message},
            ],
        )

        content = response.choices[0].message.content

        # 解析JSON
        try:
            return json.loads(content)
        except json.JSONDecodeError:
            if "```json" in content:
                json_str = content.split("```json")[1].split("```")[0].strip()
                return json.loads(json_str)
            elif "```" in content:
                json_str = content.split("```")[1].split("```")[0].strip()
                return json.loads(json_str)

    except Exception as e:
        logger.error(f"生成内容失败: {e}")

    # 返回默认内容
    return parse_outline_to_simple_data(outline_markdown, topic, subject, grade)


def parse_outline_to_simple_data(outline_markdown: str, topic: str, subject: str, grade: str) -> Dict:
    """
    解析大纲为简单的PPT数据
    """
    import re

    slides = []
    current_slide = None

    lines = outline_markdown.split('\n')

    for line in lines:
        line = line.strip()
        if not line:
            continue

        # 标题行
        if line.startswith('#'):
            if current_slide:
                slides.append(current_slide)

            title = re.sub(r'^#+\s*', '', line)
            current_slide = {
                "title": title,
                "content": []
            }

        # 列表项
        elif line.startswith(('-', '*', '1.', '2.', '3.', '4.', '5.')):
            if current_slide:
                content = re.sub(r'^[-*]\s*', '', line)
                content = re.sub(r'^\d+\.\s*', '', content)
                if content:
                    current_slide["content"].append(content)

    if current_slide:
        slides.append(current_slide)

    # 如果没有解析到内容，创建默认结构（10-15个章节）
    if not slides:
        slides = [
            {"title": "学习目标", "content": ["明确本课学习目标", "了解核心知识点", "掌握重点难点", "培养分析能力"]},
            {"title": "作者介绍", "content": ["了解作者生平背景", "认识作者创作特点", "理解作者思想情感", "学习作者写作态度"]},
            {"title": "写作背景", "content": ["了解时代背景", "认识创作动机", "理解历史语境", "把握作品主题"]},
            {"title": "导入新课", "content": ["创设情境激发兴趣", "回顾相关知识", "引出本课主题", "明确学习方向"]},
            {"title": "课文精讲", "content": ["分析文本核心内容", "解读重点段落", "品味语言特色", "理解深层含义"]},
            {"title": "重点段落分析", "content": ["找出关键语句", "分析写作手法", "理解表达效果", "体会作者情感"]},
            {"title": "写作手法", "content": ["学习修辞手法", "分析叙事技巧", "理解描写方法", "掌握表达方式"]},
            {"title": "难点解析", "content": ["突破理解难点", "分析疑难语句", "探讨深层含义", "解决学习困惑"]},
            {"title": "拓展延伸", "content": ["联系现实生活", "对比阅读分析", "深化主题理解", "培养思辨能力"]},
            {"title": "课堂互动", "content": ["小组讨论交流", "问题探究思考", "观点分享碰撞", "合作学习提升"]},
            {"title": "课堂小结", "content": ["总结本课重点", "梳理知识脉络", "巩固学习成果", "布置课后作业"]}
        ]

    return {
        "title": topic,
        "subtitle": f"{subject} {grade} 教学课件",
        "slides": slides
    }


def build_replacements_from_data(ppt_data: Dict) -> Dict[str, str]:
    """
    根据PPT数据构建替换映射

    参数:
        ppt_data: PPT数据

    返回:
        替换映射字典
    """
    replacements = {}

    # 替换标题
    if ppt_data.get("title"):
        replacements["标题"] = ppt_data["title"]
        replacements["主题"] = ppt_data["title"]
        replacements["课题"] = ppt_data["title"]

    if ppt_data.get("subtitle"):
        replacements["副标题"] = ppt_data["subtitle"]

    # 替换幻灯片内容
    slides = ppt_data.get("slides", [])

    for i, slide in enumerate(slides):
        slide_title = slide.get("title", "")
        content = slide.get("content", [])

        # 替换页面标题
        if slide_title:
            replacements[f"第{i+1}页标题"] = slide_title
            # 尝试替换常见的占位符
            if i == 0:
                replacements["导入"] = slide_title
                replacements["课程导入"] = slide_title
            elif i == 1:
                replacements["知识讲解"] = slide_title
                replacements["新课讲授"] = slide_title
            elif i == 2:
                replacements["课堂练习"] = slide_title
                replacements["巩固练习"] = slide_title
            elif i == 3:
                replacements["课堂小结"] = slide_title
                replacements["总结"] = slide_title

        # 替换内容要点
        for j, point in enumerate(content):
            if point:
                # 尝试替换常见的占位符
                replacements[f"要点{i+1}-{j+1}"] = point
                replacements[f"内容{i+1}-{j+1}"] = point

    return replacements


def create_design_spec(project_path: str, ppt_data: Dict, template_info: Dict) -> str:
    """
    创建设计规范文件

    参数:
        project_path: 项目路径
        ppt_data: PPT内容数据
        template_info: 模板信息

    返回:
        设计规范文件路径
    """
    project_path = Path(project_path)

    # 生成设计规范
    spec_content = f"""# {ppt_data.get('title', 'PPT')} - 设计规范

## 模板风格信息

- **模板名称**: {template_info.get('name', '默认')}
- **风格类型**: {template_info.get('style', '默认')}
- **主色调**: {template_info.get('colors', {}).get('primary', '#2C3E50')}
- **辅助色**: {template_info.get('colors', {}).get('secondary', '#3498DB')}
- **强调色**: {template_info.get('colors', {}).get('accent', '#E74C3C')}
- **背景色**: {template_info.get('colors', {}).get('background', '#FFFFFF')}

## 设计元素

{chr(10).join('- ' + elem for elem in template_info.get('elements', ['默认元素']))}

## 氛围特点

{template_info.get('atmosphere', '专业教育风格')}

## 页面设计建议

"""

    # 添加每页的设计建议
    for i, slide in enumerate(ppt_data.get('slides', []), 1):
        spec_content += f"""### 第{i}页: {slide.get('title', '')}

- **视觉元素**: {slide.get('visual_elements', '默认布局')}
- **设计建议**: {slide.get('design_notes', '简洁大方')}
- **演讲者备注**: {slide.get('notes', '')}

"""

    # 保存设计规范
    spec_path = project_path / "design_spec.md"
    spec_path.write_text(spec_content, encoding='utf-8')
    logger.info(f"设计规范已生成: {spec_path}")

    return str(spec_path)


def generate_enhanced_education_ppt(
    topic: str,
    subject: str = "",
    grade: str = "",
    difficulty: str = "中等",
    content_type: str = "课件",
    template_id: str = "",
    outline_markdown: str = ""
) -> Tuple[str, str]:
    """
    生成增强版教育PPT

    参数:
        topic: 主题
        subject: 学科
        grade: 年级
        difficulty: 难度
        content_type: 内容类型
        template_id: 模板ID
        outline_markdown: 大纲内容（Markdown格式）

    返回:
        (pptx_path, title)
    """
    from agent.education_templates import EDUCATION_TEMPLATES, get_template_by_name

    # 获取模板信息
    template_info = None
    for key, tmpl in EDUCATION_TEMPLATES.items():
        if tmpl["id"] == template_id or key == template_id:
            template_info = tmpl
            break

    if not template_info:
        template_info = {
            "id": "default",
            "name": "默认",
            "style": "专业教育风格",
            "colors": {"primary": "#2C3E50", "secondary": "#3498DB", "accent": "#E74C3C", "background": "#FFFFFF"},
            "elements": ["默认元素"],
            "atmosphere": "专业严谨"
        }

    logger.info(f"使用模板: {template_info['name']} ({template_info['id']})")

    # 步骤1：搜索相关内容
    logger.info("搜索相关教学内容...")
    search_results = search_content_for_topic(topic, subject)

    # 步骤2：生成优化内容（如果有大纲，使用大纲内容）
    logger.info("根据模板风格生成内容...")

    if outline_markdown:
        # 使用大纲内容生成PPT数据
        logger.info("使用大纲内容生成PPT")
        ppt_data = generate_ppt_data_from_outline(
            outline_markdown=outline_markdown,
            topic=topic,
            subject=subject,
            grade=grade,
            content_type=content_type,
            template_style=template_info['id'],
            search_results=search_results
        )
    else:
        # 生成新内容
        ppt_data = generate_content_with_style(
            topic=topic,
            subject=subject,
            grade=grade,
            content_type=content_type,
            template_style=template_info['id'],
            search_results=search_results
        )

    # 步骤3：创建项目
    project_name = f"{subject}_{grade}_{topic}" if subject and grade else topic
    project_name = f"{content_type}_{project_name}"

    logger.info(f"创建项目: {project_name}")
    from agent.ppt_master_integration import create_project
    project_path = create_project(project_name)

    # 步骤4：应用模板
    logger.info(f"应用模板: {template_info['name']}")
    apply_template_to_project(project_path, template_info['id'])

    # 步骤5：生成设计规范
    logger.info("生成设计规范...")
    create_design_spec(project_path, ppt_data, template_info)

    # 步骤6：保存内容数据
    content_path = Path(project_path) / "content.json"
    content_path.write_text(json.dumps(ppt_data, ensure_ascii=False, indent=2), encoding='utf-8')
    logger.info(f"内容数据已保存: {content_path}")

    # 步骤7：生成PPTX
    logger.info("生成PPTX文件...")
    pptx_path = generate_pptx_from_data(project_path, ppt_data, template_info)

    title = ppt_data.get('title', topic)

    return pptx_path, title


def generate_ppt_data_from_outline(
    outline_markdown: str,
    topic: str,
    subject: str,
    grade: str,
    content_type: str,
    template_style: str,
    search_results: Dict
) -> Dict:
    """
    根据大纲内容生成PPT数据

    参数:
        outline_markdown: 大纲内容（Markdown格式）
        topic: 主题
        subject: 学科
        grade: 年级
        content_type: 内容类型
        template_style: 模板风格
        search_results: 搜索结果

    返回:
        PPT数据字典
    """
    from openai import OpenAI
    import config

    client = OpenAI(
        api_key=config.OPENAI_API_KEY,
        base_url=config.OPENAI_BASE_URL,
    )

    # 格式化搜索结果
    search_context = ""
    if search_results:
        search_context = "以下是网络搜索到的相关教学资料：\n\n"
        for query, results in search_results.items():
            if results:
                search_context += f"【{query}】\n"
                for r in results[:2]:
                    search_context += f"- {r.get('title', '')}: {r.get('snippet', '')}\n"
                search_context += "\n"

    system_prompt = f"""你是一位资深教育专家，擅长根据大纲设计{subject}学科的教学PPT内容。

当前使用的模板风格：{template_style}

请根据以下大纲内容，生成详细的PPT内容：

要求：
1. 严格按照大纲结构生成内容
2. 每页内容简洁精炼，适合课堂演示（每页不超过5个要点）
3. 每个要点控制在15-20字以内
4. 使用项目符号和短句，避免大段文字
5. 每页都要有视觉元素说明和设计建议
6. 内容要生动有趣，适合{grade}学生

请严格按照以下JSON格式输出：
{{
    "title": "PPT标题",
    "subtitle": "副标题",
    "slides": [
        {{
            "title": "幻灯片标题",
            "content": ["要点1", "要点2", "要点3"],
            "visual_elements": "视觉元素说明",
            "design_notes": "设计建议",
            "notes": "演讲者备注"
        }}
    ]
}}"""

    user_message = f"""请根据以下大纲生成{content_type}的PPT内容：

主题：{topic}
学科：{subject}
年级：{grade}

大纲内容：
{outline_markdown}

{search_context}

请根据大纲内容生成详细的PPT幻灯片内容，确保每页内容简洁、生动、适合课堂演示。"""

    try:
        response = client.chat.completions.create(
            model=config.OPENAI_MODEL,
            max_tokens=4096,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_message},
            ],
        )

        content = response.choices[0].message.content

        # 解析JSON
        try:
            ppt_data = json.loads(content)
        except json.JSONDecodeError:
            if "```json" in content:
                json_str = content.split("```json")[1].split("```")[0].strip()
                ppt_data = json.loads(json_str)
            elif "```" in content:
                json_str = content.split("```")[1].split("```")[0].strip()
                ppt_data = json.loads(json_str)
            else:
                raise ValueError("无法解析AI返回的内容")

        return ppt_data

    except Exception as e:
        logger.error(f"根据大纲生成内容失败: {e}")
        # 返回基于大纲的默认内容
        return parse_outline_to_ppt_data(outline_markdown, topic, subject, grade)


def parse_outline_to_ppt_data(outline_markdown: str, topic: str, subject: str, grade: str) -> Dict:
    """
    解析大纲Markdown为PPT数据

    参数:
        outline_markdown: 大纲内容
        topic: 主题
        subject: 学科
        grade: 年级

    返回:
        PPT数据字典
    """
    import re

    slides = []
    current_slide = None

    # 简单解析Markdown大纲
    lines = outline_markdown.split('\n')

    for line in lines:
        line = line.strip()
        if not line:
            continue

        # 标题行（# ## ###）
        if line.startswith('#'):
            if current_slide:
                slides.append(current_slide)

            # 提取标题
            title = re.sub(r'^#+\s*', '', line)
            current_slide = {
                "title": title,
                "content": [],
                "visual_elements": "默认布局",
                "design_notes": "简洁大方",
                "notes": ""
            }

        # 列表项（- * 1.）
        elif line.startswith(('-', '*', '1.', '2.', '3.', '4.', '5.')):
            if current_slide:
                # 提取内容
                content = re.sub(r'^[-*]\s*', '', line)
                content = re.sub(r'^\d+\.\s*', '', content)
                if content:
                    current_slide["content"].append(content)

    # 添加最后一个幻灯片
    if current_slide:
        slides.append(current_slide)

    # 如果没有解析到内容，创建默认结构
    if not slides:
        slides = [
            {
                "title": "导入",
                "content": ["引入话题", "激发兴趣", "明确目标"],
                "visual_elements": "标题页背景",
                "design_notes": "简洁大方",
                "notes": ""
            },
            {
                "title": "知识讲解",
                "content": ["核心概念", "重点难点", "案例分析"],
                "visual_elements": "图文混排",
                "design_notes": "层次分明",
                "notes": ""
            },
            {
                "title": "课堂练习",
                "content": ["练习题目", "互动问答", "小组讨论"],
                "visual_elements": "卡片式布局",
                "design_notes": "清晰易读",
                "notes": ""
            },
            {
                "title": "课堂小结",
                "content": ["重点回顾", "知识梳理", "拓展延伸"],
                "visual_elements": "总结图表",
                "design_notes": "简洁明了",
                "notes": ""
            }
        ]

    return {
        "title": topic,
        "subtitle": f"{subject} {grade} 教学课件",
        "slides": slides
    }


def generate_pptx_from_data(project_path: str, ppt_data: Dict, template_info: Dict) -> str:
    """
    根据数据生成PPTX文件，使用模板样式

    参数:
        project_path: 项目路径
        ppt_data: PPT内容数据
        template_info: 模板信息

    返回:
        PPTX文件路径
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    project_path = Path(project_path)

    # 检查是否有模板文件
    template_file = project_path / "template.pptx"

    if template_file.exists():
        # 使用模板文件
        logger.info(f"使用模板文件: {template_file}")
        prs = Presentation(str(template_file))

        # 清空模板中的示例内容，保留样式
        # 注意：保留第一页（封面）和最后一页（结束页）的布局
    else:
        # 创建新的PPT
        logger.info("未找到模板，创建新PPT")
        prs = Presentation()

    # 设置16:9比例
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 获取颜色
    colors = template_info.get('colors', {})
    primary_color = colors.get('primary', '#2C3E50')
    secondary_color = colors.get('secondary', '#3498DB')
    accent_color = colors.get('accent', '#E74C3C')
    bg_color = colors.get('background', '#FFFFFF')

    # 转换颜色为RGB
    def hex_to_rgb(hex_color):
        hex_color = hex_color.lstrip('#')
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

    primary_rgb = hex_to_rgb(primary_color)
    secondary_rgb = hex_to_rgb(secondary_color)
    accent_rgb = hex_to_rgb(accent_color)
    bg_rgb = hex_to_rgb(bg_color)

    # 如果使用模板，清空现有幻灯片（保留布局）
    if template_file.exists():
        # 删除所有幻灯片，重新创建
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # 创建封面页
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 设置背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = primary_rgb

    # 添加装饰元素（根据模板风格）
    template_id = template_info.get('id', '')

    # 国风水墨风格 - 添加印章装饰
    if template_id == 'ink_painting':
        seal = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(11), Inches(2), Inches(1.2), Inches(1.2)
        )
        seal.fill.solid()
        seal.fill.fore_color.rgb = hex_to_rgb('#C0392B')
        seal.line.fill.background()

    # 马卡龙卡通风格 - 添加装饰圆形
    elif template_id == 'macaron_cartoon':
        for pos in [(1, 1, '#FFD93D'), (11.5, 6, '#C9B1FF')]:
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(pos[0]), Inches(pos[1]), Inches(0.8), Inches(0.8)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = hex_to_rgb(pos[2])
            circle.line.fill.background()

    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = ppt_data.get('title', '演示文稿')
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # 添加副标题
    if ppt_data.get('subtitle'):
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = ppt_data['subtitle']
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
        p.alignment = PP_ALIGN.CENTER

    # 创建内容页
    for i, slide_data in enumerate(ppt_data.get('slides', [])):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 设置背景色
        background = slide.background
        fill = background.fill
        fill.solid()

        # 根据模板风格设置背景
        if template_id in ['tech_neon']:
            fill.fore_color.rgb = hex_to_rgb('#0D0D0D')
        elif template_id in ['ink_painting']:
            fill.fore_color.rgb = hex_to_rgb('#F5F5DC')
        elif template_id in ['macaron_cartoon']:
            fill.fore_color.rgb = hex_to_rgb('#FFF5F5') if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
        else:
            fill.fore_color.rgb = bg_rgb if i % 2 == 0 else RGBColor(0xF5, 0xF5, 0xF5)

        # 添加左侧装饰条（根据模板风格）
        if template_id == 'ink_painting':
            left_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.8), Inches(0.8), Inches(0.08), Inches(5.9)
            )
            left_bar.fill.solid()
            left_bar.fill.fore_color.rgb = hex_to_rgb('#C0392B')
            left_bar.line.fill.background()
        elif template_id in ['fresh_forest', 'lab_tech']:
            left_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(0.3), Inches(7.5)
            )
            left_bar.fill.solid()
            left_bar.fill.fore_color.rgb = primary_rgb
            left_bar.line.fill.background()

        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.get('title', f'第{i+1}页')
        p.font.size = Pt(32)
        p.font.bold = True

        # 根据模板设置标题颜色
        if template_id == 'tech_neon':
            p.font.color.rgb = hex_to_rgb('#00FF88')
        elif template_id == 'ink_painting':
            p.font.color.rgb = hex_to_rgb('#2C3E50')
        else:
            p.font.color.rgb = primary_rgb

        p.alignment = PP_ALIGN.LEFT

        # 添加装饰线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(1.4), Inches(2), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = accent_rgb
        line.line.fill.background()

        # 添加内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for j, point in enumerate(slide_data.get('content', [])):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # 根据模板风格调整内容格式
            if template_id == 'ink_painting':
                p.text = point  # 国风不使用项目符号
            elif template_id == 'tech_neon':
                p.text = f">>> {point}"  # 科技风使用>>>符号
            else:
                p.text = f"• {point}"

            p.font.size = Pt(20)

            # 根据模板设置文字颜色
            if template_id == 'tech_neon':
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            else:
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(12)

        # 添加设计说明（作为备注）
        if slide_data.get('design_notes') or slide_data.get('visual_elements'):
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_text = ""
            if slide_data.get('visual_elements'):
                notes_text += f"视觉元素：{slide_data['visual_elements']}\n"
            if slide_data.get('design_notes'):
                notes_text += f"设计建议：{slide_data['design_notes']}"
            notes_tf.text = notes_text

    # 创建结束页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = primary_rgb

    # 添加感谢文字
    thank_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]

    # 根据模板风格调整结束语
    if template_id == 'tech_neon':
        p.text = "THANKS"
    elif template_id == 'macaron_cartoon':
        p.text = "太棒了！"
    elif template_id == 'warm_healing':
        p.text = "一起加油！"
    else:
        p.text = "谢谢观看"

    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # 保存文件
    output_path = project_path / "output.pptx"
    prs.save(str(output_path))
    logger.info(f"PPTX已生成: {output_path}")

    return str(output_path)
