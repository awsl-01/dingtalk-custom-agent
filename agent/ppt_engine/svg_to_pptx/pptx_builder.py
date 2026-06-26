"""
PPT Engine - PPTX构建器

将SVG页面构建为PowerPoint演示文稿。
"""

import os
from pathlib import Path
from typing import Dict, List, Any, Optional
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from .drawingml_converter import SVGToDrawingMLConverter


class PPTXBuilder:
    """PPTX构建器"""

    def __init__(self, canvas_format: str = 'ppt169'):
        """
        初始化PPTX构建器

        参数:
            canvas_format: 画布格式（ppt169/ppt43/xhs/story）
        """
        self.canvas_format = canvas_format

        # 画布尺寸（EMU）
        self.canvas_sizes = {
            'ppt169': (12192000, 6858000),  # 16:9
            'ppt43': (9144000, 6858000),    # 4:3
            'xhs': (6858000, 9144000),      # 竖版
            'story': (6858000, 12192000)    # 竖版长图
        }

        self.width, self.height = self.canvas_sizes.get(canvas_format, self.canvas_sizes['ppt169'])

        # 创建演示文稿
        self.prs = Presentation()
        self.prs.slide_width = self.width
        self.prs.slide_height = self.height

        # SVG转换器
        self.converter = SVGToDrawingMLConverter(self.width, self.height)

    def add_svg_slide(self, svg_path: str, notes: str = None) -> bool:
        """
        添加SVG幻灯片

        参数:
            svg_path: SVG文件路径
            notes: 演讲备注

        返回:
            是否成功
        """
        try:
            # 转换SVG
            result = self.converter.convert_svg_file(svg_path)
            if not result['success']:
                print(f"⚠️ SVG转换失败: {result.get('error')}")
                return False

            # 添加空白幻灯片
            slide_layout = self.prs.slide_layouts[6]  # 空白布局
            slide = self.prs.slides.add_slide(slide_layout)

            # 添加元素
            self._add_elements(slide, result['elements'])

            # 添加备注
            if notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes

            return True

        except Exception as e:
            print(f"[ERROR] Add slide failed: {e}")
            return False

    def _add_elements(self, slide, elements: List[Dict[str, Any]]):
        """递归添加元素"""
        for elem in elements:
            elem_type = elem.get('type')

            if elem_type == 'rect':
                self._add_rect(slide, elem)
            elif elem_type == 'text':
                self._add_text(slide, elem)
            elif elem_type == 'line':
                self._add_line(slide, elem)
            elif elem_type == 'circle':
                self._add_circle(slide, elem)
            elif elem_type == 'group':
                # 递归处理分组
                self._add_elements(slide, elem.get('children', []))

    def _add_rect(self, slide, elem: Dict[str, Any]):
        """添加矩形"""
        from pptx.enum.shapes import MSO_SHAPE

        x = Emu(int(elem['x'] * 914400 / 96))
        y = Emu(int(elem['y'] * 914400 / 96))
        w = Emu(int(elem['width'] * 914400 / 96))
        h = Emu(int(elem['height'] * 914400 / 96))

        if w == 0 or h == 0:
            return

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)

        # 设置圆角（如果支持）
        rx = elem.get('rx', 0)
        if rx > 0 and len(shape.adjustments) > 0:
            try:
                shape.adjustments[0] = min(rx / min(elem['width'], elem['height']), 0.5)
            except (IndexError, AttributeError):
                pass  # 某些形状不支持adjustments

        # 设置填充
        fill = elem.get('fill', 'none')
        if fill and fill != 'none':
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(fill.lstrip('#'))

        # 设置边框
        stroke = elem.get('stroke', 'none')
        if stroke and stroke != 'none':
            shape.line.color.rgb = RGBColor.from_string(stroke.lstrip('#'))
            shape.line.width = Pt(elem.get('stroke-width', 1))

    def _add_text(self, slide, elem: Dict[str, Any]):
        """添加文本"""
        x = Emu(int(elem['x'] * 914400 / 96))
        y = Emu(int(elem['y'] * 914400 / 96))
        w = Emu(2000000)  # 默认宽度
        h = Emu(200000)   # 默认高度

        text = elem.get('text', '')
        if not text:
            return

        textbox = slide.shapes.add_textbox(x, y, w, h)
        tf = textbox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text

        # 设置字体
        font = p.runs[0].font
        font.size = Pt(elem.get('font-size', 18))
        font.bold = elem.get('font-weight') == 'bold'

        # 设置颜色
        fill = elem.get('fill', '#000000')
        if fill and fill != 'none':
            font.color.rgb = RGBColor.from_string(fill.lstrip('#'))

        # 设置对齐
        anchor = elem.get('text-anchor', 'start')
        if anchor == 'middle':
            p.alignment = PP_ALIGN.CENTER
        elif anchor == 'end':
            p.alignment = PP_ALIGN.RIGHT

    def _add_line(self, slide, elem: Dict[str, Any]):
        """添加线条"""
        from pptx.enum.shapes import MSO_SHAPE

        x1 = Emu(int(elem['x1'] * 914400 / 96))
        y1 = Emu(int(elem['y1'] * 914400 / 96))
        x2 = Emu(int(elem['x2'] * 914400 / 96))
        y2 = Emu(int(elem['y2'] * 914400 / 96))

        # 使用连接符形状
        connector = slide.shapes.add_connector(
            1,  # 直线连接符
            x1, y1, x2, y2
        )

        # 设置样式
        stroke = elem.get('stroke', '#000000')
        if stroke and stroke != 'none':
            connector.line.color.rgb = RGBColor.from_string(stroke.lstrip('#'))
            connector.line.width = Pt(elem.get('stroke-width', 1))

    def _add_circle(self, slide, elem: Dict[str, Any]):
        """添加圆形"""
        from pptx.enum.shapes import MSO_SHAPE

        cx = elem['cx']
        cy = elem['cy']
        r = elem['r']

        x = Emu(int((cx - r) * 914400 / 96))
        y = Emu(int((cy - r) * 914400 / 96))
        w = Emu(int(2 * r * 914400 / 96))
        h = Emu(int(2 * r * 914400 / 96))

        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)

        # 设置填充
        fill = elem.get('fill', 'none')
        if fill and fill != 'none':
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(fill.lstrip('#'))

    def build_from_svg_dir(self, svg_dir: str, notes_dir: str = None,
                           output_path: str = None) -> str:
        """
        从SVG目录构建PPTX

        参数:
            svg_dir: SVG文件目录
            notes_dir: 备注文件目录
            output_path: 输出路径

        返回:
            输出文件路径
        """
        svg_dir = Path(svg_dir)
        if not svg_dir.exists():
            raise FileNotFoundError(f"SVG directory not found: {svg_dir}")

        # 获取所有SVG文件
        svg_files = sorted(svg_dir.glob('slide_*.svg'))

        if not svg_files:
            raise ValueError(f"SVG directory is empty: {svg_dir}")

        print(f"[INFO] Found {len(svg_files)} SVG files")

        # 添加每个SVG
        for svg_file in svg_files:
            # 读取备注
            notes = None
            if notes_dir:
                notes_file = Path(notes_dir) / f"{svg_file.stem}.md"
                if notes_file.exists():
                    notes = notes_file.read_text(encoding='utf-8')

            # 添加幻灯片
            success = self.add_svg_slide(str(svg_file), notes)
            if success:
                print(f"[OK] {svg_file.name}")
            else:
                print(f"[FAIL] {svg_file.name}")

        # 保存PPTX
        if output_path is None:
            output_path = svg_dir.parent / 'exports' / f"{svg_dir.parent.name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        self.prs.save(str(output_path))
        print(f"\n[OK] PPTX saved: {output_path}")

        return str(output_path)

    def save(self, output_path: str):
        """保存PPTX"""
        self.prs.save(output_path)
