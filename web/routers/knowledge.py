"""
知识库管理 API
支持消息和文件的展示、下载、删除、预览、上传
"""
from fastapi import APIRouter, Query, HTTPException, UploadFile, File
from fastapi.responses import FileResponse, PlainTextResponse
from typing import Optional
from pydantic import BaseModel
from datetime import datetime
import json
import os
import sys
import glob
import hashlib
import shutil

# 添加项目根目录到路径
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from web.config import KNOWLEDGE_DIR

router = APIRouter()


# 获取知识库实例的辅助函数
async def _get_knowledge_base(corp_id: str):
    """获取指定企业的知识库实例"""
    from agent.knowledge_base_v2 import get_knowledge_base
    school_dir = os.path.join(KNOWLEDGE_DIR, corp_id)
    return get_knowledge_base(school_dir, corp_id)


# 获取默认企业ID的辅助函数
def _get_default_corp_id() -> str:
    """获取默认企业ID（知识库目录下的第一个企业）"""
    if not os.path.exists(KNOWLEDGE_DIR):
        return ""
    for d in os.listdir(KNOWLEDGE_DIR):
        if os.path.isdir(os.path.join(KNOWLEDGE_DIR, d)) and not d.startswith('.'):
            return d
    return ""


# 角色权限映射：角色可以访问的内容级别
ROLE_PERMISSIONS = {
    "admin": ["public", "internal", "confidential"],      # 管理员：所有内容
    "principal": ["public", "internal", "confidential"],   # 校长：所有内容
    "director": ["public", "internal", "confidential"],    # 主任：所有内容
    "teacher": ["public", "internal"],                     # 教师：公开+内部
    "student": ["public"],                                 # 学生：仅公开
}


@router.get("/list")
async def list_knowledge(
    corp_id: Optional[str] = Query(None, description="企业ID"),
    source_type: Optional[str] = Query(None, description="来源类型: message/file"),
    date: Optional[str] = Query(None, description="日期筛选 YYYY-MM-DD"),
    user_role: Optional[str] = Query(None, description="用户角色（用于权限过滤）"),
    page: int = Query(1, ge=1, description="页码"),
    page_size: int = Query(20, ge=1, le=100, description="每页数量")
):
    """获取知识库内容列表（消息 + 文件）"""
    items = []

    # 获取用户可访问的权限级别
    allowed_permissions = ROLE_PERMISSIONS.get(user_role, ["public"]) if user_role else None

    # 获取所有企业目录
    if corp_id:
        corp_dirs = [corp_id]
    else:
        corp_dirs = [d for d in os.listdir(KNOWLEDGE_DIR)
                     if os.path.isdir(os.path.join(KNOWLEDGE_DIR, d)) and not d.startswith('.')]

    for corp_dir in corp_dirs:
        corp_path = os.path.join(KNOWLEDGE_DIR, corp_dir)

        # 加载权限配置
        permissions = _load_permissions(corp_dir)

        # 加载消息文件
        if source_type is None or source_type == "message":
            messages_dir = os.path.join(corp_path, "messages")
            if os.path.exists(messages_dir):
                items.extend(_scan_messages(messages_dir, corp_dir, date, permissions))

        # 加载文件
        if source_type is None or source_type == "file":
            files_dir = os.path.join(corp_path, "files")
            if os.path.exists(files_dir):
                items.extend(_scan_files(files_dir, corp_dir, date, permissions))

    # 根据用户权限过滤内容
    if allowed_permissions is not None:
        items = [item for item in items if item.get("permission", "public") in allowed_permissions]

    # 按时间倒序排序
    items.sort(key=lambda x: x.get("created_at", ""), reverse=True)

    # 分页
    total = len(items)
    start = (page - 1) * page_size
    end = start + page_size
    page_items = items[start:end]

    return {
        "total": total,
        "page": page,
        "page_size": page_size,
        "items": page_items
    }


def _scan_messages(directory: str, corp_id: str, date_filter: Optional[str] = None, permissions: dict = None) -> list:
    """扫描消息目录"""
    items = []

    for date_dir in os.listdir(directory):
        date_path = os.path.join(directory, date_dir)
        if not os.path.isdir(date_path):
            continue

        # 日期筛选
        if date_filter and date_dir != date_filter:
            continue

        # 扫描 .md 文件
        for file_name in os.listdir(date_path):
            if not file_name.endswith('.md'):
                continue

            file_path = os.path.join(date_path, file_name)
            # 统一使用正斜杠
            file_path = file_path.replace('\\', '/')
            try:
                stat = os.stat(file_path)

                # 读取文件内容
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()

                # 解析内容提取信息
                sender_nick = _extract_sender(content)
                preview = _extract_preview(content)

                # 获取权限
                item_id = file_name.replace('.md', '')
                permission = "public"  # 默认权限
                if permissions and item_id in permissions:
                    permission = permissions[item_id].get("permission", "public")

                items.append({
                    "id": item_id,
                    "type": "message",
                    "name": file_name,
                    "path": file_path,
                    "corp_id": corp_id,
                    "date": date_dir,
                    "sender_nick": sender_nick,
                    "preview": preview,
                    "content": content,
                    "size": stat.st_size,
                    "permission": permission,
                    "created_at": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                })
            except Exception as e:
                print(f"读取消息文件失败: {file_path}, error: {e}")

    return items


def _scan_files(directory: str, corp_id: str, date_filter: Optional[str] = None, permissions: dict = None) -> list:
    """扫描文件目录"""
    items = []

    for date_dir in os.listdir(directory):
        date_path = os.path.join(directory, date_dir)
        if not os.path.isdir(date_path):
            continue

        # 日期筛选
        if date_filter and date_dir != date_filter:
            continue

        # 扫描所有文件
        for file_name in os.listdir(date_path):
            file_path = os.path.join(date_path, file_name)
            if not os.path.isfile(file_path):
                continue

            # 统一使用正斜杠
            file_path = file_path.replace('\\', '/')

            try:
                stat = os.stat(file_path)

                # 判断文件类型
                ext = os.path.splitext(file_name)[1].lower()
                file_type = _get_file_type(ext)

                # 获取权限
                item_id = f"{date_dir}_{file_name}"
                permission = "public"  # 默认权限
                if permissions and item_id in permissions:
                    permission = permissions[item_id].get("permission", "public")

                items.append({
                    "id": item_id,
                    "type": "file",
                    "name": file_name,
                    "path": file_path,
                    "corp_id": corp_id,
                    "date": date_dir,
                    "file_type": file_type,
                    "extension": ext,
                    "size": stat.st_size,
                    "size_display": _format_size(stat.st_size),
                    "permission": permission,
                    "created_at": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                })
            except Exception as e:
                print(f"读取文件信息失败: {file_path}, error: {e}")

    return items


def _extract_sender(content: str) -> str:
    """从消息内容提取发送者"""
    # 格式: ## 💬 启拓 (13:43:10)
    if '## 💬' in content:
        try:
            sender = content.split('## 💬')[1].split('(')[0].strip()
            return sender
        except:
            pass
    return ""


def _extract_preview(content: str) -> str:
    """从消息内容提取预览文本"""
    lines = content.strip().split('\n')
    # 跳过标题行，取正文
    for line in lines[1:]:
        line = line.strip()
        if line and not line.startswith('---'):
            return line[:100]
    return content[:100]


def _load_permissions(corp_id: str) -> dict:
    """加载权限配置"""
    permissions_file = os.path.join(KNOWLEDGE_DIR, corp_id, "permissions.json")

    if not os.path.exists(permissions_file):
        return {}

    try:
        with open(permissions_file, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception as e:
        print(f"加载权限配置失败: {e}")
        return {}


def _get_file_type(ext: str) -> str:
    """根据扩展名判断文件类型"""
    type_map = {
        '.pdf': 'pdf',
        '.doc': 'word', '.docx': 'word',
        '.xls': 'excel', '.xlsx': 'excel',
        '.ppt': 'ppt', '.pptx': 'ppt',
        '.jpg': 'image', '.jpeg': 'image', '.png': 'image', '.gif': 'image',
        '.txt': 'text', '.md': 'text', '.csv': 'text',
        '.zip': 'archive', '.rar': 'archive', '.7z': 'archive',
        '.mp4': 'video', '.avi': 'video', '.mov': 'video',
        '.mp3': 'audio', '.wav': 'audio',
    }
    return type_map.get(ext, 'other')


def _format_size(size_bytes: int) -> str:
    """格式化文件大小"""
    if size_bytes < 1024:
        return f"{size_bytes} B"
    elif size_bytes < 1024 * 1024:
        return f"{size_bytes / 1024:.1f} KB"
    else:
        return f"{size_bytes / (1024 * 1024):.1f} MB"


@router.get("/stats")
async def get_knowledge_stats(corp_id: Optional[str] = None):
    """获取知识库统计"""
    stats = {
        "total_messages": 0,
        "total_files": 0,
        "total_size": 0,
        "dates": {},
        "file_types": {},
    }

    # 获取所有企业目录
    if corp_id:
        corp_dirs = [corp_id]
    else:
        corp_dirs = [d for d in os.listdir(KNOWLEDGE_DIR)
                     if os.path.isdir(os.path.join(KNOWLEDGE_DIR, d)) and not d.startswith('.')]

    for corp_dir in corp_dirs:
        corp_path = os.path.join(KNOWLEDGE_DIR, corp_dir)

        # 统计消息
        messages_dir = os.path.join(corp_path, "messages")
        if os.path.exists(messages_dir):
            for date_dir in os.listdir(messages_dir):
                date_path = os.path.join(messages_dir, date_dir)
                if not os.path.isdir(date_path):
                    continue

                md_files = [f for f in os.listdir(date_path) if f.endswith('.md')]
                stats["total_messages"] += len(md_files)
                stats["dates"][date_dir] = stats["dates"].get(date_dir, 0) + len(md_files)

        # 统计文件
        files_dir = os.path.join(corp_path, "files")
        if os.path.exists(files_dir):
            for date_dir in os.listdir(files_dir):
                date_path = os.path.join(files_dir, date_dir)
                if not os.path.isdir(date_path):
                    continue

                for file_name in os.listdir(date_path):
                    file_path = os.path.join(date_path, file_name)
                    if os.path.isfile(file_path):
                        stats["total_files"] += 1
                        size = os.path.getsize(file_path)
                        stats["total_size"] += size

                        ext = os.path.splitext(file_name)[1].lower()
                        file_type = _get_file_type(ext)
                        stats["file_types"][file_type] = stats["file_types"].get(file_type, 0) + 1

    stats["total_size_display"] = _format_size(stats["total_size"])

    return stats


@router.get("/download")
async def download_file(
    file_path: str = Query(..., description="文件路径"),
    file_name: Optional[str] = Query(None, description="下载时的文件名")
):
    """下载文件"""
    # 统一路径分隔符
    file_path = file_path.replace('\\', '/')

    # 如果是相对路径，转换为绝对路径
    if not os.path.isabs(file_path):
        project_root = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
        file_path = os.path.join(project_root, file_path)

    # 安全检查：确保路径在知识库目录内
    real_path = os.path.realpath(file_path)
    real_knowledge_dir = os.path.realpath(KNOWLEDGE_DIR)

    if not real_path.startswith(real_knowledge_dir):
        raise HTTPException(status_code=403, detail="禁止访问")

    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="文件不存在")

    # 如果是消息文件，返回文本内容
    if file_path.endswith('.md'):
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        download_name = file_name or os.path.basename(file_path)
        return PlainTextResponse(
            content=content,
            media_type='text/markdown',
            headers={'Content-Disposition': f'attachment; filename="{download_name}"'}
        )

    # 其他文件返回二进制
    download_name = file_name or os.path.basename(file_path)
    return FileResponse(
        path=file_path,
        filename=download_name,
        media_type='application/octet-stream'
    )


@router.get("/preview")
async def preview_file(file_path: str = Query(..., description="文件路径")):
    """预览文件内容"""
    # 统一路径分隔符
    file_path = file_path.replace('\\', '/')

    # 如果是相对路径，转换为绝对路径
    if not os.path.isabs(file_path):
        # 获取项目根目录
        project_root = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
        file_path = os.path.join(project_root, file_path)

    # 安全检查
    real_path = os.path.realpath(file_path)
    real_knowledge_dir = os.path.realpath(KNOWLEDGE_DIR)

    if not real_path.startswith(real_knowledge_dir):
        raise HTTPException(status_code=403, detail="禁止访问")

    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="文件不存在")

    # 根据文件类型返回预览
    ext = os.path.splitext(file_path)[1].lower()

    # 文本类文件直接返回内容
    text_extensions = ['.md', '.txt', '.csv', '.json', '.xml', '.html', '.py', '.js']
    if ext in text_extensions:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read(10000)  # 限制 10KB
            return {"type": "text", "content": content, "truncated": len(content) == 10000}
        except:
            with open(file_path, 'r', encoding='gbk', errors='ignore') as f:
                content = f.read(10000)
            return {"type": "text", "content": content, "truncated": len(content) == 10000}

    # Excel 文件
    if ext in ['.xls', '.xlsx']:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True)
            data = {}
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    if i >= 50:  # 限制 50 行
                        break
                    rows.append([str(cell) if cell else "" for cell in row])
                data[sheet_name] = rows
            wb.close()
            return {"type": "excel", "sheets": data}
        except Exception as e:
            return {"type": "error", "message": f"Excel 预览失败: {str(e)}"}

    # 图片文件
    if ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
        return {"type": "image", "path": file_path}

    # PDF 文件
    if ext == '.pdf':
        # 使用 pdfminer 提取文本
        try:
            from pdfminer.high_level import extract_text
            text = extract_text(file_path, maxpages=20)
            if text and text.strip():
                return {"type": "text", "content": text[:10000], "truncated": len(text) > 10000}
        except Exception as e:
            pass

        # pdfminer 失败，返回文件信息
        stat = os.stat(file_path)
        return {"type": "binary", "name": os.path.basename(file_path), "size": _format_size(stat.st_size), "extension": ext}

    # Word 文件
    if ext in ['.doc', '.docx']:
        try:
            from docx import Document
            doc = Document(file_path)
            paragraphs = []
            for i, para in enumerate(doc.paragraphs):
                if i >= 100:  # 限制 100 段
                    break
                if para.text.strip():
                    paragraphs.append(para.text)

            # 提取表格
            tables = []
            for table in doc.tables[:5]:  # 限制 5 个表格
                table_data = []
                for row in table.rows[:50]:  # 限制 50 行
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)

            content = "\n\n".join(paragraphs)
            return {"type": "word", "content": content[:10000], "tables": tables, "truncated": len(content) > 10000}
        except ImportError:
            stat = os.stat(file_path)
            return {"type": "binary", "name": os.path.basename(file_path), "size": _format_size(stat.st_size), "extension": ext}
        except Exception as e:
            return {"type": "error", "message": f"Word 预览失败: {str(e)}"}

    # PPT 文件
    if ext in ['.ppt', '.pptx']:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            slides = []
            for i, slide in enumerate(prs.slides):
                if i >= 50:  # 限制 50 页
                    break
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                slides.append({"slide": i + 1, "text": "\n".join(slide_text)})
            return {"type": "ppt", "slides": slides, "total_slides": len(prs.slides)}
        except ImportError:
            stat = os.stat(file_path)
            return {"type": "binary", "name": os.path.basename(file_path), "size": _format_size(stat.st_size), "extension": ext}
        except Exception as e:
            return {"type": "error", "message": f"PPT 预览失败: {str(e)}"}

    # 其他文件返回文件信息
    stat = os.stat(file_path)
    return {
        "type": "binary",
        "name": os.path.basename(file_path),
        "size": _format_size(stat.st_size),
        "extension": ext
    }


@router.delete("/{item_id}")
async def delete_knowledge_item(
    item_id: str,
    item_type: str = Query(..., description="类型: message/file"),
    corp_id: str = Query(..., description="企业ID"),
    date: str = Query(..., description="日期")
):
    """删除知识库项目（同步删除 RAG 索引）"""
    try:
        deleted_from_rag = 0

        if item_type == "message":
            # 删除消息文件
            file_path = os.path.join(KNOWLEDGE_DIR, corp_id, "messages", date, f"{item_id}.md")
            if os.path.exists(file_path):
                os.remove(file_path)
                # 同步删除 RAG 索引中的数据
                try:
                    kb = await _get_knowledge_base(corp_id)
                    deleted_from_rag = kb.delete_by_source(item_id)
                except Exception as e:
                    print(f"同步删除 RAG 索引失败: {e}")
                return {"message": "删除成功", "id": item_id, "rag_deleted": deleted_from_rag}
            else:
                raise HTTPException(status_code=404, detail="文件不存在")

        elif item_type == "file":
            # 删除上传的文件
            # item_id 格式: date_filename
            parts = item_id.split('_', 1)
            if len(parts) == 2:
                file_date, file_name = parts
                file_path = os.path.join(KNOWLEDGE_DIR, corp_id, "files", file_date, file_name)
                if os.path.exists(file_path):
                    os.remove(file_path)
                    # 同步删除 RAG 索引中的数据
                    # 尝试通过文件名查找 source_id
                    try:
                        kb = await _get_knowledge_base(corp_id)
                        # 查找包含该文件名的所有分块的 source_id
                        source_ids = set()
                        for chunk in kb._chunks:
                            if chunk.file_name == file_name:
                                source_ids.add(chunk.source_id)
                        # 删除这些 source_id 的所有分块
                        for sid in source_ids:
                            deleted_from_rag += kb.delete_by_source(sid)
                    except Exception as e:
                        print(f"同步删除 RAG 索引失败: {e}")
                    return {"message": "删除成功", "id": item_id, "rag_deleted": deleted_from_rag}
            raise HTTPException(status_code=404, detail="文件不存在")

        else:
            raise HTTPException(status_code=400, detail="无效的类型")

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"删除失败: {str(e)}")


class PermissionUpdate(BaseModel):
    """权限更新请求"""
    permission: str  # public, internal, confidential


@router.put("/{item_id}/permission")
async def update_item_permission(
    item_id: str,
    update: PermissionUpdate,
    item_type: str = Query(..., description="类型: message/file"),
    corp_id: str = Query(..., description="企业ID"),
    date: str = Query(..., description="日期")
):
    """更新知识库项目的权限等级"""
    # 验证权限值
    valid_permissions = ["public", "internal", "confidential"]
    if update.permission not in valid_permissions:
        raise HTTPException(status_code=400, detail=f"无效的权限值，可选值: {valid_permissions}")

    # 获取权限配置文件路径
    permissions_file = os.path.join(KNOWLEDGE_DIR, corp_id, "permissions.json")

    # 加载现有权限配置
    permissions = {}
    if os.path.exists(permissions_file):
        try:
            with open(permissions_file, 'r', encoding='utf-8') as f:
                permissions = json.load(f)
        except:
            pass

    # 更新权限
    permissions[item_id] = {
        "permission": update.permission,
        "item_type": item_type,
        "date": date,
    }

    # 保存权限配置
    os.makedirs(os.path.dirname(permissions_file), exist_ok=True)
    with open(permissions_file, 'w', encoding='utf-8') as f:
        json.dump(permissions, f, ensure_ascii=False, indent=2)

    return {"message": "权限更新成功", "item_id": item_id, "permission": update.permission}


@router.get("/structured")
async def get_structured_data(
    data_type: str = Query(..., description="数据类型: schedules, exams, contacts"),
    corp_id: Optional[str] = Query(None)
):
    """获取结构化数据（课表、考试、通讯录）"""
    results = []

    if corp_id:
        results = _load_structured_data(corp_id, data_type)
    else:
        # 加载所有企业
        for corp_dir in os.listdir(KNOWLEDGE_DIR):
            corp_path = os.path.join(KNOWLEDGE_DIR, corp_dir)
            if not os.path.isdir(corp_path) or corp_dir.startswith('.'):
                continue
            data = _load_structured_data(corp_dir, data_type)
            for item in data:
                item["corp_id"] = corp_dir
            results.extend(data)

    return {
        "data_type": data_type,
        "data": results,
        "total": len(results)
    }


def _load_structured_data(corp_id: str, data_type: str) -> list:
    """加载结构化数据"""
    data_file = os.path.join(KNOWLEDGE_DIR, corp_id, "structured", f"{data_type}.json")
    if not os.path.exists(data_file):
        return []

    try:
        with open(data_file, 'r', encoding='utf-8') as f:
            return json.load(f)
    except:
        return []


@router.post("/upload")
async def upload_to_knowledge(
    file: UploadFile = File(..., description="上传的文件"),
    corp_id: Optional[str] = Query(None, description="企业ID（不传则使用默认企业）"),
    tags: Optional[str] = Query(None, description="标签，逗号分隔"),
    access_level: str = Query("public", description="访问级别: public/internal/confidential")
):
    """
    上传文件到 RAG 知识库

    支持的文件类型：
    - 文档：txt, md, csv, json, xml, html
    - Office：doc, docx, xls, xlsx, ppt, pptx
    - PDF：pdf
    - 图片：jpg, jpeg, png, gif, bmp（OCR识别）
    - 代码：py, js, java, c, cpp 等
    """
    try:
        # 获取或验证企业ID
        if not corp_id:
            corp_id = _get_default_corp_id()
            if not corp_id:
                raise HTTPException(status_code=400, detail="未找到企业配置")

        # 验证企业目录存在
        corp_path = os.path.join(KNOWLEDGE_DIR, corp_id)
        if not os.path.exists(corp_path):
            raise HTTPException(status_code=404, detail=f"企业 {corp_id} 不存在")

        # 创建文件保存目录（按日期）
        from datetime import date
        today = date.today().strftime("%Y-%m-%d")
        files_dir = os.path.join(corp_path, "files", today)
        os.makedirs(files_dir, exist_ok=True)

        # 保存上传的文件
        file_name = file.filename or "upload_file"
        file_path = os.path.join(files_dir, file_name)

        # 检查文件是否已存在，添加后缀
        if os.path.exists(file_path):
            name, ext = os.path.splitext(file_name)
            file_name = f"{name}_{hashlib.md5(file_name.encode()).hexdigest()[:8]}{ext}"
            file_path = os.path.join(files_dir, file_name)

        # 读取并保存文件内容
        content = await file.read()
        with open(file_path, "wb") as f:
            f.write(content)

        file_size = len(content)

        # 提取文件文本内容
        from agent.media_handler import extract_text_from_file
        extracted_text, file_type = await extract_text_from_file(file_path)

        if not extracted_text or not extracted_text.strip():
            # 文本提取失败，只保存文件，不加入 RAG 索引
            return {
                "message": "文件已保存，但未能提取文本内容",
                "file_name": file_name,
                "file_path": file_path,
                "file_size": file_size,
                "file_type": file_type,
                "rag_indexed": False
            }

        # 添加到 RAG 索引
        from agent.knowledge_base_v2 import get_knowledge_base
        kb = get_knowledge_base(corp_path, corp_id)

        # 生成 source_id
        source_id = f"web_upload_{hashlib.md5(f'{corp_id}_{file_name}_{today}'.encode()).hexdigest()[:16]}"

        # 解析标签
        tag_list = [t.strip() for t in tags.split(",")] if tags else []

        # 添加到知识库
        chunks = await kb.add_message(
            text=extracted_text,
            source_type="file",
            source_id=source_id,
            sender_id="web_upload",
            sender_nick="Web管理后台",
            conversation_id="",
            message_type="file",
            file_name=file_name,
            file_path=file_path,
            tags=tag_list,
            file_size=file_size,
            file_type=file_type,
            access_level=access_level,
        )

        return {
            "message": "文件上传并索引成功",
            "file_name": file_name,
            "file_path": file_path,
            "file_size": file_size,
            "file_type": file_type,
            "extracted_text_length": len(extracted_text),
            "rag_indexed": True,
            "chunks_created": len(chunks) if chunks else 0
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"上传失败: {str(e)}")


@router.post("/batch-delete")
async def batch_delete_knowledge(
    item_ids: list,
    item_type: str = Query(..., description="类型: message/file"),
    corp_id: str = Query(..., description="企业ID"),
    date: str = Query(..., description="日期")
):
    """批量删除知识库项目"""
    results = []
    for item_id in item_ids:
        try:
            if item_type == "message":
                file_path = os.path.join(KNOWLEDGE_DIR, corp_id, "messages", date, f"{item_id}.md")
                if os.path.exists(file_path):
                    os.remove(file_path)
                    # 同步删除 RAG 索引
                    try:
                        kb = await _get_knowledge_base(corp_id)
                        kb.delete_by_source(item_id)
                    except Exception as e:
                        print(f"同步删除 RAG 索引失败: {e}")
                    results.append({"id": item_id, "status": "deleted"})
                else:
                    results.append({"id": item_id, "status": "not_found"})

            elif item_type == "file":
                parts = item_id.split('_', 1)
                if len(parts) == 2:
                    file_date, file_name = parts
                    file_path = os.path.join(KNOWLEDGE_DIR, corp_id, "files", file_date, file_name)
                    if os.path.exists(file_path):
                        os.remove(file_path)
                        # 同步删除 RAG 索引
                        try:
                            kb = await _get_knowledge_base(corp_id)
                            source_ids = set()
                            for chunk in kb._chunks:
                                if chunk.file_name == file_name:
                                    source_ids.add(chunk.source_id)
                            for sid in source_ids:
                                kb.delete_by_source(sid)
                        except Exception as e:
                            print(f"同步删除 RAG 索引失败: {e}")
                        results.append({"id": item_id, "status": "deleted"})
                    else:
                        results.append({"id": item_id, "status": "not_found"})

        except Exception as e:
            results.append({"id": item_id, "status": "error", "message": str(e)})

    return {"results": results, "total": len(results)}


@router.get("/rag-stats")
async def get_rag_stats(corp_id: Optional[str] = None):
    """获取 RAG 知识库统计信息"""
    try:
        if not corp_id:
            corp_id = _get_default_corp_id()
            if not corp_id:
                return {"error": "未找到企业配置"}

        kb = await _get_knowledge_base(corp_id)
        stats = kb.get_stats()

        return {
            "corp_id": corp_id,
            "total_chunks": stats.total_chunks,
            "source_types": stats.source_types,
            "categories": stats.categories,
            "top_senders": stats.top_senders[:10] if stats.top_senders else [],
        }
    except Exception as e:
        return {"error": str(e)}
